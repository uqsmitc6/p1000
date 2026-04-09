#!/usr/bin/env python3
"""
UQ Slide Compliance Tool — Web UI
===================================
Streamlit interface for the brand fixer and image audit tools.
Designed for deployment on Streamlit Community Cloud.

Usage (local):
  streamlit run app.py

Deployment:
  Push to GitHub → connect to share.streamlit.io
  Set ANTHROPIC_API_KEY in Streamlit Cloud secrets
"""

APP_VERSION = "2.3.0"

import io
import json
import os
import sys
import tempfile
import time
from pathlib import Path
from collections import defaultdict

import streamlit as st

# Ensure our scripts are importable
sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from brand_fixer import BrandFixer
from image_audit import extract_images, classify_image, generate_html_report
from ref_checker import RefChecker
from combined_pipeline import run_pipeline
from cost_logger import log_cost, render_sidebar_admin


# ─── API Key from secrets ──────────────────────────────────────────────

def get_api_key():
    """Retrieve API key from Streamlit secrets (set in Cloud dashboard)."""
    try:
        return st.secrets["ANTHROPIC_API_KEY"]
    except (KeyError, FileNotFoundError):
        return os.environ.get("ANTHROPIC_API_KEY")


# ─── Page Config ───────────────────────────────────────────────────────

st.set_page_config(
    page_title="UQ Slide Compliance Tool",
    page_icon="🟣",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ─── Custom CSS ────────────────────────────────────────────────────────

st.markdown("""
<style>
    /* UQ Purple branding */
    .stApp > header { background-color: #51247A; }

    .uq-header {
        background: linear-gradient(135deg, #51247A 0%, #962A8B 100%);
        color: white;
        padding: 1.5rem 2rem;
        border-radius: 8px;
        margin-bottom: 1.5rem;
    }
    .uq-header h1 { color: white; margin: 0 0 0.3rem 0; font-size: 1.8rem; }
    .uq-header p { color: rgba(255,255,255,0.85); margin: 0; font-size: 0.95rem; }

    /* Stat cards */
    .stat-card {
        background: white;
        border-radius: 8px;
        padding: 1rem 1.2rem;
        text-align: center;
        box-shadow: 0 1px 3px rgba(0,0,0,0.1);
    }
    .stat-card .number { font-size: 1.8rem; font-weight: bold; line-height: 1.2; }
    .stat-card .label { font-size: 0.8rem; color: #666; margin-top: 0.2rem; }

    /* Risk colours */
    .risk-critical { color: #DC2626; }
    .risk-high { color: #EA580C; }
    .risk-medium { color: #D97706; }
    .risk-low { color: #16A34A; }
    .risk-clear { color: #059669; }

    /* Change list */
    .change-item {
        padding: 0.4rem 0.6rem;
        margin: 0.2rem 0;
        border-radius: 4px;
        font-size: 0.85rem;
        border-left: 3px solid #D7D1CC;
    }
    .change-font { border-left-color: #4085C6; background: #F0F7FF; }
    .change-colour { border-left-color: #51247A; background: #F5F0FA; }
    .change-colour_flagged { border-left-color: #D97706; background: #FFFBEB; }
    .change-table { border-left-color: #16A34A; background: #F0FDF4; }
    .change-footer { border-left-color: #962A8B; background: #FDF0FA; }
    .change-heading_size { border-left-color: #E62645; background: #FFF0F3; }
    .change-bullet { border-left-color: #FBB800; background: #FFFDE7; }
    .change-body_size_flagged { border-left-color: #7C3AED; background: #F5F3FF; }
    .change-citation { border-left-color: #4085C6; background: #F0F7FF; }
    .change-reference { border-left-color: #51247A; background: #F5F0FA; }
    .change-attribution { border-left-color: #16A34A; background: #F0FDF4; }
    .change-cross_ref { border-left-color: #D97706; background: #FFFBEB; }
    .change-missing_attr { border-left-color: #E62645; background: #FFF0F3; }

    /* Hide deploy button and Streamlit branding */
    .stDeployButton { display: none; }
    #MainMenu { visibility: hidden; }
    footer { visibility: hidden; }

    /* Info box */
    .info-box {
        background: #F5F0FA;
        border: 1px solid #D7D1CC;
        border-radius: 8px;
        padding: 1rem 1.2rem;
        margin: 1rem 0;
        font-size: 0.9rem;
    }
    .info-box strong { color: #51247A; }
</style>
""", unsafe_allow_html=True)


# ─── Header ────────────────────────────────────────────────────────────

st.markdown(f"""
<div class="uq-header">
    <h1>UQ Slide Compliance Tool <span style="font-size: 0.5em; opacity: 0.7; font-weight: normal;">v{APP_VERSION}</span></h1>
    <p>Brand formatting fixer &amp; image copyright audit for executive education decks</p>
</div>
""", unsafe_allow_html=True)


# ─── Sidebar (minimal — just info) ────────────────────────────────────

with st.sidebar:
    # ── v5 Dependency Diagnostic ──
    st.markdown("### v5 Engine Status")
    _diag_api = bool(get_api_key())
    try:
        import anthropic as _anth_check
        _diag_anthropic = True
    except ImportError:
        _diag_anthropic = False
    try:
        import fitz as _fitz_check
        _diag_fitz = True
    except ImportError:
        _diag_fitz = False
    import shutil
    _diag_libre = shutil.which("libreoffice") is not None

    _checks = [
        ("API key", _diag_api),
        ("anthropic pkg", _diag_anthropic),
        ("PyMuPDF (fitz)", _diag_fitz),
        ("LibreOffice", _diag_libre),
    ]
    for label, ok in _checks:
        icon = "✅" if ok else "❌"
        st.markdown(f"{icon} {label}")
    if all(ok for _, ok in _checks):
        st.success("v5 AI QA ready")
    else:
        missing = [label for label, ok in _checks if not ok]
        st.warning(f"v5 AI QA disabled — missing: {', '.join(missing)}")
    st.markdown("---")

    st.markdown("### How to use")
    st.markdown(
        "**Brand Fixer** — Upload a `.pptx`, click the button, "
        "download your fixed file. No API key needed."
    )
    st.markdown(
        "**Image Audit** — Upload a `.pptx` and the tool will extract "
        "every image and classify its copyright risk using AI."
    )
    st.markdown(
        "**Reference Checker** — Scans citations, references, and image "
        "attributions. Auto-fixes formatting where possible."
    )
    st.markdown(
        "**Full Compliance Check** — Runs all three checks in sequence "
        "and produces a unified report. One upload, one click."
    )
    st.markdown("---")
    st.markdown(
        "**Tips:**\n"
        "- Start with 'Extract only' to preview images before running the full audit\n"
        "- For large decks, set a limit (e.g. 10) to test before running all images\n"
        "- The brand fixer flags uncertain colours for you to review manually\n"
        "- Use 'Full Compliance Check' to run everything at once"
    )
    st.markdown("---")

    # API status indicator
    api_key = get_api_key()
    if api_key:
        st.success("API key configured", icon="✅")
    else:
        st.warning("API key not configured — image audit classification unavailable", icon="⚠️")

    st.markdown(
        "<small style='color: #999;'>Built for UQ Business School<br/>"
        "Learning Design Team</small>",
        unsafe_allow_html=True,
    )

    st.markdown("---")
    render_sidebar_admin()


# ─── Main Tabs ─────────────────────────────────────────────────────────

tab1, tab2, tab3, tab4 = st.tabs(["Brand Fixer", "Image Audit", "Reference Checker", "Full Compliance Check"])


# ═══════════════════════════════════════════════════════════════════════
# TAB 1: Brand Fixer
# ═══════════════════════════════════════════════════════════════════════

with tab1:
    st.markdown("#### Upload a slide deck to fix brand formatting")

    st.markdown("""
    <div class="info-box">
        <strong>What this does:</strong> Corrects fonts to Arial, fixes text colours
        to the UQ palette, restyles tables with UQ brand colours, normalises heading
        sizes, and standardises bullets. Your original file is never modified —
        you'll download a new fixed copy.
    </div>
    """, unsafe_allow_html=True)

    uploaded_file = st.file_uploader(
        "Choose a .pptx file",
        type=["pptx"],
        key="brand_upload",
        help="Upload the slide deck you want to fix",
    )

    if uploaded_file is not None:
        # Show file info
        file_size_mb = len(uploaded_file.getvalue()) / (1024 * 1024)
        st.caption(f"Uploaded: **{uploaded_file.name}** ({file_size_mb:.1f} MB)")

        if st.button("Fix Brand Formatting", type="primary", key="run_brand"):
            # Save uploaded file to temp
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp.write(uploaded_file.getvalue())
                tmp_path = tmp.name

            try:
                # Load and fix
                prs = Presentation(tmp_path)
                fixer = BrandFixer(prs, report=True)

                # Progress updates
                progress = st.progress(0, text="Fixing fonts...")
                fixer.fix_fonts()
                progress.progress(14, text="Fixing text colours...")
                fixer.fix_colours()
                progress.progress(28, text="Restyling tables...")
                fixer.fix_tables()
                progress.progress(42, text="Standardising footers...")
                fixer.fix_footers()
                progress.progress(56, text="Normalising heading sizes...")
                fixer.fix_heading_sizes()
                progress.progress(70, text="Checking body text sizes...")
                fixer.flag_body_text_sizes()
                progress.progress(84, text="Fixing bullet styles...")
                fixer.fix_bullets()
                progress.progress(100, text="Done!")

                # Save fixed file to buffer
                output_buffer = io.BytesIO()
                prs.save(output_buffer)
                output_bytes = output_buffer.getvalue()

                report = fixer.generate_report()
                total = report["total_changes"]
                num_slides = len(prs.slides)

                # Store results in session state so they survive reruns
                st.session_state["brand_result"] = {
                    "output_bytes": output_bytes,
                    "report": report,
                    "total": total,
                    "num_slides": num_slides,
                    "fixed_name": uploaded_file.name.replace(".pptx", "_FIXED.pptx"),
                }

            except Exception as e:
                st.error(f"Something went wrong: {e}")
                with st.expander("Error details"):
                    import traceback
                    st.code(traceback.format_exc())
            finally:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass

        # ── Display results from session state (persists across reruns) ──
        if "brand_result" in st.session_state:
            result = st.session_state["brand_result"]
            report = result["report"]
            total = result["total"]

            st.markdown("---")
            st.markdown("### Results")

            if total == 0:
                st.success(
                    "No changes needed — this deck is already brand-compliant!",
                    icon="✅",
                )
            else:
                stats = report["summary"]
                cols = st.columns(8)
                stat_items = [
                    ("Font fixes", stats.get("font", 0), "#4085C6"),
                    ("Colour fixes", stats.get("colour", 0), "#51247A"),
                    ("Flagged", stats.get("colour_flagged", 0), "#D97706"),
                    ("Tables", stats.get("table", 0), "#16A34A"),
                    ("Footers", stats.get("footer", 0), "#962A8B"),
                    ("Headings", stats.get("heading_size", 0), "#E62645"),
                    ("Body size", stats.get("body_size_flagged", 0), "#7C3AED"),
                    ("Bullets", stats.get("bullet", 0), "#FBB800"),
                ]
                for col, (label, count, colour) in zip(cols, stat_items):
                    col.markdown(
                        f"<div class='stat-card'>"
                        f"<div class='number' style='color:{colour};'>{count}</div>"
                        f"<div class='label'>{label}</div></div>",
                        unsafe_allow_html=True,
                    )

                st.markdown(f"**{total} total changes** across {result['num_slides']} slides")

                # ── Download button ──
                st.download_button(
                    label=f"Download {result['fixed_name']}",
                    data=result["output_bytes"],
                    file_name=result["fixed_name"],
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    type="primary",
                )

                # ── Flagged items (show first — these need attention) ──
                flagged = [
                    c for c in report["changes"]
                    if c["category"] == "colour_flagged"
                ]
                if flagged:
                    with st.expander(
                        f"⚠️  {len(flagged)} colours flagged for review",
                        expanded=True,
                    ):
                        st.markdown(
                            "These non-UQ colours were **not auto-corrected** "
                            "because they may be intentional accent colours. "
                            "Please check these slides and decide whether to "
                            "keep or change them:"
                        )
                        # Deduplicate by colour + slide for cleaner display
                        seen = set()
                        for c in flagged:
                            key = (c["slide"], c["detail"])
                            if key not in seen:
                                seen.add(key)
                                st.markdown(
                                    f"<div class='change-item change-colour_flagged'>"
                                    f"Slide {c['slide']}: {c['detail']}</div>",
                                    unsafe_allow_html=True,
                                )

                # ── Body text size flags ──
                body_flagged = [
                    c for c in report["changes"]
                    if c["category"] == "body_size_flagged"
                ]
                if body_flagged:
                    with st.expander(
                        f"📏  {len(body_flagged)} body text sizes flagged for review",
                        expanded=False,
                    ):
                        st.markdown(
                            "These body text runs are outside the recommended "
                            "12–24pt range. They were **not auto-corrected** "
                            "as small text may be intentional (captions, footnotes). "
                            "Review and adjust if needed:"
                        )
                        for c in body_flagged:
                            st.markdown(
                                f"<div class='change-item change-body_size_flagged'>"
                                f"Slide {c['slide']}: {c['detail']}</div>",
                                unsafe_allow_html=True,
                            )

                # ── Detailed changes (collapsed) ──
                if report["changes"]:
                    with st.expander(
                        f"View all {total} changes", expanded=False
                    ):
                        by_slide = defaultdict(list)
                        for change in report["changes"]:
                            by_slide[change["slide"]].append(change)

                        for slide_num in sorted(by_slide.keys()):
                            changes = by_slide[slide_num]
                            st.markdown(
                                f"**Slide {slide_num}** "
                                f"({len(changes)} change{'s' if len(changes) != 1 else ''})"
                            )
                            for c in changes:
                                cat = c["category"]
                                st.markdown(
                                    f"<div class='change-item change-{cat}'>"
                                    f"{c['detail']}</div>",
                                    unsafe_allow_html=True,
                                )


# ═══════════════════════════════════════════════════════════════════════
# TAB 2: Image Audit
# ═══════════════════════════════════════════════════════════════════════

with tab2:
    st.markdown("#### Upload a slide deck to audit images for copyright risk")

    st.markdown("""
    <div class="info-box">
        <strong>What this does:</strong> Extracts every image from the deck and
        uses AI to classify each one by type (stock photo, screenshot, diagram, etc.)
        and copyright risk level (Critical → Clear). You'll get a downloadable
        report showing which images need attention.<br/><br/>
        <strong>Why it matters:</strong> Executive education programs are commercial —
        they fall outside the statutory education licence. Every uncleared copyrighted
        image is genuine legal exposure.
    </div>
    """, unsafe_allow_html=True)

    audit_file = st.file_uploader(
        "Choose a .pptx file",
        type=["pptx"],
        key="audit_upload",
        help="Upload the slide deck you want to audit",
    )

    if audit_file is not None:
        file_size_mb = len(audit_file.getvalue()) / (1024 * 1024)
        st.caption(f"Uploaded: **{audit_file.name}** ({file_size_mb:.1f} MB)")

        col_a, col_b = st.columns(2)
        with col_a:
            limit = st.number_input(
                "Limit images (0 = all)",
                min_value=0,
                max_value=500,
                value=0,
                help="Limit to first N images — useful for testing. Set to 0 to audit all.",
            )
        with col_b:
            extract_only = st.checkbox(
                "Extract only (no AI classification)",
                help="Just see what images are in the deck without running the AI. Free and fast.",
            )

        if st.button("Run Image Audit", type="primary", key="run_audit"):
            # Validate API key if doing classification
            if not extract_only:
                api_key = get_api_key()
                if not api_key:
                    st.error(
                        "The API key hasn't been configured yet. "
                        "Ask Sean to set it up, or tick 'Extract only' to preview images without AI."
                    )
                    st.stop()

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp.write(audit_file.getvalue())
                tmp_path = tmp.name

            try:
                # ── Extract ──
                with st.spinner("Extracting images from presentation..."):
                    with tempfile.TemporaryDirectory() as extract_dir:
                        img_limit = limit if limit > 0 else None
                        images = extract_images(
                            tmp_path, output_dir=extract_dir, limit=img_limit
                        )

                if not images:
                    st.info("No images found in this presentation.")
                    st.stop()

                num_unique_slides = len(set(i['slide_number'] for i in images))

                # ── Extract-only mode ──
                if extract_only:
                    # Store in session state
                    st.session_state["audit_result"] = {
                        "mode": "extract_only",
                        "images": images,
                        "num_unique_slides": num_unique_slides,
                        "source_name": audit_file.name,
                    }

                # ── Full AI classification ──
                else:
                    st.markdown("---")

                    import anthropic
                    client = anthropic.Anthropic(api_key=api_key)

                    classifications = []
                    progress = st.progress(0)
                    status_text = st.empty()

                    for i, img_info in enumerate(images):
                        status_text.markdown(
                            f"Classifying image **{i+1}/{len(images)}**: "
                            f"`{img_info['filename']}` (slide {img_info['slide_number']})"
                        )
                        progress.progress(i / len(images))

                        result = classify_image(client, img_info)
                        classifications.append(result)

                        # Brief pause for rate limiting
                        if i < len(images) - 1:
                            time.sleep(0.3)

                    progress.progress(1.0)
                    status_text.empty()

                    # Generate self-contained HTML report with embedded images
                    image_bytes_map = {
                        img["filename"]: img["image_bytes"]
                        for img in images if "image_bytes" in img
                    }
                    images_clean = [
                        {k: v for k, v in img.items() if k != "image_bytes"}
                        for img in images
                    ]
                    with tempfile.NamedTemporaryFile(
                        suffix=".html", delete=False, mode="w"
                    ) as report_tmp:
                        report_path = report_tmp.name

                    generate_html_report(
                        images_clean, classifications,
                        audit_file.name, report_path, None,
                        image_bytes_map=image_bytes_map,
                    )
                    html_content = Path(report_path).read_text()
                    try:
                        os.unlink(report_path)
                    except Exception:
                        pass

                    # Build JSON data
                    risk_counts = defaultdict(int)
                    for cls in classifications:
                        if "error" not in cls:
                            risk_counts[cls.get("risk_level", "UNKNOWN")] += 1

                    json_data = {
                        "metadata": {
                            "input_file": audit_file.name,
                            "total_images": len(images),
                        },
                        "summary": {"risk_counts": dict(risk_counts)},
                        "images": [
                            {
                                **{k: v for k, v in img.items() if k != "image_bytes"},
                                **cls,
                            }
                            for img, cls in zip(images, classifications)
                        ],
                    }

                    # Calculate cost and log
                    from image_audit import COST_PER_M_INPUT, COST_PER_M_OUTPUT
                    total_input_tokens = sum(
                        c.get("api_tokens_used", {}).get("input", 0)
                        for c in classifications
                    )
                    total_output_tokens = sum(
                        c.get("api_tokens_used", {}).get("output", 0)
                        for c in classifications
                    )
                    cost_usd = (
                        (total_input_tokens / 1_000_000) * COST_PER_M_INPUT
                        + (total_output_tokens / 1_000_000) * COST_PER_M_OUTPUT
                    )
                    log_cost(
                        tool="Image Audit",
                        filename=audit_file.name,
                        num_images=len(images),
                        input_tokens=total_input_tokens,
                        output_tokens=total_output_tokens,
                        cost_usd=cost_usd,
                    )

                    # Store in session state
                    st.session_state["audit_result"] = {
                        "mode": "classified",
                        "images": images,
                        "classifications": classifications,
                        "risk_counts": dict(risk_counts),
                        "html_content": html_content,
                        "json_data": json_data,
                        "num_unique_slides": num_unique_slides,
                        "source_name": audit_file.name,
                    }

            except Exception as e:
                st.error(f"Something went wrong: {e}")
                with st.expander("Error details"):
                    import traceback
                    st.code(traceback.format_exc())
            finally:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass

        # ── Display results from session state (persists across reruns) ──
        if "audit_result" in st.session_state:
            ar = st.session_state["audit_result"]
            images = ar["images"]

            st.success(
                f"Found **{len(images)}** unique images across "
                f"{ar['num_unique_slides']} slides"
            )

            if ar["mode"] == "extract_only":
                st.markdown("---")
                st.markdown("### Extracted Images")
                st.caption("Showing all images found in the deck. Run with AI classification to get risk ratings.")

                for img_info in images:
                    col1, col2 = st.columns([1, 3])
                    with col1:
                        try:
                            from PIL import Image as PILImage
                            pil_img = PILImage.open(io.BytesIO(img_info["image_bytes"]))
                            st.image(pil_img, width=180)
                        except Exception:
                            st.caption("*(preview unavailable)*")
                    with col2:
                        st.markdown(
                            f"**Slide {img_info['slide_number']}** — "
                            f"`{img_info['filename']}`"
                        )
                        st.caption(
                            f"{img_info['width']}×{img_info['height']}px "
                            f"  |  {img_info['content_type']}"
                        )
                        if img_info.get("slide_title"):
                            st.caption(f"Slide title: {img_info['slide_title']}")
                    st.markdown("---")

            elif ar["mode"] == "classified":
                classifications = ar["classifications"]
                risk_counts = ar["risk_counts"]

                st.success("Classification complete!")

                # ── Risk Summary ──
                st.markdown("### Risk Summary")

                risk_items = [
                    ("Critical", risk_counts.get("CRITICAL", 0), "critical", "#DC2626"),
                    ("High", risk_counts.get("HIGH", 0), "high", "#EA580C"),
                    ("Medium", risk_counts.get("MEDIUM", 0), "medium", "#D97706"),
                    ("Low", risk_counts.get("LOW", 0), "low", "#16A34A"),
                    ("Clear", risk_counts.get("CLEAR", 0), "clear", "#059669"),
                ]

                cols = st.columns(5)
                for col, (label, count, cls_name, colour) in zip(cols, risk_items):
                    col.markdown(
                        f"<div class='stat-card'>"
                        f"<div class='number' style='color:{colour};'>{count}</div>"
                        f"<div class='label'>{label}</div></div>",
                        unsafe_allow_html=True,
                    )

                # Key message based on results
                critical_high = risk_counts.get("CRITICAL", 0) + risk_counts.get("HIGH", 0)
                if critical_high > 0:
                    st.warning(
                        f"**{critical_high} image{'s' if critical_high != 1 else ''} "
                        f"flagged as Critical or High risk** — "
                        f"these likely need replacement or licence verification.",
                        icon="⚠️",
                    )

                # ── Downloads ──
                st.markdown("---")
                st.markdown("### Download Report")

                dl_col1, dl_col2 = st.columns(2)

                with dl_col1:
                    report_name = ar["source_name"].replace(".pptx", "_audit_report.html")
                    st.download_button(
                        label="Download HTML Report",
                        data=ar["html_content"],
                        file_name=report_name,
                        mime="text/html",
                        type="primary",
                    )

                with dl_col2:
                    json_name = ar["source_name"].replace(".pptx", "_audit_data.json")
                    st.download_button(
                        label="Download JSON Data",
                        data=json.dumps(ar["json_data"], indent=2, default=str),
                        file_name=json_name,
                        mime="application/json",
                    )

                # ── Image cards ──
                st.markdown("---")
                st.markdown("### Image Details")
                st.caption("Sorted by risk level (highest first)")

                risk_order = {
                    "CRITICAL": 0, "HIGH": 1, "MEDIUM": 2,
                    "LOW": 3, "CLEAR": 4, "UNKNOWN": 5,
                }
                paired = sorted(
                    zip(images, classifications),
                    key=lambda x: risk_order.get(
                        x[1].get("risk_level", "UNKNOWN"), 5
                    ),
                )

                # Risk filter
                show_levels = st.multiselect(
                    "Filter by risk level:",
                    ["CRITICAL", "HIGH", "MEDIUM", "LOW", "CLEAR"],
                    default=["CRITICAL", "HIGH", "MEDIUM", "LOW", "CLEAR"],
                )

                for img_info, cls in paired:
                    risk_level = cls.get("risk_level", "UNKNOWN")
                    if risk_level not in show_levels:
                        continue

                    col1, col2 = st.columns([1, 3])
                    with col1:
                        try:
                            from PIL import Image as PILImage
                            pil_img = PILImage.open(
                                io.BytesIO(img_info["image_bytes"])
                            )
                            st.image(pil_img, width=200)
                        except Exception:
                            st.caption("*(preview unavailable)*")

                    with col2:
                        risk_display = cls.get("risk_level", "ERROR")
                        img_type = cls.get("image_type", "ERROR")
                        action = cls.get("recommended_action", "N/A")

                        # Header line
                        risk_colour = {
                            "CRITICAL": "#DC2626", "HIGH": "#EA580C",
                            "MEDIUM": "#D97706", "LOW": "#16A34A",
                            "CLEAR": "#059669",
                        }.get(risk_display, "#6B7280")

                        st.markdown(
                            f"**Slide {img_info['slide_number']}** &nbsp; "
                            f"<span style='background:{risk_colour};color:white;"
                            f"padding:2px 8px;border-radius:4px;font-size:0.8rem;'>"
                            f"{risk_display}</span> &nbsp; "
                            f"<span style='background:#E5E7EB;padding:2px 8px;"
                            f"border-radius:4px;font-size:0.8rem;'>{img_type}</span>",
                            unsafe_allow_html=True,
                        )

                        if "error" in cls:
                            st.error(cls["error"])
                        else:
                            st.markdown(cls.get("reasoning", ""))

                            if cls.get("content_description"):
                                st.caption(f"Content: {cls['content_description']}")

                            if cls.get("alt_text"):
                                st.caption(f"Suggested alt text: *{cls['alt_text']}*")

                            st.caption(f"Recommended action: **{action}**")

                            flags = []
                            if cls.get("watermark_text"):
                                flags.append(f"Watermark: {cls['watermark_text']}")
                            if cls.get("copyright_notice"):
                                flags.append(f"Copyright: {cls['copyright_notice']}")
                            if cls.get("brand_visible"):
                                flags.append(f"Brand: {cls['brand_visible']}")
                            if flags:
                                st.warning(" | ".join(flags))

                    st.markdown("---")


# ═══════════════════════════════════════════════════════════════════════
# TAB 3: Reference Checker
# ═══════════════════════════════════════════════════════════════════════

with tab3:
    st.markdown("#### Upload a slide deck to check references and image attributions")

    st.markdown("""
    <div class="info-box">
        <strong>What this does:</strong> Scans all slide text for in-text citations,
        reference lists, and image attribution text. Checks APA 7 formatting,
        standardises image attributions (e.g. to "Source: Adobe Stock {ID}"),
        cross-references citations against the reference list, and flags slides
        with images but no attribution.<br/><br/>
        <strong>Auto-fixes:</strong> Standardises attribution formats, fixes
        "&amp;" vs "and" in citations, adds missing periods after "et al".
        Downloads include both the fixed file and a detailed report.
    </div>
    """, unsafe_allow_html=True)

    ref_file = st.file_uploader(
        "Choose a .pptx file",
        type=["pptx"],
        key="ref_upload",
        help="Upload the slide deck you want to check",
    )

    if ref_file is not None:
        file_size_mb = len(ref_file.getvalue()) / (1024 * 1024)
        st.caption(f"Uploaded: **{ref_file.name}** ({file_size_mb:.1f} MB)")

        col_mode_a, col_mode_b = st.columns(2)
        with col_mode_a:
            report_only = st.checkbox(
                "Report only (don't modify file)",
                help="Just scan and report — don't make any changes to the file",
            )

        if st.button("Check References", type="primary", key="run_ref"):
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp.write(ref_file.getvalue())
                tmp_path = tmp.name

            try:
                prs = Presentation(tmp_path)
                checker = RefChecker(prs, report=True)

                progress = st.progress(0, text="Scanning citations...")
                checker.scan_citations()
                progress.progress(25, text="Scanning reference lists...")
                checker.scan_references()
                progress.progress(50, text="Checking image attributions...")
                checker.scan_attributions()
                progress.progress(75, text="Cross-referencing...")
                checker.cross_reference()

                if not report_only:
                    progress.progress(85, text="Applying fixes...")
                    checker.fix_attributions()
                    checker.fix_citations()

                progress.progress(100, text="Done!")

                report = checker.generate_report()

                # Save fixed file if not report-only
                output_bytes = None
                if not report_only:
                    output_buffer = io.BytesIO()
                    prs.save(output_buffer)
                    output_bytes = output_buffer.getvalue()

                # Store in session state
                st.session_state["ref_result"] = {
                    "report": report,
                    "output_bytes": output_bytes,
                    "report_only": report_only,
                    "source_name": ref_file.name,
                    "fixed_name": ref_file.name.replace(".pptx", "_REFFIXED.pptx"),
                    "num_slides": len(prs.slides),
                }

            except Exception as e:
                st.error(f"Something went wrong: {e}")
                with st.expander("Error details"):
                    import traceback
                    st.code(traceback.format_exc())
            finally:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass

        # ── Display results from session state ──
        if "ref_result" in st.session_state:
            rr = st.session_state["ref_result"]
            report = rr["report"]

            st.markdown("---")
            st.markdown("### Results")

            # ── Summary stats ──
            summary = report["summary"]
            cols = st.columns(5)
            stat_items = [
                ("Citations", summary["citations_found"], "#4085C6"),
                ("References", summary["references_found"], "#51247A"),
                ("Issues", report["total_issues"], "#D97706"),
                ("Auto-fixed", report["total_changes"], "#16A34A"),
                ("Ref slides", len(summary["ref_slides"]), "#962A8B"),
            ]
            for col, (label, count, colour) in zip(cols, stat_items):
                col.markdown(
                    f"<div class='stat-card'>"
                    f"<div class='number' style='color:{colour};'>{count}</div>"
                    f"<div class='label'>{label}</div></div>",
                    unsafe_allow_html=True,
                )

            if report["total_issues"] == 0:
                st.success(
                    "No reference or attribution issues found!",
                    icon="✅",
                )
            else:
                # Issue breakdown
                by_cat = summary.get("issues_by_category", {})
                by_sev = summary.get("issues_by_severity", {})

                if by_sev.get("warning", 0) > 0 or by_sev.get("error", 0) > 0:
                    warn_count = by_sev.get("warning", 0) + by_sev.get("error", 0)
                    st.warning(
                        f"**{warn_count} issue{'s' if warn_count != 1 else ''} "
                        f"need attention** — see details below.",
                        icon="⚠️",
                    )

                # ── Download buttons ──
                dl_cols = st.columns(2)
                with dl_cols[0]:
                    if not rr["report_only"] and rr["output_bytes"]:
                        st.download_button(
                            label=f"Download {rr['fixed_name']}",
                            data=rr["output_bytes"],
                            file_name=rr["fixed_name"],
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            type="primary",
                        )

                with dl_cols[1]:
                    json_name = rr["source_name"].replace(".pptx", "_ref_report.json")
                    st.download_button(
                        label="Download JSON Report",
                        data=json.dumps(report, indent=2, default=str),
                        file_name=json_name,
                        mime="application/json",
                    )

                # ── Changes made (if auto-fixed) ──
                if report["changes"]:
                    with st.expander(
                        f"✅ {report['total_changes']} auto-fixes applied",
                        expanded=True,
                    ):
                        for c in report["changes"]:
                            cat = c["category"]
                            st.markdown(
                                f"<div class='change-item change-{cat}'>"
                                f"Slide {c['slide']}: {c['detail']}</div>",
                                unsafe_allow_html=True,
                            )

                # ── Issues by category ──

                # Missing attributions
                missing = [i for i in report["issues"]
                           if i["category"] == "missing_attr"]
                if missing:
                    with st.expander(
                        f"🖼️ {len(missing)} slides with images but no attribution",
                        expanded=True,
                    ):
                        st.markdown(
                            "These slides have images but no "
                            "'Source: ...' text. Consider adding attribution:"
                        )
                        for i in missing:
                            st.markdown(
                                f"<div class='change-item change-missing_attr'>"
                                f"Slide {i['slide']}: {i['description']}</div>",
                                unsafe_allow_html=True,
                            )

                # Cross-reference issues
                xref_warnings = [i for i in report["issues"]
                                 if i["category"] == "cross_ref"
                                 and i["severity"] == "warning"]
                xref_info = [i for i in report["issues"]
                             if i["category"] == "cross_ref"
                             and i["severity"] == "info"]

                if xref_warnings:
                    with st.expander(
                        f"🔗 {len(xref_warnings)} orphaned citations "
                        f"(cited but not in references)",
                        expanded=True,
                    ):
                        for i in xref_warnings:
                            detail = i['description']
                            if i.get('original'):
                                detail += f" — `{i['original']}`"
                            st.markdown(
                                f"<div class='change-item change-cross_ref'>"
                                f"Slide {i['slide']}: {detail}</div>",
                                unsafe_allow_html=True,
                            )

                if xref_info:
                    with st.expander(
                        f"📚 {len(xref_info)} references never cited in slides",
                        expanded=False,
                    ):
                        st.caption(
                            "These are in the reference list but never cited "
                            "in slide text. They may be 'further reading' — "
                            "review and decide."
                        )
                        for i in xref_info:
                            st.markdown(
                                f"<div class='change-item change-cross_ref'>"
                                f"Slide {i['slide']}: {i['description']}</div>",
                                unsafe_allow_html=True,
                            )

                # Citation formatting issues
                cite_issues = [i for i in report["issues"]
                               if i["category"] == "citation"]
                if cite_issues:
                    with st.expander(
                        f"📝 {len(cite_issues)} citation formatting issues",
                        expanded=False,
                    ):
                        for i in cite_issues:
                            detail = i['description']
                            if i.get('original') and i.get('suggested'):
                                detail += (f"<br/>"
                                           f"<code>{i['original']}</code> → "
                                           f"<code>{i['suggested']}</code>")
                            st.markdown(
                                f"<div class='change-item change-citation'>"
                                f"Slide {i['slide']}: {detail}</div>",
                                unsafe_allow_html=True,
                            )

                # Attribution formatting issues
                attr_issues = [i for i in report["issues"]
                               if i["category"] == "attribution"]
                if attr_issues:
                    with st.expander(
                        f"🏷️ {len(attr_issues)} attribution formatting issues",
                        expanded=False,
                    ):
                        for i in attr_issues:
                            detail = i['description']
                            if i.get('original') and i.get('suggested'):
                                detail += (f"<br/>"
                                           f"<code>{i['original']}</code> → "
                                           f"<code>{i['suggested']}</code>")
                            st.markdown(
                                f"<div class='change-item change-attribution'>"
                                f"Slide {i['slide']}: {detail}</div>",
                                unsafe_allow_html=True,
                            )

                # Reference formatting issues
                ref_issues = [i for i in report["issues"]
                              if i["category"] == "reference"]
                if ref_issues:
                    with st.expander(
                        f"📖 {len(ref_issues)} reference formatting issues",
                        expanded=False,
                    ):
                        for i in ref_issues:
                            detail = i['description']
                            if i.get('original'):
                                detail += f"<br/><code>{i['original']}</code>"
                            if i.get('suggested'):
                                detail += f" → <code>{i['suggested']}</code>"
                            st.markdown(
                                f"<div class='change-item change-reference'>"
                                f"Slide {i['slide']}: {detail}</div>",
                                unsafe_allow_html=True,
                            )


# ═══════════════════════════════════════════════════════════════════════
# TAB 4: Full Compliance Check (Combined Pipeline)
# ═══════════════════════════════════════════════════════════════════════

with tab4:
    st.markdown("#### Upload a slide deck for a full compliance check")

    st.markdown("""
    <div class="info-box">
        <strong>What this does:</strong> Runs all compliance checks in sequence:
        <ol style="margin: 0.5rem 0 0 1.2rem; font-size: 0.9rem;">
            <li><strong>Layout Auto-Apply</strong> — rebuilds each slide using the correct UQ template layout</li>
            <li><strong>Brand Fixer</strong> — fonts, colours, tables, footers, headings, bullets</li>
            <li><strong>Reference Checker</strong> — citations, reference lists, image attributions</li>
            <li><strong>Image Audit</strong> — copyright risk classification for every image</li>
        </ol>
        <br/>
        You'll get a single fixed PPTX (layout + brand + reference fixes applied), a self-contained
        HTML image audit report, and a unified compliance summary — all from one upload.
    </div>
    """, unsafe_allow_html=True)

    combo_file = st.file_uploader(
        "Choose a .pptx file",
        type=["pptx"],
        key="combo_upload",
        help="Upload the slide deck you want to run through the full pipeline",
    )

    if combo_file is not None:
        file_size_mb = len(combo_file.getvalue()) / (1024 * 1024)
        st.caption(f"Uploaded: **{combo_file.name}** ({file_size_mb:.1f} MB)")

        col_opt1, col_opt2, col_opt3 = st.columns(3)
        with col_opt1:
            combo_img_limit = st.number_input(
                "Limit images (0 = all)",
                min_value=0,
                max_value=500,
                value=0,
                key="combo_img_limit",
                help="Limit image audit to first N images (0 = audit all). Useful for testing.",
            )
        with col_opt2:
            combo_skip_images = st.checkbox(
                "Skip image audit (brand + refs only)",
                key="combo_skip_images",
                help="Run brand fixer and reference checker only — no AI image classification.",
            )
        with col_opt3:
            combo_layout_mode = st.selectbox(
                "Template reformat",
                options=[
                    "v5 — Smart reformat + AI QA (recommended, uses API)",
                    "v4 — Smart reformat (free, no AI)",
                    "Skip",
                    "v2 — Legacy recipe engine",
                ],
                index=0,
                key="combo_layout_mode",
                help="v5: v4 + auto-fit + Claude Vision QA per slide (~$0.05/slide). "
                     "v4: Extracts content and injects into new template (free). "
                     "Skip: don't apply layouts. v2: Legacy recipe-based engine. "
                     "AI Vision: Claude analyses each slide image for best layout (most accurate, uses API credits).",
            )

        if st.button("Run Full Compliance Check", type="primary", key="run_combo"):
            # Validate API key if needed (image audit OR v5 AI QA)
            api_key = get_api_key()
            _needs_key_for_images = not combo_skip_images
            _needs_key_for_v5 = "v5" in combo_layout_mode
            if (_needs_key_for_images or _needs_key_for_v5) and not api_key:
                _reason = []
                if _needs_key_for_v5:
                    _reason.append("v5 AI QA")
                if _needs_key_for_images:
                    _reason.append("image audit")
                st.error(
                    f"An Anthropic API key is required for {' and '.join(_reason)}. "
                    "Ask Sean to set it up in Streamlit secrets, "
                    "or switch to v4 / tick 'Skip image audit' to run without API calls."
                )
                st.stop()

            progress_bar = st.progress(0)
            status_text = st.empty()

            def combo_progress(pct, msg):
                progress_bar.progress(min(pct / 100, 1.0))
                status_text.markdown(f"**{msg}**")

            try:
                skip_layout = combo_layout_mode == "Skip"
                if "v5" in combo_layout_mode:
                    layout_engine = "v5"
                elif "v4" in combo_layout_mode:
                    layout_engine = "v4"
                else:
                    layout_engine = "v2"
                skip_layout_vision = True

                # v5 needs API key for QA; also pass if image audit needs it
                needs_api = (layout_engine == "v5") or (not combo_skip_images)
                effective_api_key = api_key if needs_api else None

                results = run_pipeline(
                    pptx_bytes=combo_file.getvalue(),
                    filename=combo_file.name,
                    api_key=effective_api_key,
                    image_limit=combo_img_limit if combo_img_limit > 0 else None,
                    skip_image_audit=combo_skip_images,
                    skip_layout=skip_layout,
                    skip_layout_vision=skip_layout_vision,
                    layout_engine=layout_engine,
                    progress_callback=combo_progress,
                )

                progress_bar.progress(1.0)
                status_text.empty()

                # Log cost if image audit ran
                ir = results["image_report"]
                if ir and ir.get("cost_usd"):
                    log_cost(
                        tool="Full Compliance Check",
                        filename=combo_file.name,
                        num_images=ir.get("total_images", 0),
                        input_tokens=ir.get("tokens", {}).get("input", 0),
                        output_tokens=ir.get("tokens", {}).get("output", 0),
                        cost_usd=ir["cost_usd"],
                    )

                # Log layout cost if applicable
                lr = results.get("layout_report")
                if lr and lr.get("cost_usd"):
                    log_cost(
                        tool="Layout Auto-Apply",
                        filename=combo_file.name,
                        num_images=lr.get("total_slides", 0),
                        input_tokens=lr.get("tokens", {}).get("input", 0),
                        output_tokens=lr.get("tokens", {}).get("output", 0),
                        cost_usd=lr["cost_usd"],
                    )

                # Store in session state
                st.session_state["combo_result"] = {
                    "output_bytes": results["output_bytes"],
                    "brand_report": results["brand_report"],
                    "ref_report": results["ref_report"],
                    "image_report": results["image_report"],
                    "image_html": results["image_html"],
                    "image_data": results["image_data"],
                    "layout_report": results.get("layout_report"),
                    "design_report": results.get("design_report"),
                    "qa_report": results.get("qa_report"),
                    "summary": results["summary"],
                    "source_name": combo_file.name,
                    "fixed_name": combo_file.name.replace(".pptx", "_COMPLIANT.pptx"),
                    "skip_images": combo_skip_images,
                }

            except Exception as e:
                st.error(f"Something went wrong: {e}")
                with st.expander("Error details"):
                    import traceback
                    st.code(traceback.format_exc())

        # ── Display results from session state ──
        if "combo_result" in st.session_state:
            cr = st.session_state["combo_result"]
            summary = cr["summary"]
            brand_report = cr["brand_report"]
            ref_report = cr["ref_report"]
            image_report = cr["image_report"]
            layout_report = cr.get("layout_report")

            st.markdown("---")
            st.markdown("### Compliance Summary")

            # ── Overview stats ──
            cols = st.columns(5)
            cols[0].markdown(
                f"<div class='stat-card'>"
                f"<div class='number' style='color:#51247A;'>{summary['num_slides']}</div>"
                f"<div class='label'>Slides</div></div>",
                unsafe_allow_html=True,
            )

            # Layout stats
            if summary.get("layout"):
                layout_err = summary["layout"].get("error")
                if layout_err:
                    cols[1].markdown(
                        f"<div class='stat-card'>"
                        f"<div class='number' style='color:#E62645;'>!</div>"
                        f"<div class='label'>Layout error</div></div>",
                        unsafe_allow_html=True,
                    )
                    st.warning(f"Layout auto-apply failed: {layout_err}")
                else:
                    layout_changed = summary["layout"].get("changed", 0)
                    cols[1].markdown(
                        f"<div class='stat-card'>"
                        f"<div class='number' style='color:#962A8B;'>{layout_changed}</div>"
                        f"<div class='label'>Layouts fixed</div></div>",
                        unsafe_allow_html=True,
                    )
            else:
                cols[1].markdown(
                    f"<div class='stat-card'>"
                    f"<div class='number' style='color:#999;'>—</div>"
                    f"<div class='label'>Layouts (skipped)</div></div>",
                    unsafe_allow_html=True,
                )

            cols[2].markdown(
                f"<div class='stat-card'>"
                f"<div class='number' style='color:#4085C6;'>"
                f"{summary['brand']['total_changes']}</div>"
                f"<div class='label'>Brand fixes</div></div>",
                unsafe_allow_html=True,
            )
            cols[3].markdown(
                f"<div class='stat-card'>"
                f"<div class='number' style='color:#16A34A;'>"
                f"{summary['references']['total_changes']}</div>"
                f"<div class='label'>Ref fixes</div></div>",
                unsafe_allow_html=True,
            )
            if summary.get("images"):
                cols[4].markdown(
                    f"<div class='stat-card'>"
                    f"<div class='number' style='color:#D97706;'>"
                    f"{summary['images']['total_images']}</div>"
                    f"<div class='label'>Images audited</div></div>",
                    unsafe_allow_html=True,
                )
            else:
                cols[4].markdown(
                    f"<div class='stat-card'>"
                    f"<div class='number' style='color:#999;'>—</div>"
                    f"<div class='label'>Images (skipped)</div></div>",
                    unsafe_allow_html=True,
                )

            # ── Download section ──
            st.markdown("---")
            st.markdown("### Downloads")

            dl_cols = st.columns(3 if not cr["skip_images"] else 2)

            with dl_cols[0]:
                st.download_button(
                    label=f"Download {cr['fixed_name']}",
                    data=cr["output_bytes"],
                    file_name=cr["fixed_name"],
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    type="primary",
                )
                st.caption("Brand + reference fixes applied")

            with dl_cols[1]:
                json_summary = {
                    "summary": summary,
                    "brand": brand_report,
                    "references": ref_report,
                }
                # Include v5 engine reports if present
                if cr.get("layout_report"):
                    json_summary["layout_report"] = cr["layout_report"]
                if cr.get("design_report"):
                    json_summary["design_report"] = cr["design_report"]
                if cr.get("qa_report"):
                    json_summary["qa_report"] = cr["qa_report"]
                if image_report:
                    # Include image report without raw classification data to keep JSON manageable
                    json_summary["images"] = {
                        "total_images": image_report.get("total_images", 0),
                        "risk_counts": image_report.get("risk_counts", {}),
                    }
                json_name = cr["source_name"].replace(".pptx", "_compliance_report.json")
                st.download_button(
                    label="Download JSON Report",
                    data=json.dumps(json_summary, indent=2, default=str),
                    file_name=json_name,
                    mime="application/json",
                )
                st.caption("Full compliance data")

            if not cr["skip_images"] and cr.get("image_html"):
                with dl_cols[2]:
                    html_name = cr["source_name"].replace(".pptx", "_image_audit.html")
                    st.download_button(
                        label="Download Image Audit (HTML)",
                        data=cr["image_html"],
                        file_name=html_name,
                        mime="text/html",
                    )
                    st.caption("Self-contained with embedded images")

            # ── Layout Auto-Apply Details ──
            if layout_report and "results" in layout_report:
                st.markdown("---")
                changed_layouts = [r for r in layout_report["results"] if r.get("changed")]
                with st.expander(
                    f"📐 Layout Auto-Apply — {layout_report['rebuilt']} rebuilt, "
                    f"{len(changed_layouts)} layouts changed",
                    expanded=False,
                ):
                    if layout_report.get("error"):
                        st.error(f"Layout error: {layout_report['error']}")
                    elif not changed_layouts:
                        st.success("All slides already using correct UQ template layouts!", icon="✅")
                    else:
                        lcols = st.columns(3)
                        lcols[0].markdown(
                            f"<div class='stat-card'>"
                            f"<div class='number' style='color:#962A8B;'>{layout_report['rebuilt']}</div>"
                            f"<div class='label'>Rebuilt</div></div>",
                            unsafe_allow_html=True,
                        )
                        lcols[1].markdown(
                            f"<div class='stat-card'>"
                            f"<div class='number' style='color:#4085C6;'>{len(changed_layouts)}</div>"
                            f"<div class='label'>Changed</div></div>",
                            unsafe_allow_html=True,
                        )
                        lcols[2].markdown(
                            f"<div class='stat-card'>"
                            f"<div class='number' style='color:#E62645;'>{layout_report.get('failed', 0)}</div>"
                            f"<div class='label'>Failed</div></div>",
                            unsafe_allow_html=True,
                        )

                        st.markdown("**Layout changes:**")
                        for r in changed_layouts:
                            conf = f"{r['confidence']:.0%}" if r.get('confidence') else "?"
                            st.markdown(
                                f"<div class='change-item change-colour'>"
                                f"Slide {r['slide']}: "
                                f"<code>{r['from']}</code> → <code>{r['to']}</code> "
                                f"(confidence: {conf})</div>",
                                unsafe_allow_html=True,
                            )

            # ── Design Analysis Report ──
            design_report = cr.get("design_report")
            if design_report:
                st.markdown("---")
                with st.expander(
                    f"🔍 Design Analysis — {len(design_report)} slides flagged for review",
                    expanded=len(design_report) > 0,
                ):
                    if not design_report:
                        st.success("No design issues detected!", icon="✅")
                    else:
                        for item in design_report:
                            slide_num = item["slide"]
                            slide_type = item.get("type", "unknown")
                            flags = item.get("flags", [])
                            for flag in flags:
                                # Colour-code by severity
                                if flag.startswith("OVERSET"):
                                    st.markdown(
                                        f"<div class='change-item' style='border-left-color:#E62645;'>"
                                        f"<strong>Slide {slide_num}</strong> ({slide_type}): {flag}</div>",
                                        unsafe_allow_html=True,
                                    )
                                elif flag.startswith("DENSE") or flag.startswith("IMG_TEXT"):
                                    st.markdown(
                                        f"<div class='change-item' style='border-left-color:#D97706;'>"
                                        f"<strong>Slide {slide_num}</strong> ({slide_type}): {flag}</div>",
                                        unsafe_allow_html=True,
                                    )
                                else:
                                    st.markdown(
                                        f"<div class='change-item' style='border-left-color:#4085C6;'>"
                                        f"<strong>Slide {slide_num}</strong> ({slide_type}): {flag}</div>",
                                        unsafe_allow_html=True,
                                    )

            # ── AI Quality Assessment (v5) ──
            qa_report = cr.get("qa_report")
            if qa_report and qa_report.get("qa_summary"):
                qs = qa_report["qa_summary"]
                st.markdown("---")
                with st.expander(
                    f"🤖 AI Quality Assessment — "
                    f"{qs['auto_approved']} approved, "
                    f"{qs['needs_review']} review, "
                    f"{qs['needs_manual_fix']} fix "
                    f"(avg score: {qs['average_quality_score']}/10, "
                    f"${qs['total_cost_usd']:.2f})",
                    expanded=True,
                ):
                    # Summary stats
                    qcols = st.columns(5)
                    qa_stats = [
                        ("Approved", qs["auto_approved"], "#16A34A"),
                        ("Review", qs["needs_review"], "#D97706"),
                        ("Fix", qs["needs_manual_fix"], "#E62645"),
                        ("Overflow", qs["text_overflow_count"], "#962A8B"),
                        ("Avg Score", f"{qs['average_quality_score']}/10", "#4085C6"),
                    ]
                    for col, (label, val, colour) in zip(qcols, qa_stats):
                        col.markdown(
                            f"<div class='stat-card'>"
                            f"<div class='number' style='color:{colour};'>{val}</div>"
                            f"<div class='label'>{label}</div></div>",
                            unsafe_allow_html=True,
                        )

                    # Before/after comparisons
                    comparisons = qa_report.get("comparisons", [])
                    qa_results_list = qa_report.get("qa_results", [])

                    if comparisons:
                        st.markdown("### Slide-by-slide review")

                        # Filter controls
                        qa_filter = st.radio(
                            "Show:",
                            ["All", "Needs review", "Needs fix", "Overflow"],
                            horizontal=True, key="qa_filter",
                        )

                        for comp in comparisons:
                            # Apply filter
                            if qa_filter == "Needs review" and comp["recommendation"] != "needs_review":
                                continue
                            if qa_filter == "Needs fix" and comp["recommendation"] != "needs_manual_fix":
                                continue
                            if qa_filter == "Overflow" and not comp.get("text_overflow"):
                                continue

                            rec = comp["recommendation"]
                            icon = {"auto_approve": "✅", "needs_review": "⚠️",
                                    "needs_manual_fix": "❌"}.get(rec, "?")
                            score = comp.get("quality_score", "?")

                            with st.expander(
                                f"{icon} Slide {comp['slide_number']} — "
                                f"{comp['slide_type']} → {comp['target_layout']} "
                                f"(score: {score}/10)",
                                expanded=(rec != "auto_approve"),
                            ):
                                # Before/after images side by side
                                img_cols = st.columns(2)
                                with img_cols[0]:
                                    st.caption("**Original**")
                                    st.image(
                                        f"data:image/png;base64,{comp['original_b64']}",
                                        use_container_width=True,
                                    )
                                with img_cols[1]:
                                    st.caption("**Reformatted**")
                                    st.image(
                                        f"data:image/png;base64,{comp['reformatted_b64']}",
                                        use_container_width=True,
                                    )

                                # AI assessment
                                st.markdown(f"**AI assessment:** {comp.get('summary', '')}")

                                if comp.get("missing_content"):
                                    st.warning(f"Missing content: {comp['missing_content']}")
                                if comp.get("text_overflow"):
                                    st.error("Text overflow detected")
                                if comp.get("fix_suggestions"):
                                    st.markdown("**Suggestions:**")
                                    for sug in comp["fix_suggestions"]:
                                        st.markdown(f"- {sug}")

            # ── Brand Fixer Details ──
            st.markdown("---")
            with st.expander(
                f"🎨 Brand Fixer — {brand_report['total_changes']} changes",
                expanded=False,
            ):
                if brand_report["total_changes"] == 0:
                    st.success("No brand changes needed!", icon="✅")
                else:
                    stats = brand_report["summary"]
                    bcols = st.columns(7)
                    brand_stat_items = [
                        ("Font", stats.get("font", 0), "#4085C6"),
                        ("Colour", stats.get("colour", 0), "#51247A"),
                        ("Flagged", stats.get("colour_flagged", 0), "#D97706"),
                        ("Tables", stats.get("table", 0), "#16A34A"),
                        ("Footers", stats.get("footer", 0), "#962A8B"),
                        ("Headings", stats.get("heading_size", 0), "#E62645"),
                        ("Bullets", stats.get("bullet", 0), "#FBB800"),
                    ]
                    for col, (label, count, colour) in zip(bcols, brand_stat_items):
                        col.markdown(
                            f"<div class='stat-card'>"
                            f"<div class='number' style='color:{colour};'>{count}</div>"
                            f"<div class='label'>{label}</div></div>",
                            unsafe_allow_html=True,
                        )

                    # Flagged colours
                    flagged = [c for c in brand_report["changes"]
                               if c["category"] == "colour_flagged"]
                    if flagged:
                        st.markdown(f"**⚠️ {len(flagged)} colours flagged for review:**")
                        seen = set()
                        for c in flagged:
                            key = (c["slide"], c["detail"])
                            if key not in seen:
                                seen.add(key)
                                st.markdown(
                                    f"<div class='change-item change-colour_flagged'>"
                                    f"Slide {c['slide']}: {c['detail']}</div>",
                                    unsafe_allow_html=True,
                                )

            # ── Reference Checker Details ──
            with st.expander(
                f"📚 Reference Checker — {ref_report['total_issues']} issues, "
                f"{ref_report['total_changes']} fixes",
                expanded=False,
            ):
                if ref_report["total_issues"] == 0:
                    st.success("No reference issues found!", icon="✅")
                else:
                    ref_summary = ref_report["summary"]
                    rcols = st.columns(4)
                    ref_stat_items = [
                        ("Citations", ref_summary["citations_found"], "#4085C6"),
                        ("References", ref_summary["references_found"], "#51247A"),
                        ("Issues", ref_report["total_issues"], "#D97706"),
                        ("Auto-fixed", ref_report["total_changes"], "#16A34A"),
                    ]
                    for col, (label, count, colour) in zip(rcols, ref_stat_items):
                        col.markdown(
                            f"<div class='stat-card'>"
                            f"<div class='number' style='color:{colour};'>{count}</div>"
                            f"<div class='label'>{label}</div></div>",
                            unsafe_allow_html=True,
                        )

                    # Show changes
                    if ref_report["changes"]:
                        st.markdown(f"**Auto-fixes applied:**")
                        for c in ref_report["changes"]:
                            cat = c["category"]
                            st.markdown(
                                f"<div class='change-item change-{cat}'>"
                                f"Slide {c['slide']}: {c['detail']}</div>",
                                unsafe_allow_html=True,
                            )

                    # Show key issues
                    missing = [i for i in ref_report["issues"]
                               if i["category"] == "missing_attr"]
                    if missing:
                        st.markdown(f"**🖼️ {len(missing)} slides with images but no attribution**")
                        for i in missing:
                            st.markdown(
                                f"<div class='change-item change-missing_attr'>"
                                f"Slide {i['slide']}: {i['description']}</div>",
                                unsafe_allow_html=True,
                            )

                    xref_warnings = [i for i in ref_report["issues"]
                                     if i["category"] == "cross_ref"
                                     and i["severity"] == "warning"]
                    if xref_warnings:
                        st.markdown(
                            f"**🔗 {len(xref_warnings)} orphaned citations "
                            f"(cited but not in references)**"
                        )
                        for i in xref_warnings:
                            detail = i['description']
                            if i.get('original'):
                                detail += f" — `{i['original']}`"
                            st.markdown(
                                f"<div class='change-item change-cross_ref'>"
                                f"Slide {i['slide']}: {detail}</div>",
                                unsafe_allow_html=True,
                            )

            # ── Image Audit Details ──
            if not cr["skip_images"] and image_report:
                risk_counts = image_report.get("risk_counts", {})
                with st.expander(
                    f"🖼️ Image Audit — {image_report.get('total_images', 0)} images",
                    expanded=False,
                ):
                    if image_report.get("total_images", 0) == 0:
                        st.info("No images found in this presentation.")
                    elif not risk_counts:
                        st.info(image_report.get("note", "No classifications available."))
                    else:
                        icols = st.columns(5)
                        risk_items = [
                            ("Critical", risk_counts.get("CRITICAL", 0), "#DC2626"),
                            ("High", risk_counts.get("HIGH", 0), "#EA580C"),
                            ("Medium", risk_counts.get("MEDIUM", 0), "#D97706"),
                            ("Low", risk_counts.get("LOW", 0), "#16A34A"),
                            ("Clear", risk_counts.get("CLEAR", 0), "#059669"),
                        ]
                        for col, (label, count, colour) in zip(icols, risk_items):
                            col.markdown(
                                f"<div class='stat-card'>"
                                f"<div class='number' style='color:{colour};'>{count}</div>"
                                f"<div class='label'>{label}</div></div>",
                                unsafe_allow_html=True,
                            )

                        critical_high = (
                            risk_counts.get("CRITICAL", 0)
                            + risk_counts.get("HIGH", 0)
                        )
                        if critical_high > 0:
                            st.warning(
                                f"**{critical_high} image{'s' if critical_high != 1 else ''} "
                                f"flagged as Critical or High risk** — "
                                f"these likely need replacement or licence verification.",
                                icon="⚠️",
                            )

                        st.caption(
                            "Download the HTML report above for full image-by-image "
                            "detail with embedded thumbnails."
                        )
