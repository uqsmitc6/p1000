"""
UQ Slide Compliance Tool — v5 Engine
=====================================
v4 core + AI-powered Quality Assurance + Auto-fit safety net

Wraps the v4 content extraction + placeholder injection engine
and adds three new capabilities:

  1. AUTO-FIT SAFETY NET: Sets MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE on every
     text placeholder so PowerPoint auto-shrinks text that overflows.

  2. AI VISUAL QA: Renders each output slide as an image, sends it to
     Claude Vision alongside the original slide render, and asks for a
     structured quality assessment (overflow, missing content, layout
     issues, readability problems).

  3. DESIGN RECOMMENDATIONS: AI suggests per-slide improvements —
     split long slides, change layout, reposition images, etc.

The v5 engine produces the same output PPTX as v4 but with auto-fit
applied, plus a detailed AI QA report that can drive a human review UI.

Cost: ~$0.03-0.06 per slide pair (Claude Sonnet with two images).
For a 150-slide deck: ~$4.50-$9.00 total.
"""

import io
import os
import re
import json
import base64
import tempfile
import subprocess
import concurrent.futures
from dataclasses import dataclass, field
from typing import Optional, List, Dict, Any, Tuple
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import MSO_AUTO_SIZE

# v4 is the core engine — we import and wrap it
from v4_engine import (
    run_v4_pipeline,
    extract_slide_content,
    classify_slide,
    select_layout,
    inject_content,
    clear_unused_placeholders,
    analyse_design,
    SlideContent,
    SlideResult,
    _remove_all_slides,
    _deduplicate_zip,
)

try:
    import anthropic
except ImportError:
    anthropic = None

try:
    import fitz  # PyMuPDF for PDF→image rendering
except ImportError:
    fitz = None


# ============================================================================
# Configuration
# ============================================================================

CLAUDE_MODEL = "claude-sonnet-4-20250514"
MAX_PARALLEL_QA = 5  # concurrent API calls for QA
RENDER_DPI = 150  # slide render resolution (px width ≈ 1440)

# Cost tracking (approximate, Claude Sonnet pricing)
COST_PER_INPUT_TOKEN = 3.0 / 1_000_000   # $3 per 1M input tokens
COST_PER_OUTPUT_TOKEN = 15.0 / 1_000_000  # $15 per 1M output tokens


# ============================================================================
# Step 1: Auto-fit safety net
# ============================================================================

def apply_autofit_to_presentation(pptx_bytes: bytes) -> bytes:
    """Apply MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE to all text placeholders.

    This tells PowerPoint to auto-shrink text that overflows a placeholder
    when the file is opened. It's a safety net — not a fix, but it prevents
    the worst overflow issues from being visible to end users.
    """
    prs = Presentation(io.BytesIO(pptx_bytes))

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                # Only set auto-fit on placeholders, not freeform shapes
                # (freeform shapes should keep their original sizing)
                try:
                    if shape.placeholder_format is not None:
                        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except (ValueError, AttributeError):
                    pass

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ============================================================================
# Step 2: Slide rendering
# ============================================================================

def render_slides_to_images(
    pptx_bytes: bytes,
    dpi: int = RENDER_DPI,
) -> List[bytes]:
    """Render each slide of a PPTX to PNG images via LibreOffice + PyMuPDF.

    Returns a list of PNG image bytes, one per slide.
    """
    if fitz is None:
        raise ImportError("PyMuPDF (fitz) is required for slide rendering. "
                          "Install with: pip install PyMuPDF")

    # Write PPTX to temp file
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(pptx_bytes)
        tmp_pptx = tmp.name

    # Convert to PDF via LibreOffice headless
    tmp_dir = tempfile.mkdtemp()
    try:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf",
             tmp_pptx, "--outdir", tmp_dir],
            capture_output=True, text=True, timeout=300,
        )

        pdf_name = os.path.splitext(os.path.basename(tmp_pptx))[0] + ".pdf"
        pdf_path = os.path.join(tmp_dir, pdf_name)

        if not os.path.exists(pdf_path):
            raise RuntimeError(
                f"LibreOffice PDF conversion failed: {result.stderr}")

        # Render PDF pages to images
        doc = fitz.open(pdf_path)
        images = []
        # Calculate zoom to achieve desired DPI (PDF default is 72 DPI)
        zoom = dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)

        for page_idx in range(doc.page_count):
            page = doc[page_idx]
            pix = page.get_pixmap(matrix=mat)
            images.append(pix.tobytes("png"))

        doc.close()
        return images

    finally:
        # Cleanup
        os.unlink(tmp_pptx)
        for f in os.listdir(tmp_dir):
            os.unlink(os.path.join(tmp_dir, f))
        os.rmdir(tmp_dir)


# ============================================================================
# Step 3: AI Visual QA
# ============================================================================

QA_SYSTEM_PROMPT = """You are a presentation quality assurance specialist for the University of Queensland Business School. You are reviewing automatically reformatted PowerPoint slides to check that content has been correctly transferred from an old template to a new UQ branded template.

The new template uses:
- Purple colour scheme (UQ purple #51247A as primary)
- White backgrounds on content slides, purple backgrounds on covers/section dividers/thank you slides
- Title at top, content below, UQ logo in top-right corner
- Clean, professional academic styling"""

QA_SLIDE_PROMPT = """Compare these two slide images. The FIRST image is the ORIGINAL slide (old branding). The SECOND image is the REFORMATTED slide (new UQ template).

Slide {slide_num} of {total_slides}
Classified as: {slide_type}
Template layout used: {target_layout}
Extracted content: {char_count} characters of text, {image_count} images, {table_count} tables
Design flags from heuristic analysis: {design_flags}

Assess the reformatted slide against these criteria:

1. **CONTENT COMPLETENESS**: Is ALL text from the original present in the reformatted version? Is any content missing, truncated, or cut off?
2. **TEXT OVERFLOW**: Does any text appear to overflow its container, get cut off at edges, or overlap with other elements?
3. **IMAGE TRANSFER**: Are all meaningful images from the original present? Are they reasonably sized and positioned? (Decorative branding elements from the old template should NOT be carried over.)
4. **LAYOUT APPROPRIATENESS**: Is the chosen layout ({target_layout}) a good match for this content? Would a different layout work better?
5. **READABILITY**: Is text readable? Good contrast against background? Appropriate font sizes?
6. **OVERALL QUALITY**: On a scale of 1-10, how well does the reformatted slide represent the original content?

Respond in this exact JSON format:
{{
  "content_complete": true/false,
  "missing_content": "description of what's missing, or null",
  "text_overflow": true/false,
  "overflow_details": "description of overflow, or null",
  "images_correct": true/false,
  "image_issues": "description, or null",
  "layout_appropriate": true/false,
  "suggested_layout": "alternative layout name, or null",
  "readability_ok": true/false,
  "readability_issues": "description, or null",
  "quality_score": 1-10,
  "recommendation": "auto_approve" / "needs_review" / "needs_manual_fix",
  "fix_suggestions": ["list", "of", "specific", "actionable", "suggestions"],
  "summary": "one-sentence summary of the slide's quality"
}}"""


@dataclass
class QAResult:
    """Result of AI quality assessment for a single slide."""
    slide_number: int
    content_complete: bool = True
    missing_content: Optional[str] = None
    text_overflow: bool = False
    overflow_details: Optional[str] = None
    images_correct: bool = True
    image_issues: Optional[str] = None
    layout_appropriate: bool = True
    suggested_layout: Optional[str] = None
    readability_ok: bool = True
    readability_issues: Optional[str] = None
    quality_score: int = 7
    recommendation: str = "auto_approve"  # auto_approve, needs_review, needs_manual_fix
    fix_suggestions: List[str] = field(default_factory=list)
    summary: str = ""
    # Cost tracking
    input_tokens: int = 0
    output_tokens: int = 0
    cost_usd: float = 0.0
    error: Optional[str] = None


def run_ai_qa(
    client,
    original_images: List[bytes],
    reformatted_images: List[bytes],
    slide_results: List[SlideResult],
    total_slides: int,
    max_parallel: int = MAX_PARALLEL_QA,
    progress_callback=None,
) -> List[QAResult]:
    """Run AI visual QA on all slides using Claude Vision.

    Sends original + reformatted image pairs to Claude for comparison.
    Returns a QAResult for each slide.
    """
    qa_results = []

    def progress(detail, pct):
        if progress_callback:
            progress_callback("qa", detail, pct)

    # Build QA tasks
    tasks = []
    for i in range(min(len(original_images), len(reformatted_images))):
        sr = slide_results[i] if i < len(slide_results) else None
        tasks.append((i, original_images[i], reformatted_images[i], sr))

    # Process in parallel batches
    completed = 0
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_parallel) as executor:
        futures = {}
        for task in tasks:
            idx, orig_img, reformat_img, sr = task
            future = executor.submit(
                _qa_single_slide, client, idx, orig_img, reformat_img,
                sr, total_slides,
            )
            futures[future] = idx

        for future in concurrent.futures.as_completed(futures):
            idx = futures[future]
            completed += 1
            progress(
                f"QA slide {completed}/{len(tasks)}...",
                completed / len(tasks),
            )
            try:
                result = future.result()
                qa_results.append(result)
            except Exception as e:
                qa_results.append(QAResult(
                    slide_number=idx + 1,
                    error=str(e),
                    recommendation="needs_review",
                    summary=f"QA failed: {e}",
                ))

    # Sort by slide number
    qa_results.sort(key=lambda r: r.slide_number)
    return qa_results


def _qa_single_slide(
    client,
    slide_idx: int,
    original_png: bytes,
    reformatted_png: bytes,
    slide_result: Optional[SlideResult],
    total_slides: int,
) -> QAResult:
    """Run QA on a single slide pair."""
    slide_num = slide_idx + 1

    # Build prompt with slide metadata
    sr = slide_result or SlideResult(
        slide_number=slide_num, source_layout="Unknown",
        slide_type="unknown", target_layout="Unknown",
    )

    prompt = QA_SLIDE_PROMPT.format(
        slide_num=slide_num,
        total_slides=total_slides,
        slide_type=sr.slide_type,
        target_layout=sr.target_layout,
        char_count=sr.text_chars,
        image_count=sr.image_count,
        table_count=sr.table_count,
        design_flags=", ".join(sr.design_flags) if sr.design_flags else "None",
    )

    # Encode images
    orig_b64 = base64.b64encode(original_png).decode("utf-8")
    reformat_b64 = base64.b64encode(reformatted_png).decode("utf-8")

    try:
        response = client.messages.create(
            model=CLAUDE_MODEL,
            max_tokens=1024,
            system=QA_SYSTEM_PROMPT,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "ORIGINAL slide (old branding):",
                        },
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": "image/png",
                                "data": orig_b64,
                            },
                        },
                        {
                            "type": "text",
                            "text": "REFORMATTED slide (new UQ template):",
                        },
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": "image/png",
                                "data": reformat_b64,
                            },
                        },
                        {
                            "type": "text",
                            "text": prompt,
                        },
                    ],
                }
            ],
        )

        # Parse response
        raw_text = response.content[0].text.strip()
        input_tokens = response.usage.input_tokens
        output_tokens = response.usage.output_tokens
        cost = (input_tokens * COST_PER_INPUT_TOKEN +
                output_tokens * COST_PER_OUTPUT_TOKEN)

        # Extract JSON from response (may have markdown fencing)
        json_match = re.search(r'\{[\s\S]*\}', raw_text)
        if not json_match:
            return QAResult(
                slide_number=slide_num,
                error=f"Could not parse JSON from response: {raw_text[:200]}",
                recommendation="needs_review",
                input_tokens=input_tokens,
                output_tokens=output_tokens,
                cost_usd=cost,
            )

        data = json.loads(json_match.group())

        return QAResult(
            slide_number=slide_num,
            content_complete=data.get("content_complete", True),
            missing_content=data.get("missing_content"),
            text_overflow=data.get("text_overflow", False),
            overflow_details=data.get("overflow_details"),
            images_correct=data.get("images_correct", True),
            image_issues=data.get("image_issues"),
            layout_appropriate=data.get("layout_appropriate", True),
            suggested_layout=data.get("suggested_layout"),
            readability_ok=data.get("readability_ok", True),
            readability_issues=data.get("readability_issues"),
            quality_score=data.get("quality_score", 5),
            recommendation=data.get("recommendation", "needs_review"),
            fix_suggestions=data.get("fix_suggestions", []),
            summary=data.get("summary", ""),
            input_tokens=input_tokens,
            output_tokens=output_tokens,
            cost_usd=cost,
        )

    except json.JSONDecodeError as e:
        return QAResult(
            slide_number=slide_num,
            error=f"JSON parse error: {e}",
            recommendation="needs_review",
        )
    except Exception as e:
        return QAResult(
            slide_number=slide_num,
            error=f"API error: {e}",
            recommendation="needs_review",
        )


# ============================================================================
# Step 4: Before/After image pairs for human review
# ============================================================================

def generate_comparison_data(
    original_images: List[bytes],
    reformatted_images: List[bytes],
    qa_results: List[QAResult],
    slide_results: List[SlideResult],
) -> List[Dict[str, Any]]:
    """Generate comparison data for the Streamlit review UI.

    Each entry contains base64 images for before/after display,
    plus QA results and metadata.
    """
    comparisons = []

    for i in range(min(len(original_images), len(reformatted_images))):
        qa = qa_results[i] if i < len(qa_results) else None
        sr = slide_results[i] if i < len(slide_results) else None

        comparisons.append({
            "slide_number": i + 1,
            "original_b64": base64.b64encode(original_images[i]).decode("utf-8"),
            "reformatted_b64": base64.b64encode(reformatted_images[i]).decode("utf-8"),
            "slide_type": sr.slide_type if sr else "unknown",
            "source_layout": sr.source_layout if sr else "unknown",
            "target_layout": sr.target_layout if sr else "unknown",
            "quality_score": qa.quality_score if qa else None,
            "recommendation": qa.recommendation if qa else "unknown",
            "summary": qa.summary if qa else "",
            "fix_suggestions": qa.fix_suggestions if qa else [],
            "design_flags": sr.design_flags if sr else [],
            "text_overflow": qa.text_overflow if qa else False,
            "content_complete": qa.content_complete if qa else None,
            "missing_content": qa.missing_content if qa else None,
        })

    return comparisons


# ============================================================================
# Main v5 Pipeline
# ============================================================================

def run_v5_pipeline(
    source_pptx_bytes: bytes,
    template_path: str,
    api_key: Optional[str] = None,
    skip_ai_qa: bool = False,
    max_parallel_qa: int = MAX_PARALLEL_QA,
    progress_callback=None,
) -> dict:
    """Run the v5 pipeline: v4 reformat + auto-fit + AI QA.

    Args:
        source_pptx_bytes: Raw bytes of the source PPTX
        template_path: Path to the UQ template PPTX
        api_key: Anthropic API key (required for AI QA)
        skip_ai_qa: If True, skip the AI QA pass (just do v4 + auto-fit)
        max_parallel_qa: Max concurrent API calls for QA
        progress_callback: Optional callable(step, detail, pct)

    Returns dict with:
        output_pptx_bytes: The reformatted PPTX with auto-fit applied
        v4_results: Per-slide results from v4 engine
        v4_summary: Summary from v4 engine
        design_report: Heuristic design flags from v4
        qa_results: Per-slide AI QA results (or None if skipped)
        qa_summary: QA summary stats (or None)
        comparisons: Before/after image data for review UI (or None)
        cost: Total API cost in USD
    """
    def progress(step, detail, pct):
        if progress_callback:
            progress_callback(step, detail, pct)

    total_cost = 0.0
    total_input_tokens = 0
    total_output_tokens = 0

    # ── Phase 1: Run v4 core engine ──────────────────────────────────
    progress("v4", "Running content extraction and injection...", 0.05)

    def v4_progress(step, detail, v4_pct):
        # Map v4 progress (0-1) to overall progress (0.05-0.40)
        mapped = 0.05 + v4_pct * 0.35
        progress("v4", detail, mapped)

    v4_result = run_v4_pipeline(
        source_pptx_bytes=source_pptx_bytes,
        template_path=template_path,
        progress_callback=v4_progress,
    )

    output_bytes = v4_result["output_pptx_bytes"]
    v4_results = v4_result["results"]
    v4_summary = v4_result["summary"]
    design_report = v4_result["design_report"]

    # ── Phase 2: Apply auto-fit safety net ───────────────────────────
    progress("autofit", "Applying auto-fit text shrink safety net...", 0.42)
    output_bytes = apply_autofit_to_presentation(output_bytes)
    # Re-deduplicate after modification
    output_bytes = _deduplicate_zip(output_bytes)

    # ── Phase 3: AI Visual QA ────────────────────────────────────────
    qa_results = None
    qa_summary = None
    comparisons = None

    if skip_ai_qa:
        progress("qa", "AI QA skipped (skip_ai_qa=True)", 0.90)
    elif not api_key:
        progress("qa", "AI QA skipped — no API key provided", 0.90)
    elif anthropic is None:
        progress("qa", "AI QA skipped — 'anthropic' package not installed (add to requirements.txt)", 0.90)
    elif fitz is None:
        progress("qa", "AI QA skipped — 'PyMuPDF' package not installed (add PyMuPDF to requirements.txt)", 0.90)

    if not skip_ai_qa and api_key:
        if anthropic is None:
            pass  # Already logged above
        elif fitz is None:
            pass  # Already logged above
        else:
            # Render source slides
            progress("render", "Rendering original slides...", 0.45)
            try:
                original_images = render_slides_to_images(source_pptx_bytes)
            except Exception as e:
                progress("render", f"Failed to render originals: {e}", 0.50)
                original_images = []

            # Render output slides
            progress("render", "Rendering reformatted slides...", 0.50)
            try:
                reformatted_images = render_slides_to_images(output_bytes)
            except Exception as e:
                progress("render", f"Failed to render reformatted: {e}", 0.55)
                reformatted_images = []

            if original_images and reformatted_images:
                # Run AI QA
                progress("qa", "Starting AI quality assessment...", 0.55)
                client = anthropic.Anthropic(api_key=api_key)

                def qa_progress(step, detail, qa_pct):
                    mapped = 0.55 + qa_pct * 0.35
                    progress("qa", detail, mapped)

                qa_results = run_ai_qa(
                    client=client,
                    original_images=original_images,
                    reformatted_images=reformatted_images,
                    slide_results=v4_results,
                    total_slides=v4_summary["total_slides"],
                    max_parallel=max_parallel_qa,
                    progress_callback=qa_progress,
                )

                # Aggregate QA stats
                total_input_tokens = sum(r.input_tokens for r in qa_results)
                total_output_tokens = sum(r.output_tokens for r in qa_results)
                total_cost = sum(r.cost_usd for r in qa_results)

                auto_approved = sum(
                    1 for r in qa_results if r.recommendation == "auto_approve")
                needs_review = sum(
                    1 for r in qa_results if r.recommendation == "needs_review")
                needs_fix = sum(
                    1 for r in qa_results if r.recommendation == "needs_manual_fix")
                avg_score = (sum(r.quality_score for r in qa_results)
                             / len(qa_results)) if qa_results else 0

                qa_summary = {
                    "total_assessed": len(qa_results),
                    "auto_approved": auto_approved,
                    "needs_review": needs_review,
                    "needs_manual_fix": needs_fix,
                    "average_quality_score": round(avg_score, 1),
                    "text_overflow_count": sum(
                        1 for r in qa_results if r.text_overflow),
                    "content_incomplete_count": sum(
                        1 for r in qa_results if not r.content_complete),
                    "layout_inappropriate_count": sum(
                        1 for r in qa_results if not r.layout_appropriate),
                    "total_cost_usd": round(total_cost, 4),
                    "total_input_tokens": total_input_tokens,
                    "total_output_tokens": total_output_tokens,
                }

                # Generate comparison data for UI
                progress("compare", "Generating comparison data...", 0.92)
                comparisons = generate_comparison_data(
                    original_images, reformatted_images,
                    qa_results, v4_results,
                )
            else:
                progress("qa", "Skipping AI QA — rendering failed", 0.90)

    # ── Done ─────────────────────────────────────────────────────────
    progress("done", "Complete", 1.0)

    return {
        "output_pptx_bytes": output_bytes,
        "v4_results": v4_results,
        "v4_summary": v4_summary,
        "design_report": design_report,
        "qa_results": qa_results,
        "qa_summary": qa_summary,
        "comparisons": comparisons,
        "cost": {
            "total_usd": round(total_cost, 4),
            "input_tokens": total_input_tokens,
            "output_tokens": total_output_tokens,
        },
    }


# ============================================================================
# CLI
# ============================================================================

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 3:
        print("Usage: python v5_engine.py <source.pptx> <template.pptx> [output.pptx]")
        print()
        print("Environment variables:")
        print("  ANTHROPIC_API_KEY  — required for AI QA pass")
        print("  SKIP_AI_QA=1       — skip AI QA (just v4 + auto-fit)")
        sys.exit(1)

    source_path = sys.argv[1]
    template_path = sys.argv[2]
    output_path = sys.argv[3] if len(sys.argv) > 3 else "v5_output.pptx"

    api_key = os.environ.get("ANTHROPIC_API_KEY")
    skip_qa = os.environ.get("SKIP_AI_QA", "").strip() == "1"

    if not api_key and not skip_qa:
        print("Warning: No ANTHROPIC_API_KEY set. AI QA will be skipped.")
        print("Set SKIP_AI_QA=1 to suppress this warning.\n")
        skip_qa = True

    with open(source_path, "rb") as f:
        source_bytes = f.read()

    def print_progress(step, detail, pct):
        print(f"  [{pct:5.1%}] [{step:8s}] {detail}")

    result = run_v5_pipeline(
        source_pptx_bytes=source_bytes,
        template_path=template_path,
        api_key=api_key,
        skip_ai_qa=skip_qa,
        progress_callback=print_progress,
    )

    with open(output_path, "wb") as f:
        f.write(result["output_pptx_bytes"])

    # Print v4 summary
    s = result["v4_summary"]
    print(f"\n=== v4 Engine ===")
    print(f"Rebuilt: {s['success']}/{s['total_slides']}")
    print(f"Failed: {s['failed']}")
    print(f"Flagged (heuristic): {s['flagged_for_review']}")

    # Print design flags
    if result["design_report"]:
        print(f"\n=== Design Flags ===")
        for item in result["design_report"]:
            for flag in item["flags"]:
                print(f"  Slide {item['slide']}: {flag}")

    # Print QA summary
    if result["qa_summary"]:
        q = result["qa_summary"]
        print(f"\n=== AI Quality Assessment ===")
        print(f"Assessed: {q['total_assessed']} slides")
        print(f"Auto-approved: {q['auto_approved']}")
        print(f"Needs review: {q['needs_review']}")
        print(f"Needs manual fix: {q['needs_manual_fix']}")
        print(f"Average quality score: {q['average_quality_score']}/10")
        print(f"Text overflow detected: {q['text_overflow_count']}")
        print(f"Content incomplete: {q['content_incomplete_count']}")
        print(f"API cost: ${q['total_cost_usd']:.4f}")

        # Per-slide QA detail
        print(f"\n=== Per-Slide QA ===")
        for r in result["qa_results"]:
            icon = {"auto_approve": "✅", "needs_review": "⚠️",
                    "needs_manual_fix": "❌"}.get(r.recommendation, "?")
            print(f"  {icon} Slide {r.slide_number} "
                  f"(score: {r.quality_score}/10): {r.summary}")
            if r.fix_suggestions:
                for sug in r.fix_suggestions:
                    print(f"      → {sug}")

    print(f"\nOutput saved to: {output_path}")
    print(f"Total cost: ${result['cost']['total_usd']:.4f}")
