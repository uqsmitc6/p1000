"""
UQ Slide Compliance Tool — v4 Engine
=====================================
Content Extraction + Placeholder Injection

Architecture:
  1. EXTRACT: Pull structured content from each source slide
     (title, subtitle, body text, images, tables — as raw data, not shapes)
  2. CLASSIFY: Determine slide type using heuristics + optional AI vision
  3. MAP: Select the best template layout for each slide type
  4. INJECT: Place extracted content INTO the template's existing placeholders
  5. CLEAN: Clear unused placeholders so no ghost text shows through
  6. REVIEW: (hook) AI design analysis pass — flag overset text, bad layouts, etc.

Key difference from v3: We NEVER copy shapes from source to target.
Content is extracted as data and injected into template placeholders,
which means template formatting (colours, fonts, sizes) is inherited
automatically and old branding simply cannot transfer.
"""

import io
import os
import re
import zipfile
import tempfile
from dataclasses import dataclass, field
from typing import Optional, List, Dict, Tuple, Any
from collections import Counter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# ============================================================================
# Data Models — extracted content as structured data
# ============================================================================

@dataclass
class TextRun:
    """A single run of text with formatting."""
    text: str
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    font_size: Optional[float] = None  # in points
    hyperlink: Optional[str] = None


@dataclass
class Paragraph:
    """A paragraph containing one or more text runs."""
    runs: List[TextRun] = field(default_factory=list)
    level: int = 0
    alignment: Optional[int] = None  # PP_ALIGN value

    @property
    def text(self) -> str:
        return "".join(r.text for r in self.runs)

    @property
    def char_count(self) -> int:
        return len(self.text)


@dataclass
class TextBlock:
    """A block of text (from a placeholder or text box)."""
    paragraphs: List[Paragraph] = field(default_factory=list)
    source_type: str = ""  # "placeholder", "textbox", "group_child"
    source_ph_idx: Optional[int] = None
    source_ph_type: Optional[str] = None
    left: int = 0  # EMU position for spatial analysis
    top: int = 0
    width: int = 0
    height: int = 0

    @property
    def full_text(self) -> str:
        return "\n".join(p.text for p in self.paragraphs)

    @property
    def char_count(self) -> int:
        return sum(p.char_count for p in self.paragraphs)

    @property
    def is_title_like(self) -> bool:
        """Heuristic: short text near the top of the slide."""
        text = self.full_text.strip()
        if not text:
            return False
        # Short text, near top, not a bullet list
        is_short = len(text) < 120
        is_near_top = self.top < Emu(2.5 * 914400)  # top 2.5 inches
        has_no_bullets = len(self.paragraphs) <= 2
        return is_short and is_near_top and has_no_bullets


@dataclass
class ExtractedImage:
    """An image extracted from a slide."""
    image_bytes: bytes
    content_type: str  # e.g. "image/png", "image/jpeg"
    width: int = 0   # EMU
    height: int = 0  # EMU
    left: int = 0    # EMU
    top: int = 0     # EMU
    source_type: str = ""  # "placeholder", "freeform", "group_child"

    @property
    def width_inches(self) -> float:
        return self.width / 914400

    @property
    def height_inches(self) -> float:
        return self.height / 914400

    @property
    def is_logo_sized(self) -> bool:
        """Heuristic: likely a logo/branding element if small."""
        return self.width_inches < 2.5 and self.height_inches < 1.5

    @property
    def is_decorative(self) -> bool:
        """Heuristic: decorative/branding if small and in corner."""
        if not self.is_logo_sized:
            return False
        slide_w = 13.333 * 914400  # 16:9 widescreen
        slide_h = 7.5 * 914400
        in_corner = (
            (self.left < Emu(1.5 * 914400) or self.left > slide_w - Emu(3 * 914400))
            and (self.top < Emu(1.5 * 914400) or self.top > slide_h - Emu(1.5 * 914400))
        )
        return in_corner


@dataclass
class ExtractedTable:
    """A table extracted from a slide."""
    rows: List[List[str]] = field(default_factory=list)
    col_count: int = 0
    row_count: int = 0
    width: int = 0   # EMU
    height: int = 0  # EMU


@dataclass
class SlideContent:
    """All content extracted from a single source slide."""
    slide_number: int
    source_layout_name: str

    # Extracted content
    title: Optional[TextBlock] = None
    subtitle: Optional[TextBlock] = None
    body_blocks: List[TextBlock] = field(default_factory=list)
    images: List[ExtractedImage] = field(default_factory=list)
    tables: List[ExtractedTable] = field(default_factory=list)

    # Classification (set during classify step)
    slide_type: str = ""  # cover, section_divider, thank_you, content, image_content, two_column, title_only, blank, quote
    target_layout: str = ""  # template layout name to use

    # Design analysis flags (set during review step)
    design_flags: List[str] = field(default_factory=list)

    @property
    def total_text_chars(self) -> int:
        total = 0
        if self.title:
            total += self.title.char_count
        if self.subtitle:
            total += self.subtitle.char_count
        for b in self.body_blocks:
            total += b.char_count
        return total

    @property
    def content_images(self) -> List[ExtractedImage]:
        """Images that are actual content, not branding/logos."""
        return [img for img in self.images if not img.is_decorative]

    @property
    def has_substantial_text(self) -> bool:
        return self.total_text_chars > 50

    @property
    def has_images(self) -> bool:
        return len(self.content_images) > 0

    @property
    def primary_body_text(self) -> str:
        """Combined body text for injection into a single placeholder."""
        return "\n\n".join(b.full_text for b in self.body_blocks if b.full_text.strip())


# ============================================================================
# STEP 1: Content Extraction
# ============================================================================

# Old UQ branding patterns to filter out
OLD_BRAND_PATTERNS = [
    r'uq\s*business\s*school',
    r'executive\s*education',
    r'cricos.*provider',
    r'teqsa.*provider',
    r'the\s*university\s*of\s*queensland',
    r'uq\.edu\.au',
]
OLD_BRAND_RE = re.compile('|'.join(OLD_BRAND_PATTERNS), re.IGNORECASE)


def extract_slide_content(slide, slide_number: int) -> SlideContent:
    """Extract all content from a source slide as structured data.

    Returns a SlideContent object with title, subtitle, body blocks,
    images, and tables — but NO shapes, NO XML, NO positioning data
    that would carry old branding.
    """
    layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
    content = SlideContent(
        slide_number=slide_number,
        source_layout_name=layout_name,
    )

    # --- Extract from placeholders first ---
    for shape in slide.shapes:
        try:
            pf = shape.placeholder_format
            if not pf:
                raise ValueError("not a placeholder")
            # Handle broken placeholder indices (e.g. 0xFFFFFFFF)
            if pf.idx is not None and pf.idx > 100:
                # Treat as freeform — this is a corrupted placeholder
                _extract_from_freeform(shape, content)
            else:
                _extract_from_placeholder(shape, pf, content)
        except (ValueError, AttributeError):
            # Not a placeholder — handle as freeform shape
            _extract_from_freeform(shape, content)

    # --- If no title was found in placeholders, try to infer one ---
    if not content.title:
        _infer_title(content)

    return content


def _extract_from_placeholder(shape, pf, content: SlideContent):
    """Extract content from a placeholder shape."""
    ptype = str(pf.type) if pf.type is not None else ""

    # Skip structural placeholders
    if any(skip in ptype for skip in ['FOOTER', 'SLIDE_NUMBER', 'DATE_TIME']):
        return

    # Title
    if 'TITLE' in ptype or 'CENTER_TITLE' in ptype:
        if shape.has_text_frame and shape.text_frame.text.strip():
            content.title = _extract_text_block(shape, "placeholder", pf.idx, ptype)
        return

    # Subtitle
    if 'SUBTITLE' in ptype:
        if shape.has_text_frame and shape.text_frame.text.strip():
            content.subtitle = _extract_text_block(shape, "placeholder", pf.idx, ptype)
        return

    # Check for image inside placeholder
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    r_ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
    blips = shape._element.findall(f'.//{{{a_ns}}}blip')
    if blips:
        for blip in blips:
            img = _extract_image_from_blip(blip, shape, r_ns)
            if img:
                img.source_type = "placeholder"
                content.images.append(img)
        return

    # Check for table inside placeholder
    if shape.has_table:
        content.tables.append(_extract_table(shape))
        return

    # Body text
    if shape.has_text_frame and shape.text_frame.text.strip():
        block = _extract_text_block(shape, "placeholder", pf.idx, ptype)
        # Filter out old branding text
        if not _is_brand_text(block.full_text):
            content.body_blocks.append(block)


def _extract_from_freeform(shape, content: SlideContent):
    """Extract content from a non-placeholder shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    shape_type = shape.shape_type

    # --- Images ---
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        img = _extract_image_from_picture(shape)
        if img and not img.is_decorative:
            content.images.append(img)
        return

    # --- Tables ---
    if shape.has_table:
        content.tables.append(_extract_table(shape))
        return

    # --- Groups: recurse into children ---
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        _extract_from_group(shape, content)
        return

    # --- Text boxes and auto shapes with text ---
    if shape.has_text_frame and shape.text_frame.text.strip():
        block = _extract_text_block(shape, "textbox")
        text = block.full_text.strip()
        # Filter out old branding
        if _is_brand_text(text):
            return
        # Filter out very short decorative text (single chars, arrows, etc.)
        if len(text) <= 2 and not text.isdigit():
            return
        content.body_blocks.append(block)


def _extract_from_group(group_shape, content: SlideContent):
    """Recursively extract content from a grouped shape."""
    try:
        for child in group_shape.shapes:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = _extract_image_from_picture(child)
                if img:
                    img.source_type = "group_child"
                    if not img.is_decorative:
                        content.images.append(img)
            elif child.has_text_frame and child.text_frame.text.strip():
                block = _extract_text_block(child, "group_child")
                if not _is_brand_text(block.full_text):
                    content.body_blocks.append(block)
            elif hasattr(child, 'shapes'):
                # Nested group
                _extract_from_group(child, content)
    except Exception:
        pass  # Groups can be tricky — don't crash


def _extract_text_block(
    shape, source_type: str,
    ph_idx: Optional[int] = None,
    ph_type: Optional[str] = None,
) -> TextBlock:
    """Extract a TextBlock from any shape with a text frame."""
    block = TextBlock(
        source_type=source_type,
        source_ph_idx=ph_idx,
        source_ph_type=ph_type,
        left=shape.left or 0,
        top=shape.top or 0,
        width=shape.width or 0,
        height=shape.height or 0,
    )

    for src_para in shape.text_frame.paragraphs:
        para = Paragraph(
            level=src_para.level if src_para.level is not None else 0,
            alignment=src_para.alignment,
        )
        for src_run in src_para.runs:
            run = TextRun(
                text=src_run.text,
                bold=src_run.font.bold,
                italic=src_run.font.italic,
                font_size=src_run.font.size.pt if src_run.font.size else None,
            )
            # Extract hyperlink if present
            try:
                if src_run.hyperlink and src_run.hyperlink.address:
                    run.hyperlink = src_run.hyperlink.address
            except Exception:
                pass
            para.runs.append(run)
        block.paragraphs.append(para)

    return block


def _extract_image_from_picture(shape) -> Optional[ExtractedImage]:
    """Extract image bytes from a freeform picture shape."""
    try:
        image = shape.image
        return ExtractedImage(
            image_bytes=image.blob,
            content_type=image.content_type,
            width=shape.width or 0,
            height=shape.height or 0,
            left=shape.left or 0,
            top=shape.top or 0,
            source_type="freeform",
        )
    except Exception:
        return None


def _extract_image_from_blip(blip, shape, r_ns: str) -> Optional[ExtractedImage]:
    """Extract image bytes from a blip element inside a placeholder."""
    try:
        old_rId = blip.get(f'{r_ns}embed')
        if not old_rId:
            return None
        slide_part = shape._element.getparent()
        # Walk up to find the slide part
        while slide_part is not None and not hasattr(slide_part, 'part'):
            slide_part = slide_part.getparent()

        # Get from the shape's slide part
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        part = shape.part if hasattr(shape, 'part') else None
        if part is None:
            return None

        if old_rId not in part.rels:
            return None

        rel = part.rels[old_rId]
        image_part = rel.target_part

        return ExtractedImage(
            image_bytes=image_part.blob,
            content_type=image_part.content_type if hasattr(image_part, 'content_type') else 'image/png',
            width=shape.width or 0,
            height=shape.height or 0,
            left=shape.left or 0,
            top=shape.top or 0,
            source_type="placeholder",
        )
    except Exception:
        return None


def _extract_table(shape) -> ExtractedTable:
    """Extract table data as rows of strings."""
    table = shape.table
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip())
        rows.append(cells)

    return ExtractedTable(
        rows=rows,
        col_count=len(table.columns),
        row_count=len(table.rows),
        width=shape.width or 0,
        height=shape.height or 0,
    )


def _is_brand_text(text: str) -> bool:
    """Check if text is old branding that should be filtered out.

    Only filters SHORT text that is PRIMARILY branding content.
    Does NOT filter text that merely contains a UQ email address
    or domain as part of real content (e.g. contact details).
    """
    if not text:
        return True  # Empty text is effectively brand filler
    text = text.strip()
    if not text:
        return True

    # Long text is real content even if it mentions UQ
    if len(text) > 150:
        return False

    # If text contains an @ sign, it's likely contact info — keep it
    if '@' in text:
        return False

    # If text contains a phone number pattern, it's contact info
    if re.search(r'\+?\d[\d\s\-]{7,}', text):
        return False

    # Only filter if the text is PREDOMINANTLY branding
    # (i.e. the branding pattern covers most of the text)
    match = OLD_BRAND_RE.search(text)
    if not match:
        return False

    # If the match covers more than 40% of the text, it's branding
    match_len = match.end() - match.start()
    return match_len / len(text) > 0.4


def _infer_title(content: SlideContent):
    """If no title was found in placeholders, try to find one in body blocks."""
    if not content.body_blocks:
        return

    # Sort by top position — highest block that looks title-like
    candidates = [b for b in content.body_blocks if b.is_title_like]
    if candidates:
        candidates.sort(key=lambda b: b.top)
        content.title = candidates[0]
        content.body_blocks.remove(candidates[0])


# ============================================================================
# STEP 2: Slide Classification
# ============================================================================

def classify_slide(content: SlideContent, slide_idx: int, total_slides: int):
    """Classify a slide's type based on its extracted content.

    Sets content.slide_type to one of:
      cover, section_divider, thank_you, acknowledgement, references,
      contents_page, quote, content_with_image, two_column,
      content_text_only, title_only, blank
    """
    title_text = content.title.full_text.strip().lower() if content.title else ""
    is_first = slide_idx == 0
    is_last = slide_idx == total_slides - 1

    # --- Position-based overrides ---
    if is_first:
        content.slide_type = "cover"
        return

    # --- Keyword-based classification ---
    if _matches_keywords(title_text, ["thank", "thanks", "questions", "q&a", "q & a"]):
        content.slide_type = "thank_you"
        return

    # Contact details slide — has email addresses and phone numbers
    all_text = content.primary_body_text.lower()
    email_count = len(re.findall(r'[\w.-]+@[\w.-]+', all_text))
    phone_count = len(re.findall(r'\+?\d[\d\s\-]{7,}', all_text))
    if email_count >= 2 or (email_count >= 1 and phone_count >= 1):
        if _matches_keywords(title_text, ["contact", "team", "facilitator", "presenter", "instructor"]) or not content.title:
            content.slide_type = "thank_you"  # Use Thank You layout for contact slides
            return

    if _matches_keywords(title_text, ["acknowledg", "country", "traditional", "custodian",
                                       "elder", "first nation", "aboriginal", "torres strait"]):
        content.slide_type = "acknowledgement"
        return

    if _matches_keywords(title_text, ["reference", "bibliograph", "reading list", "further reading"]):
        content.slide_type = "references"
        return

    if _matches_keywords(title_text, ["content", "agenda", "overview", "outline", "today",
                                       "schedule", "program", "programme"]):
        content.slide_type = "contents_page"
        return

    # Section divider heuristics
    layout_lower = content.source_layout_name.lower()
    if "section" in layout_lower or "divider" in layout_lower:
        content.slide_type = "section_divider"
        return

    # Very short title, no body, no images → likely section divider
    if (content.title and content.title.char_count < 40
            and not content.body_blocks and not content.has_images):
        # Check if it's just a number or very short phrase
        if content.title.char_count < 15 or title_text.replace(" ", "").isdigit():
            content.slide_type = "section_divider"
            return

    # Quote detection
    body_text = content.primary_body_text
    if _looks_like_quote(body_text):
        content.slide_type = "quote"
        return

    # --- Content-based classification ---
    has_images = content.has_images
    has_text = content.has_substantial_text
    num_body_blocks = len(content.body_blocks)
    has_table = len(content.tables) > 0

    if has_table:
        content.slide_type = "content_with_table"
        return

    if not has_text and not has_images:
        content.slide_type = "blank"
        return

    if content.title and not has_text and not has_images:
        # Title only with no body
        body_text_total = sum(b.char_count for b in content.body_blocks)
        if body_text_total < 10:
            content.slide_type = "title_only"
            return

    if has_images and has_text:
        content.slide_type = "content_with_image"
        return

    if has_images and not has_text:
        content.slide_type = "image_only"
        return

    # Two-column detection: multiple body blocks side by side
    if num_body_blocks >= 2:
        blocks_sorted = sorted(content.body_blocks, key=lambda b: b.left)
        left_most = blocks_sorted[0].left
        right_most = blocks_sorted[-1].left
        # If blocks are spread horizontally, it's two-column
        if right_most - left_most > Emu(4 * 914400):  # 4+ inches apart
            content.slide_type = "two_column"
            return

    content.slide_type = "content_text_only"


def _matches_keywords(text: str, keywords: list) -> bool:
    return any(kw in text for kw in keywords)


def _looks_like_quote(text: str) -> bool:
    if not text:
        return False
    text = text.strip()
    # Starts with quote mark, or has attribution pattern
    if text.startswith('"') or text.startswith('\u201c') or text.startswith('\u2018'):
        return True
    # Short text with an attribution line (dash + name)
    lines = text.split('\n')
    if len(lines) >= 2:
        last_line = lines[-1].strip()
        if last_line.startswith('-') or last_line.startswith('\u2013') or last_line.startswith('\u2014'):
            return True
    return False


# ============================================================================
# STEP 3: Layout Selection
# ============================================================================

# Map slide_type → preferred template layout name
SLIDE_TYPE_TO_LAYOUT = {
    "cover": "Cover 1",
    "section_divider": "Section Divider",
    "thank_you": "Thank You",
    "acknowledgement": "Title and Content",
    "references": "Title and Content",
    "contents_page": "Title and Content",
    "quote": "Quote 1",
    "content_with_image": "Text with Image Half",
    "content_with_table": "Title and Content",
    "two_column": "Two Content",
    "content_text_only": "Title and Content",
    "image_only": "Title Only",
    "title_only": "Title Only",
    "blank": "Blank Branded",
}

# Refined image layout selection based on image count/size
IMAGE_LAYOUT_RULES = {
    1: {
        "large": "Text with Image Half",      # Big photo → half-half
        "medium": "Text with Image One Third",  # Medium → one third
        "small": "Title and Content",           # Small diagram → just content
    },
    2: "Two Content",
    3: "Three Column Text & Images",
    4: "Text with 4 Images",
}

# Layout name → layout variant mapping from old templates
LAYOUT_NAME_MAP = {
    # Strip numeric prefixes (1_, 2_, 3_, 8_) during lookup
}


def select_layout(content: SlideContent, available_layouts: Dict[str, Any]) -> str:
    """Select the best template layout for this slide's content.

    Uses the slide_type classification plus content analysis to pick
    the most appropriate of the 46 template layouts.
    """
    # Start with the type-based default
    preferred = SLIDE_TYPE_TO_LAYOUT.get(content.slide_type, "Title and Content")

    # Refine for image content
    if content.slide_type == "content_with_image":
        images = content.content_images
        if len(images) == 1:
            img = images[0]
            if img.width_inches > 5:
                preferred = "Text with Image Half"
            elif img.width_inches > 3:
                preferred = "Text with Image One Third"
            else:
                preferred = "Title and Content"
        elif len(images) == 2:
            preferred = "Two Content"
        elif len(images) >= 4:
            preferred = "Text with 4 Images" if "Text with 4 Images" in available_layouts else "Title and Content"
        elif len(images) == 3:
            preferred = "Three Column Text & Images" if "Three Column Text & Images" in available_layouts else "Two Content"

    # Also try to match by source layout name (old template → new)
    # BUT: type-based classification takes priority for certain slide types
    # where we KNOW the content type regardless of source layout
    TYPE_PRIORITY_OVERRIDES = {
        "cover", "thank_you", "section_divider", "references",
        "acknowledgement", "quote",
    }
    # Restrictive layouts that should NEVER override the type-based selection
    # when the slide has body content — these layouts lack body placeholders
    # and would cause content loss.
    RESTRICTIVE_LAYOUTS = {
        "Title Only", "Blank Branded", "Icons & Text", "Process Diagram",
        "Image collage", "Three Pullouts", "Multi-layout 1", "Multi-layout 2",
    }
    if content.slide_type not in TYPE_PRIORITY_OVERRIDES:
        source_name = content.source_layout_name
        mapped = _map_source_layout(source_name)
        if mapped and mapped in available_layouts:
            # Use mapped name if it's a better match than the generic type default
            if mapped != "Title and Content":  # Don't override with the generic fallback
                # Don't override to a restrictive layout if the slide has body content
                if mapped in RESTRICTIVE_LAYOUTS and content.body_blocks:
                    pass  # Keep the type-based selection
                else:
                    preferred = mapped

    # Validate layout exists in template
    if preferred not in available_layouts:
        preferred = "Title and Content"
    if preferred not in available_layouts:
        # Last resort: first available content layout
        for name in available_layouts:
            if "Content" in name or "Title" in name:
                preferred = name
                break

    content.target_layout = preferred
    return preferred


def _map_source_layout(source_name: str) -> Optional[str]:
    """Map an old template layout name to the new template layout name."""
    # Strip numeric prefix
    stripped = re.sub(r'^\d+_', '', source_name)

    # Direct mappings for known variants
    KNOWN_MAPS = {
        "Cover 1": "Cover 1",
        "Cover 2": "Cover 2",
        "Cover 3": "Cover 3",
        "Cover 4": "Thank You",
        "Section Divider": "Section Divider",
        "Section Divider 2": "Section Divider",
        "Title and Content": "Title and Content",
        "Two Content": "Two Content",
        "Title Only": "Title Only",
        "Blank Branded": "Blank Branded",
        "Text with Image One Third": "Text with Image One Third",
        "Text with Image One Third Alt": "Text with Image One Third Alt",
        "Text with Image Half": "Text with Image Half",
        "Text with Image Half Alt": "Text with Image Half Alt",
        "Text with Image Two Thirds": "Text with Image Two Thirds",
        "Text with Image Two Thirds Alt": "Text with Image Two Thirds Alt",
        "Text with Dark Purple Block": "Text with Dark Purple Block",
        "Text with Neutral Block": "Text with Neutral Block",
        "Text with Grey Block": "Text with Grey Block",
        "Quote 1": "Quote 1",
        "Quote 2": "Quote 2",
        "Thank You": "Thank You",
        "Picture with Pullout": "Picture with Pullout",
        "Picture with Caption": "Picture with Caption",
        "Contents 1": "Contents 1",
        "Contents 2": "Contents 2",
        "Title and Table": "Title and Table",
        "Three Pullouts": "Three Pullouts",
        "Icons & Text": "Icons & Text",
    }

    if source_name in KNOWN_MAPS:
        return KNOWN_MAPS[source_name]
    if stripped in KNOWN_MAPS:
        return KNOWN_MAPS[stripped]
    return None


# ============================================================================
# STEP 4: Placeholder Injection
# ============================================================================

# Template placeholder index maps by layout name
# Format: { layout_name: { role: ph_idx } }
# Roles: "title", "subtitle", "body", "body2", "picture", "picture2"
PLACEHOLDER_MAP = {
    "Cover 1": {"title": 0, "subtitle": 10, "supertitle": 11},
    "Cover 2": {"title": 0, "subtitle": 10, "supertitle": 11},
    "Cover 3": {"title": 0, "subtitle": 10, "supertitle": 11},
    "Section Divider": {"title": 0, "section_num": 11, "body": 13},
    "Title and Content": {"title": 0, "subtitle": 31, "body": 10},
    "Two Content": {"title": 0, "subtitle": 31, "body": 10, "body2": 32},
    "Title Only": {"title": 0},
    "Blank Branded": {},
    "Text with Image Half": {"title": 0, "subtitle": 31, "body": 10, "picture": 32},
    "Text with Image Half Alt": {"title": 0, "subtitle": 31, "body": 10, "picture": 32},
    "Text with Image One Third": {"title": 0, "subtitle": 31, "body": 10, "picture": 34},
    "Text with Image One Third Alt": {"title": 0, "subtitle": 31, "body": 10, "picture": 34},
    "Text with Image Two Thirds": {"title": 0, "subtitle": 31, "body": 10, "picture": 32},
    "Text with Image Two Thirds Alt": {"title": 0, "subtitle": 31, "body": 10, "picture": 32},
    "Quote 1": {"title": 0, "body": 10},
    "Quote 2": {"title": 0, "body": 10},
    "Thank You": {"title": 0, "name": 10, "role": 16, "email": 17, "phone": 18},
    "Text with Dark Purple Block": {"title": 0, "subtitle": 31, "body": 10},
    "Text with Neutral Block": {"title": 0, "subtitle": 31, "body": 10},
    "Text with Grey Block": {"title": 0, "subtitle": 31, "body": 10},
    "Three Column Text & Images": {"title": 0, "body": 10},
    "Text with 4 Images": {"title": 0, "body": 10},
    "Title and Table": {"title": 0, "subtitle": 31, "body": 10},
    "Contents 1": {"title": 0, "body": 10},
    "Contents 2": {"title": 0, "body": 10},
    "Picture with Pullout": {"title": 0, "picture": 32},
    "Picture with Caption": {"title": 0, "body": 10, "picture": 32},
    "Three Pullouts": {"title": 0},
    "Icons & Text": {"title": 0},
    "Process Diagram": {"title": 0},
    "Multi-layout 1": {"title": 0, "body": 10},
    "Multi-layout 2": {"title": 0, "body": 10},
    "Image collage": {"title": 0},
    "One Third Two Third Title and Content": {"title": 0, "subtitle": 31, "body": 10, "body2": 32},
    "Two Third One Third Title and Content": {"title": 0, "subtitle": 31, "body": 10, "body2": 32},
    "Two Content Layout Horizontal": {"title": 0, "subtitle": 31, "body": 10, "body2": 32},
}


def inject_content(slide, content: SlideContent, layout_name: str):
    """Inject extracted content INTO the template slide's placeholders.

    This is the core of v4: content goes INTO existing placeholders,
    inheriting their formatting (colours, fonts, sizes, positions).
    """
    ph_map = PLACEHOLDER_MAP.get(layout_name, {})

    # Build a dict of actual placeholder objects on this slide, by index
    placeholders = {}
    for ph in slide.placeholders:
        placeholders[ph.placeholder_format.idx] = ph

    # --- Inject title ---
    if content.title and "title" in ph_map:
        ph_idx = ph_map["title"]
        if ph_idx in placeholders:
            _inject_text_block(placeholders[ph_idx], content.title, preserve_format=True)

    # --- Inject subtitle ---
    if content.subtitle and "subtitle" in ph_map:
        ph_idx = ph_map["subtitle"]
        if ph_idx in placeholders:
            _inject_text_block(placeholders[ph_idx], content.subtitle, preserve_format=True)

    # --- Inject body content ---
    # Skip generic body injection for structural slides handled below
    if content.body_blocks and "body" in ph_map and content.slide_type not in ("cover", "thank_you", "section_divider"):
        ph_idx = ph_map["body"]
        if ph_idx in placeholders:
            _inject_body_blocks(placeholders[ph_idx], content.body_blocks, content.slide_type)

    # --- Inject second body (two-column layouts) ---
    if "body2" in ph_map and len(content.body_blocks) >= 2:
        ph_idx = ph_map["body2"]
        if ph_idx in placeholders:
            # Split body blocks between body and body2
            # The first inject already put all blocks in body — re-do with split
            mid = len(content.body_blocks) // 2
            left_blocks = content.body_blocks[:mid]
            right_blocks = content.body_blocks[mid:]

            # Re-inject body (left column) with only first half
            body_ph_idx = ph_map["body"]
            if body_ph_idx in placeholders:
                _inject_body_blocks(placeholders[body_ph_idx], left_blocks, content.slide_type)
            _inject_body_blocks(placeholders[ph_idx], right_blocks, content.slide_type)

    # --- Inject images ---
    if content.content_images:
        _inject_images(slide, content, ph_map, placeholders)

    # --- Inject tables ---
    if content.tables:
        _inject_tables(slide, content.tables)

    # --- Handle structural slides specially ---
    if content.slide_type == "cover":
        _inject_cover(slide, content, ph_map, placeholders)
    elif content.slide_type == "thank_you":
        _inject_thank_you(slide, content, ph_map, placeholders)
    elif content.slide_type == "section_divider":
        _inject_section_divider(slide, content, ph_map, placeholders)


def _inject_text_block(placeholder, block: TextBlock, preserve_format: bool = True):
    """Inject a TextBlock into a placeholder, preserving template formatting."""
    tf = placeholder.text_frame
    tf.clear()

    for i, para in enumerate(block.paragraphs):
        if not para.text.strip():
            # Keep empty paragraphs as spacing
            if i > 0:
                tf.add_paragraph()
            continue

        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()

        # Set paragraph level for indentation
        if para.level > 0:
            p.level = para.level

        # Don't override alignment — let template handle it
        # (unless source had explicit alignment that matters)

        for j, run in enumerate(para.runs):
            if not run.text:
                continue
            if j == 0 and i == 0 and p.runs:
                r = p.runs[0]
            else:
                r = p.add_run()
            r.text = run.text

            # Preserve bold/italic but NOT font size or colour
            # (let template handle those)
            if run.bold is not None:
                r.font.bold = run.bold
            if run.italic is not None:
                r.font.italic = run.italic

            # Add hyperlink if present
            if run.hyperlink:
                try:
                    r.hyperlink.address = run.hyperlink
                except Exception:
                    pass


def _inject_body_blocks(placeholder, blocks: List[TextBlock], slide_type: str):
    """Inject multiple body text blocks into a single placeholder."""
    tf = placeholder.text_frame
    tf.clear()

    first_para = True
    for block in blocks:
        for para in block.paragraphs:
            if not para.text.strip() and first_para:
                continue  # Skip leading empty paragraphs

            if not first_para:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                first_para = False

            if para.level > 0:
                p.level = para.level

            for j, run in enumerate(para.runs):
                if not run.text:
                    continue
                r = p.add_run()
                r.text = run.text
                if run.bold is not None:
                    r.font.bold = run.bold
                if run.italic is not None:
                    r.font.italic = run.italic
                if run.hyperlink:
                    try:
                        r.hyperlink.address = run.hyperlink
                    except Exception:
                        pass

    if first_para:
        # Nothing was injected — clear the placeholder completely
        tf.clear()


def _inject_images(slide, content: SlideContent, ph_map: dict, placeholders: dict):
    """Inject images into picture placeholders or as positioned shapes."""
    images = content.content_images

    # Try picture placeholder first
    if "picture" in ph_map:
        ph_idx = ph_map["picture"]
        if ph_idx in placeholders and images:
            img = images[0]
            _insert_image_into_placeholder(placeholders[ph_idx], img)
            images = images[1:]  # remaining images

    # Additional images: add as positioned freeform shapes in content area
    if images:
        _add_remaining_images(slide, images, content.target_layout)


def _insert_image_into_placeholder(placeholder, img: ExtractedImage):
    """Insert an image into a PICTURE placeholder."""
    try:
        with tempfile.NamedTemporaryFile(
            suffix=_ext_from_content_type(img.content_type),
            delete=False
        ) as tmp:
            tmp.write(img.image_bytes)
            tmp_path = tmp.name

        try:
            placeholder.insert_picture(tmp_path)
        finally:
            os.unlink(tmp_path)
    except Exception:
        pass  # If picture insertion fails, just skip


def _add_remaining_images(slide, images: List[ExtractedImage], layout_name: str):
    """Add extra images as freeform shapes positioned in the content area."""
    # Position in the lower portion of the content area
    content_top = Emu(2.5 * 914400)  # Below title area
    content_left = Emu(0.52 * 914400)
    max_width = Emu(12 * 914400)
    max_height = Emu(4 * 914400)

    # Distribute images horizontally
    num_images = len(images)
    if num_images == 0:
        return

    img_width = min(max_width // num_images, Emu(5 * 914400))
    img_height = min(max_height, Emu(3.5 * 914400))

    for i, img in enumerate(images[:4]):  # Max 4 extra images
        left = content_left + i * (img_width + Emu(0.2 * 914400))
        top = content_top + Emu(2 * 914400)  # Below text

        try:
            with tempfile.NamedTemporaryFile(
                suffix=_ext_from_content_type(img.content_type),
                delete=False
            ) as tmp:
                tmp.write(img.image_bytes)
                tmp_path = tmp.name

            try:
                slide.shapes.add_picture(tmp_path, left, top, img_width, img_height)
            finally:
                os.unlink(tmp_path)
        except Exception:
            pass


def _inject_tables(slide, tables: List[ExtractedTable]):
    """Inject tables as native PowerPoint table shapes."""
    for table_data in tables[:1]:  # Max 1 table per slide
        if not table_data.rows:
            continue

        rows = table_data.row_count
        cols = table_data.col_count
        if rows == 0 or cols == 0:
            continue

        # Position below title area
        left = Emu(0.52 * 914400)
        top = Emu(2.5 * 914400)
        width = Emu(12 * 914400)
        height = Emu(4 * 914400)

        try:
            tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            tbl = tbl_shape.table

            for r_idx, row_data in enumerate(table_data.rows):
                for c_idx, cell_text in enumerate(row_data):
                    if c_idx < cols and r_idx < rows:
                        tbl.cell(r_idx, c_idx).text = cell_text
        except Exception:
            pass


def _inject_cover(slide, content: SlideContent, ph_map: dict, placeholders: dict):
    """Special injection for cover slides."""
    # Cover 1 has: title (0), subtitle (10), supertitle (11)
    # supertitle (idx 11) sits ABOVE the title — programme name or category
    # subtitle (idx 10) sits BELOW the title — presenter name, date, etc.

    if not content.body_blocks:
        return

    # If we have a subtitle from the source, put it in subtitle placeholder
    if content.subtitle and "subtitle" in ph_map:
        ph_idx = ph_map["subtitle"]
        if ph_idx in placeholders:
            _inject_text_block(placeholders[ph_idx], content.subtitle)

    # Body blocks on a cover are usually presenter info
    # Put them in the subtitle area (below title)
    if content.body_blocks and "subtitle" in ph_map:
        ph_idx = ph_map["subtitle"]
        if ph_idx in placeholders:
            # Combine body blocks into subtitle
            combined = TextBlock(paragraphs=[])
            for block in content.body_blocks:
                combined.paragraphs.extend(block.paragraphs)
            _inject_text_block(placeholders[ph_idx], combined)


def _inject_thank_you(slide, content: SlideContent, ph_map: dict, placeholders: dict):
    """Special injection for Thank You / Contact slides."""
    # Thank You has: title (0), name (10), role (16), email (17), phone (18)
    # Title already injected by main inject_content.
    #
    # For contact slides, body_blocks contain text boxes with:
    # names, roles, emails, phone numbers. Try to parse and place them.

    if not content.body_blocks:
        return

    # Collect all body text lines
    all_lines = []
    for block in content.body_blocks:
        text = block.full_text.strip()
        if text:
            all_lines.append(text)

    if not all_lines:
        return

    # Try to identify structured contact info
    names = []
    roles = []
    emails = []
    phones = []
    other = []

    for line in all_lines:
        line_stripped = line.strip()
        if not line_stripped:
            continue
        if '@' in line_stripped:
            emails.append(line_stripped)
        elif re.search(r'\+?\d[\d\s\-]{7,}', line_stripped):
            phones.append(line_stripped)
        elif any(title in line_stripped.lower() for title in
                 ['professor', 'associate', 'director', 'manager', 'lead',
                  'coordinator', 'lecturer', 'senior', 'executive', 'industry',
                  'dr ', 'dr.', 'head of']):
            roles.append(line_stripped)
        elif len(line_stripped) < 50 and not any(c.isdigit() for c in line_stripped):
            # Short text without numbers — likely a name
            names.append(line_stripped)
        else:
            other.append(line_stripped)

    # Inject into Thank You placeholders
    # name placeholder (idx 10)
    if names and "name" in ph_map and ph_map["name"] in placeholders:
        placeholders[ph_map["name"]].text_frame.text = "\n".join(names)

    # role placeholder (idx 16)
    if roles and "role" in ph_map and ph_map["role"] in placeholders:
        placeholders[ph_map["role"]].text_frame.text = "\n".join(roles)

    # email placeholder (idx 17)
    if emails and "email" in ph_map and ph_map["email"] in placeholders:
        placeholders[ph_map["email"]].text_frame.text = "\n".join(emails)

    # phone placeholder (idx 18)
    if phones and "phone" in ph_map and ph_map["phone"] in placeholders:
        placeholders[ph_map["phone"]].text_frame.text = "\n".join(phones)

    # If we couldn't parse anything structured, just dump all text into name placeholder
    if not names and not roles and not emails and not phones and other:
        if "name" in ph_map and ph_map["name"] in placeholders:
            placeholders[ph_map["name"]].text_frame.text = "\n".join(other)


def _inject_section_divider(slide, content: SlideContent, ph_map: dict, placeholders: dict):
    """Special injection for section divider slides."""
    # Section Divider has: title (0), section_num (11), body (13)
    # Try to extract a section number from the title
    if content.title:
        title_text = content.title.full_text.strip()
        # Check if title is JUST a number (e.g. "00", "01")
        if title_text.replace(" ", "").isdigit():
            num = title_text.strip()
            # Put number in section_num placeholder only
            if "section_num" in ph_map and ph_map["section_num"] in placeholders:
                placeholders[ph_map["section_num"]].text_frame.text = num
            # Clear the title (body text becomes the real title content)
            if "title" in ph_map and ph_map["title"] in placeholders:
                # Use first body block as the title if available
                if content.body_blocks:
                    _inject_text_block(placeholders[ph_map["title"]], content.body_blocks[0])
                    content.body_blocks = content.body_blocks[1:]  # Remove used block
                else:
                    placeholders[ph_map["title"]].text_frame.clear()
        else:
            # Title starts with a number but has more text
            match = re.match(r'^(\d+)\s*[.:\-]?\s*(.*)', title_text)
            if match:
                num = match.group(1)
                rest = match.group(2)
                if "section_num" in ph_map and ph_map["section_num"] in placeholders:
                    placeholders[ph_map["section_num"]].text_frame.text = num
                if rest and "title" in ph_map and ph_map["title"] in placeholders:
                    placeholders[ph_map["title"]].text_frame.text = rest

    # Body text in body placeholder (remaining blocks after title extraction)
    if content.body_blocks and "body" in ph_map:
        ph_idx = ph_map["body"]
        if ph_idx in placeholders:
            _inject_body_blocks(placeholders[ph_idx], content.body_blocks, "section_divider")


# ============================================================================
# STEP 5: Cleanup — clear unused placeholders
# ============================================================================

# Placeholder types that should always be cleared if unused
CLEARABLE_TYPES = {'TITLE', 'CENTER_TITLE', 'SUBTITLE', 'BODY', 'OBJECT'}


def clear_unused_placeholders(slide, content: SlideContent, layout_name: str):
    """Clear any placeholder that wasn't filled with content.

    This prevents ghost text like 'Click to add title' from showing.
    """
    ph_map = PLACEHOLDER_MAP.get(layout_name, {})
    used_indices = set()

    # Track which placeholders were actually used
    if content.title and "title" in ph_map:
        used_indices.add(ph_map["title"])
    if content.subtitle and "subtitle" in ph_map:
        used_indices.add(ph_map["subtitle"])
    if content.body_blocks and "body" in ph_map:
        used_indices.add(ph_map["body"])
    if len(content.body_blocks) >= 2 and "body2" in ph_map:
        used_indices.add(ph_map["body2"])
    if content.content_images and "picture" in ph_map:
        used_indices.add(ph_map["picture"])

    # Special structural slides: mark their specific placeholders used
    # ONLY if content was actually injected into them
    if content.slide_type == "section_divider":
        if content.title:
            for role in ["section_num", "title"]:
                if role in ph_map:
                    used_indices.add(ph_map[role])
        if content.body_blocks:
            if "body" in ph_map:
                used_indices.add(ph_map["body"])

    if content.slide_type == "cover":
        # subtitle is filled by _inject_cover if body_blocks exist
        if content.body_blocks or content.subtitle:
            if "subtitle" in ph_map:
                used_indices.add(ph_map["subtitle"])
        # supertitle is NOT filled by default — leave it to be cleared
        # (unless we explicitly filled it)

    if content.slide_type == "thank_you":
        # These are filled by _inject_thank_you — mark all as used
        # since the handler decides what to put where
        all_body_text = content.primary_body_text
        if all_body_text.strip():
            for role in ["name", "role", "email", "phone"]:
                if role in ph_map:
                    used_indices.add(ph_map[role])

    # Clear unused placeholders
    for ph in slide.placeholders:
        try:
            idx = ph.placeholder_format.idx
            ptype = str(ph.placeholder_format.type) if ph.placeholder_format.type else ""

            # Skip footer and slide number — those are structural
            if any(skip in ptype for skip in ['FOOTER', 'SLIDE_NUMBER', 'DATE_TIME']):
                continue

            # If this placeholder wasn't used, clear it
            if idx not in used_indices:
                if ph.has_text_frame:
                    ph.text_frame.clear()
                    # Also try to make the placeholder invisible
                    # by removing its default text
                    _clear_placeholder_default_text(ph)
        except (ValueError, AttributeError):
            pass


def _clear_placeholder_default_text(placeholder):
    """Remove default/ghost text from an unused placeholder."""
    try:
        sp = placeholder._element
        # Clear all <a:p> elements to remove "Click to add..." text
        a_ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        txBody = sp.find(f'{a_ns}txBody')
        if txBody is not None:
            for p in txBody.findall(f'{a_ns}p'):
                # Remove all runs
                for r in p.findall(f'{a_ns}r'):
                    p.remove(r)
                # Remove endParaRPr default text
                for endPara in p.findall(f'{a_ns}endParaRPr'):
                    pass  # Keep endParaRPr but ensure no text
    except Exception:
        pass


# ============================================================================
# STEP 6: Design Analysis (hook for AI review)
# ============================================================================

def analyse_design(content: SlideContent) -> List[str]:
    """Analyse extracted content for design issues.

    Returns a list of flag strings. These flags can be used to:
    - Warn the user about overset text
    - Suggest slide splits
    - Flag missing titles
    - Note excessive bullet points

    This is the LOCAL heuristic analysis. A separate AI vision pass
    can be added on top for deeper analysis.
    """
    flags = []

    # Overset text detection
    body_chars = sum(b.char_count for b in content.body_blocks)
    if body_chars > 800:
        flags.append(f"OVERSET: {body_chars} chars of body text — likely overflows placeholder. Consider splitting into 2+ slides.")
    elif body_chars > 500:
        flags.append(f"DENSE: {body_chars} chars of body text — may be tight. Review for overflow.")

    # Too many bullet points
    bullet_count = sum(
        len([p for p in b.paragraphs if p.level > 0])
        for b in content.body_blocks
    )
    if bullet_count > 10:
        flags.append(f"BULLETS: {bullet_count} bullet points — consider consolidating or splitting.")

    # Missing title
    if not content.title and content.slide_type not in ["blank"]:
        flags.append("NO_TITLE: Slide has no identifiable title.")

    # Image + heavy text conflict
    if content.has_images and body_chars > 400:
        flags.append("IMG_TEXT_CONFLICT: Slide has both images and heavy text — likely needs a text-only or image-only layout.")

    # Table on a slide with other content
    if content.tables and body_chars > 100:
        flags.append("TABLE_CROWDED: Table slide also has significant body text.")

    content.design_flags = flags
    return flags


# ============================================================================
# Main Pipeline
# ============================================================================

@dataclass
class SlideResult:
    slide_number: int
    source_layout: str
    slide_type: str
    target_layout: str
    text_chars: int = 0
    image_count: int = 0
    table_count: int = 0
    design_flags: List[str] = field(default_factory=list)
    status: str = "success"
    error: str = ""


def run_v4_pipeline(
    source_pptx_bytes: bytes,
    template_path: str,
    progress_callback=None,
    skip_design_analysis: bool = False,
) -> dict:
    """Run the v4 Content Extraction + Placeholder Injection pipeline.

    Steps:
      1. Extract content from each source slide
      2. Classify each slide's type
      3. Select the best template layout
      4. Create new slide from template layout
      5. Inject extracted content into template placeholders
      6. Clear unused placeholders
      7. Run design analysis (optional)

    Returns dict with:
      output_pptx_bytes, results, summary, design_report
    """
    results = []
    design_report = []

    def progress(step, detail, pct):
        if progress_callback:
            progress_callback(step, detail, pct)

    # Load source
    progress("load", "Loading source presentation...", 0.05)
    source_prs = Presentation(io.BytesIO(source_pptx_bytes))
    num_slides = len(source_prs.slides)

    # Load template
    progress("create", "Creating target from template...", 0.10)
    target_prs = Presentation(template_path)
    _remove_all_slides(target_prs)

    # Build layout lookup
    available_layouts = {}
    for layout in target_prs.slide_layouts:
        available_layouts[layout.name] = layout

    # --- Process each slide ---
    extracted_contents = []

    for slide_idx in range(num_slides):
        pct = 0.12 + (0.35 * slide_idx / num_slides)
        progress("extract", f"Extracting slide {slide_idx + 1}/{num_slides}...", pct)

        source_slide = source_prs.slides[slide_idx]
        content = extract_slide_content(source_slide, slide_idx + 1)

        # Classify
        classify_slide(content, slide_idx, num_slides)

        # Select layout
        select_layout(content, available_layouts)

        # Design analysis
        if not skip_design_analysis:
            flags = analyse_design(content)
            if flags:
                design_report.append({
                    "slide": slide_idx + 1,
                    "type": content.slide_type,
                    "flags": flags,
                })

        extracted_contents.append(content)

    # --- Build output slides ---
    for slide_idx, content in enumerate(extracted_contents):
        pct = 0.50 + (0.45 * slide_idx / num_slides)
        progress("build", f"Building slide {slide_idx + 1}/{num_slides}...", pct)

        layout_name = content.target_layout
        template_layout = available_layouts.get(layout_name)

        if not template_layout:
            template_layout = available_layouts.get("Title and Content")
            layout_name = "Title and Content"

        result = SlideResult(
            slide_number=content.slide_number,
            source_layout=content.source_layout_name,
            slide_type=content.slide_type,
            target_layout=layout_name,
            text_chars=content.total_text_chars,
            image_count=len(content.content_images),
            table_count=len(content.tables),
            design_flags=content.design_flags,
        )

        try:
            # Create slide from template layout
            new_slide = target_prs.slides.add_slide(template_layout)

            # Inject content
            inject_content(new_slide, content, layout_name)

            # Clear unused placeholders
            clear_unused_placeholders(new_slide, content, layout_name)

        except Exception as e:
            result.status = "failed"
            result.error = str(e)

        results.append(result)

    # --- Save ---
    progress("save", "Saving presentation...", 0.96)
    output_buf = io.BytesIO()
    target_prs.save(output_buf)
    clean_bytes = _deduplicate_zip(output_buf.getvalue())

    # --- Summary ---
    type_counts = Counter(r.slide_type for r in results)
    layout_counts = Counter(r.target_layout for r in results)
    flagged_slides = sum(1 for r in results if r.design_flags)

    summary = {
        "total_slides": num_slides,
        "success": sum(1 for r in results if r.status == "success"),
        "failed": sum(1 for r in results if r.status == "failed"),
        "slide_types": dict(type_counts),
        "layout_distribution": dict(layout_counts),
        "flagged_for_review": flagged_slides,
        "total_images": sum(r.image_count for r in results),
        "total_tables": sum(r.table_count for r in results),
    }

    progress("done", "Complete", 1.0)

    return {
        "output_pptx_bytes": clean_bytes,
        "results": results,
        "summary": summary,
        "design_report": design_report,
    }


# ============================================================================
# Utility Functions
# ============================================================================

def _remove_all_slides(prs):
    """Remove all slides from a presentation (keep layouts)."""
    while len(prs.slides) > 0:
        try:
            rId = prs.slides._sldIdLst[0].get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            if rId:
                prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        except Exception:
            break


def _deduplicate_zip(pptx_bytes: bytes) -> bytes:
    """Remove duplicate ZIP entries that python-pptx creates."""
    src = io.BytesIO(pptx_bytes)
    dst = io.BytesIO()
    with zipfile.ZipFile(src, "r") as zin:
        entries = {}
        for info in zin.infolist():
            entries[info.filename] = (info, zin.read(info.filename))
        with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
            for filename, (info, data) in entries.items():
                zout.writestr(info, data)
    return dst.getvalue()


def _ext_from_content_type(content_type: str) -> str:
    """Get file extension from MIME content type."""
    mapping = {
        "image/png": ".png",
        "image/jpeg": ".jpg",
        "image/gif": ".gif",
        "image/bmp": ".bmp",
        "image/tiff": ".tiff",
        "image/svg+xml": ".svg",
        "image/x-wmf": ".wmf",
        "image/x-emf": ".emf",
    }
    return mapping.get(content_type, ".png")


# ============================================================================
# CLI
# ============================================================================

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 3:
        print("Usage: python v4_engine.py <source.pptx> <template.pptx> [output.pptx]")
        sys.exit(1)

    source_path = sys.argv[1]
    template_path = sys.argv[2]
    output_path = sys.argv[3] if len(sys.argv) > 3 else "v4_output.pptx"

    with open(source_path, "rb") as f:
        source_bytes = f.read()

    def print_progress(step, detail, pct):
        print(f"  [{pct:5.1%}] {detail}")

    result = run_v4_pipeline(source_bytes, template_path, print_progress)

    with open(output_path, "wb") as f:
        f.write(result["output_pptx_bytes"])

    s = result["summary"]
    print(f"\nDone! {s['success']}/{s['total_slides']} slides rebuilt.")
    print(f"Failed: {s['failed']}")
    print(f"Images: {s['total_images']}, Tables: {s['total_tables']}")
    print(f"Flagged for design review: {s['flagged_for_review']}")

    if result["design_report"]:
        print("\n--- Design Report ---")
        for item in result["design_report"]:
            print(f"  Slide {item['slide']} ({item['type']}):")
            for flag in item["flags"]:
                print(f"    ⚠ {flag}")

    print("\nSlide type distribution:")
    for stype, count in sorted(s["slide_types"].items(), key=lambda x: -x[1]):
        print(f"  {count:3d}x {stype}")

    print("\nLayout distribution:")
    for name, count in sorted(s["layout_distribution"].items(), key=lambda x: -x[1]):
        print(f"  {count:3d}x {name}")
