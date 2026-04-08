#!/usr/bin/env python3
"""
UQ Slide Compliance Tool — Reference & Attribution Checker
============================================================
Scans a .pptx for:
  1. In-text citations → validates APA 7 format
  2. Reference list entries → validates APA 7 format, auto-fixes where possible
  3. Image attributions → standardises to "Source: Adobe Stock {ID}" etc.
  4. Cross-references citations ↔ reference list (flags orphans)
  5. Flags slides with images but no attribution text

Can auto-fix the PPTX file and produce a change report.

Usage:
    python ref_checker.py input.pptx --output fixed.pptx --report
    python ref_checker.py input.pptx --report-only   # don't modify, just report
"""

import argparse
import json
import re
import sys
from collections import defaultdict
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN


# ─── Constants ────────────────────────────────────────────────────────

# Known stock/image sources that UQ holds licences for
LICENSED_SOURCES = {
    "adobe stock": "Adobe Stock",
    "shutterstock": "Shutterstock",
    "microsoft stock": "Microsoft Stock",
}

# Standard attribution format: "Source: {Provider} {ID}"
STANDARD_ATTR_PREFIX = "Source:"

# APA 7 in-text citation patterns
# (Author, Year) | (Author & Author, Year) | (Author et al., Year)
# Also handles multiple citations: (Author, Year; Author, Year)
CITE_PATTERN = re.compile(
    r'\('
    r'([A-Z][A-Za-z\'\-]+(?:\s(?:&|and)\s[A-Z][A-Za-z\'\-]+)?'
    r'(?:\s+et\s+al\.)?'
    r',?\s*\d{4}[a-z]?'
    r'(?:\s*;\s*'  # semicolon separator for multiple citations
    r'[A-Z][A-Za-z\'\-]+(?:\s(?:&|and)\s[A-Z][A-Za-z\'\-]+)?'
    r'(?:\s+et\s+al\.)?'
    r',?\s*\d{4}[a-z]?)*'
    r')\)'
)

# Simpler pattern to extract individual author-year pairs from a citation
AUTHOR_YEAR_PATTERN = re.compile(
    r'([A-Z][A-Za-z\'\-]+(?:\s(?:&|and)\s[A-Z][A-Za-z\'\-]+)?'
    r'(?:\s+et\s+al\.)?)'
    r',?\s*(\d{4})[a-z]?'
)

# Narrative citation: Author (Year) or Author and Author (Year)
NARRATIVE_CITE_PATTERN = re.compile(
    r'([A-Z][A-Za-z\'\-]+(?:\s(?:&|and|&)\s[A-Z][A-Za-z\'\-]+)?'
    r'(?:\s+et\s+al\.)?)'
    r'\s*\((\d{4})[a-z]?\)'
)

# Reference list entry — broad pattern to catch the start of entries
# Author, A. B. (Year). or Author, A. B., & Author, C. D. (Year).
REF_ENTRY_PATTERN = re.compile(
    r'([A-Z][A-Za-z\'\-]+,\s[A-Z]\.(?:\s[A-Z]\.)?'
    r'(?:,?\s(?:&|and)\s[A-Z][A-Za-z\'\-]+,\s[A-Z]\.(?:\s[A-Z]\.)?)*'
    r'(?:,?\s(?:et\sal\.))?)'
    r'\s*\((\d{4}(?:,\s\w+\s\d{1,2})?|n\.d\.)\)\.\s*'
)

# "Adapted from" pattern — common in academic slides
ADAPTED_PATTERN = re.compile(
    r'(?:Adapted\s+(?:from|by))\s*:?\s*'
    r'([A-Z][A-Za-z\'\-]+(?:,\s[A-Z]\.(?:\s[A-Z]\.)?)?'
    r'(?:(?:,?\s(?:&|and)\s)[A-Z][A-Za-z\'\-]+(?:,\s[A-Z]\.(?:\s[A-Z]\.)?)?)*)'
    r'\s*\((\d{4}(?:,\s\w+\s\d{1,2})?|n\.d\.)\)',
    re.IGNORECASE,
)

# Image attribution patterns
ATTR_PATTERNS = [
    # "Source: Adobe Stock 123456"
    re.compile(r'Source:\s*Adobe\s+Stock\s+(\d+)', re.IGNORECASE),
    # "Image licensed through Adobe Stock: 123456"
    re.compile(r'Image\s+licensed\s+through\s+Adobe\s+Stock:\s*(\d+)', re.IGNORECASE),
    # "Image Source: Adobe Stock 123456"
    re.compile(r'Image\s+Source:\s*Adobe\s+Stock\s+(\d+)', re.IGNORECASE),
    # "Source: Adobe Stock" (no ID)
    re.compile(r'Source:\s*Adobe\s+Stock\b(?!\s*\d)', re.IGNORECASE),
    # Shutterstock variants
    re.compile(r'(?:Source:|Image\s+licensed\s+through)\s*Shutterstock:\s*(\d+)', re.IGNORECASE),
    # Microsoft Stock variants
    re.compile(r'(?:Source:|Image[s]?\s+(?:from|by|licensed\s+through))?\s*:?\s*Microsoft\s+Stock(?:\s+Image)?', re.IGNORECASE),
    # Wikimedia Commons
    re.compile(r'Source:\s*(.+?(?:Wikimedia\s+Commons|CC\s+BY[^\s]*).+)', re.IGNORECASE),
    # Public domain
    re.compile(r'Source:\s*Public\s+domain', re.IGNORECASE),
    # Flickr
    re.compile(r'(?:Image\s+from|Source:)\s*(.+?Flickr.+)', re.IGNORECASE),
    # Generic "Source: ..."
    re.compile(r'Source:\s*(.+)', re.IGNORECASE),
]

# Adobe Stock with ID — for standardisation
ADOBE_STOCK_WITH_ID = re.compile(
    r'(?:Source:\s*|Image\s+(?:licensed\s+through|Source:\s*|from\s+))?\s*Adobe\s+Stock[:\s]*(\d+)',
    re.IGNORECASE,
)

# Adobe Stock without ID
ADOBE_STOCK_NO_ID = re.compile(
    r'(?:Source:\s*|Image\s+(?:licensed\s+through|from\s+))?\s*Adobe\s+Stock\b(?!\s*\d)',
    re.IGNORECASE,
)

# Shutterstock with ID
SHUTTERSTOCK_WITH_ID = re.compile(
    r'(?:Source:\s*|Image\s+licensed\s+through\s*)Shutterstock[:\s]*(\d+)',
    re.IGNORECASE,
)

# Microsoft Stock
MICROSOFT_STOCK = re.compile(
    r'(?:Source:\s*|Image[s]?\s+(?:from|by)\s+|Image\s+licensed\s+through\s*)?Microsoft\s+Stock(?:\s+Image)?',
    re.IGNORECASE,
)


# ─── Helper Functions ─────────────────────────────────────────────────

def extract_surname(author_str):
    """Extract the primary surname from an author string for matching."""
    # "Kotter, J. P." → "Kotter"
    # "Kotter" → "Kotter"
    parts = author_str.split(",")
    return parts[0].strip()


def normalise_author_for_matching(author_str):
    """Normalise author string for fuzzy matching between citations and refs.
    Returns a list of surnames to match against (handles multi-author citations)."""
    s = author_str.strip()
    s = re.sub(r'\s+et\s+al\.?', '', s)
    s = re.sub(r'\s+', ' ', s)
    # Get first surname only (the primary lookup key)
    return extract_surname(s).lower()


def find_images_on_slide(slide):
    """Return a list of image shapes on a slide."""
    images = []
    for shape in slide.shapes:
        try:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                images.append(shape)
            elif hasattr(shape, 'image'):
                images.append(shape)
        except Exception:
            pass
    return images


def find_attribution_text_on_slide(slide):
    """Find all text on a slide that looks like an image attribution."""
    attributions = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue

        # Check each line
        for line in text.split('\n'):
            line = line.strip()
            if not line:
                continue
            line_lower = line.lower()

            # Check for known attribution patterns
            is_attr = False
            attr_type = None
            stock_id = None

            # Adobe Stock with ID
            m = ADOBE_STOCK_WITH_ID.search(line)
            if m:
                is_attr = True
                attr_type = "adobe_stock"
                stock_id = m.group(1)
            # Shutterstock with ID
            elif SHUTTERSTOCK_WITH_ID.search(line):
                m = SHUTTERSTOCK_WITH_ID.search(line)
                is_attr = True
                attr_type = "shutterstock"
                stock_id = m.group(1)
            # Microsoft Stock
            elif MICROSOFT_STOCK.search(line):
                is_attr = True
                attr_type = "microsoft_stock"
            # Wikimedia Commons
            elif 'wikimedia commons' in line_lower or 'cc by' in line_lower:
                is_attr = True
                attr_type = "wikimedia"
            # Public domain
            elif 'public domain' in line_lower:
                is_attr = True
                attr_type = "public_domain"
            # Flickr
            elif 'flickr' in line_lower:
                is_attr = True
                attr_type = "flickr"
            # Generic "Source: ..."
            elif line_lower.startswith('source:'):
                is_attr = True
                attr_type = "generic"
            # "Image from/by/licensed..."
            elif re.match(r'image\s+(?:from|by|licensed)', line_lower):
                is_attr = True
                attr_type = "generic"

            if is_attr:
                attributions.append({
                    "text": line,
                    "type": attr_type,
                    "stock_id": stock_id,
                    "shape_name": shape.name,
                    "shape": shape,
                })

    return attributions


def standardise_attribution(text, attr_type, stock_id):
    """Return the standardised form of an attribution string."""
    if attr_type == "adobe_stock" and stock_id:
        return f"Source: Adobe Stock {stock_id}"
    elif attr_type == "shutterstock" and stock_id:
        return f"Source: Shutterstock {stock_id}"
    elif attr_type == "microsoft_stock":
        return "Source: Microsoft Stock"
    # For wikimedia, public domain, flickr, generic — leave as-is but ensure "Source: " prefix
    elif not text.lower().startswith("source:"):
        # Strip common prefixes before adding "Source: "
        cleaned = re.sub(
            r'^(?:Image\s+(?:from|by|licensed\s+through)\s*:?\s*)',
            '', text, flags=re.IGNORECASE
        ).strip()
        return f"Source: {cleaned}" if cleaned else text
    return text


# ─── APA 7 Validation Helpers ─────────────────────────────────────────

class APA7Issue:
    """Represents a single APA 7 formatting issue."""
    def __init__(self, slide_num, category, description, original=None,
                 suggested=None, severity="warning", shape_name=None):
        self.slide_num = slide_num
        self.category = category  # citation, reference, attribution, cross_ref, missing_attr
        self.description = description
        self.original = original
        self.suggested = suggested
        self.severity = severity  # error, warning, info
        self.shape_name = shape_name

    def to_dict(self):
        d = {
            "slide": self.slide_num,
            "category": self.category,
            "description": self.description,
            "severity": self.severity,
        }
        if self.original:
            d["original"] = self.original
        if self.suggested:
            d["suggested"] = self.suggested
        if self.shape_name:
            d["shape"] = self.shape_name
        return d


def check_citation_format(citation_text, slide_num):
    """Check an in-text citation for APA 7 compliance. Returns list of issues."""
    issues = []
    inner = citation_text[1:-1]  # strip outer parens

    # Check for "and" vs "&" — APA 7 uses "&" in parenthetical citations
    if " and " in inner and "(" == citation_text[0]:
        issues.append(APA7Issue(
            slide_num, "citation",
            f'Use "&" instead of "and" in parenthetical citations',
            original=citation_text,
            suggested=citation_text.replace(" and ", " & "),
            severity="warning",
        ))

    # Check for missing comma before year
    # e.g., "(Kotter 1996)" should be "(Kotter, 1996)"
    if re.search(r'[a-z]\s+\d{4}', inner) and ',' not in inner:
        issues.append(APA7Issue(
            slide_num, "citation",
            f'Missing comma between author and year',
            original=citation_text,
            severity="warning",
        ))

    # Check for "et al" without period
    if 'et al' in inner and 'et al.' not in inner:
        issues.append(APA7Issue(
            slide_num, "citation",
            f'"et al" should have a period: "et al."',
            original=citation_text,
            suggested=citation_text.replace('et al', 'et al.'),
            severity="warning",
        ))

    return issues


def check_reference_entry_format(text, slide_num):
    """Check a reference list entry for common APA 7 issues. Returns list of issues."""
    issues = []

    # Check: Year should be in parentheses
    # Look for patterns like "Author, A. B. 2020." (missing parens around year)
    if re.search(r'[A-Z]\.\s+\d{4}\.', text):
        issues.append(APA7Issue(
            slide_num, "reference",
            "Year should be in parentheses: (Year).",
            original=text[:80],
            severity="warning",
        ))

    # Check: DOI should be a URL format (https://doi.org/...)
    doi_match = re.search(r'doi:\s*(10\.\d+/\S+)', text, re.IGNORECASE)
    if doi_match:
        issues.append(APA7Issue(
            slide_num, "reference",
            "DOI should be in URL format: https://doi.org/...",
            original=f"doi: {doi_match.group(1)}",
            suggested=f"https://doi.org/{doi_match.group(1)}",
            severity="warning",
        ))

    # Check: "Retrieved from" is not used in APA 7 for DOIs/URLs
    if 'retrieved from' in text.lower():
        issues.append(APA7Issue(
            slide_num, "reference",
            '"Retrieved from" is deprecated in APA 7 — use the URL directly',
            original=text[:80],
            severity="warning",
        ))

    # Check: Publisher location is not needed in APA 7
    # Pattern: "City, State: Publisher" or "City: Publisher"
    if re.search(r'(?:New York|London|San Francisco|Boston|Chicago|Thousand Oaks|Palo Alto),?\s*(?:[A-Z]{2})?:\s*[A-Z]', text):
        issues.append(APA7Issue(
            slide_num, "reference",
            "Publisher location not needed in APA 7 — just use publisher name",
            original=text[:80],
            severity="info",
        ))

    # Check: Edition format should be (Xth ed.)
    ed_match = re.search(r'(\d+)\s*(?:st|nd|rd|th)\s+edition', text, re.IGNORECASE)
    if ed_match:
        ordinal = ed_match.group(1)
        suffix = {1: "st", 2: "nd", 3: "rd"}.get(int(ordinal) % 10, "th")
        if int(ordinal) in [11, 12, 13]:
            suffix = "th"
        issues.append(APA7Issue(
            slide_num, "reference",
            f'Edition format should be "({ordinal}{suffix} ed.)" not "{ed_match.group(0)}"',
            original=ed_match.group(0),
            suggested=f"({ordinal}{suffix} ed.)",
            severity="warning",
        ))

    return issues


# ─── Main Checker Class ───────────────────────────────────────────────

class RefChecker:
    """Scans a presentation for reference and attribution issues."""

    def __init__(self, prs, report=True):
        self.prs = prs
        self.do_report = report
        self.issues = []
        self.changes = []
        self.citations = []        # (slide_num, author, year, full_text)
        self.references = []       # (slide_num, author_surname, year, full_text)
        self.ref_slides = []       # slide numbers that contain reference lists
        self.attribution_fixes = 0

    def log_issue(self, issue):
        """Record an issue found during checking."""
        self.issues.append(issue)

    def log_change(self, slide_num, category, detail):
        """Record a change made to the file."""
        self.changes.append({
            "slide": slide_num,
            "category": category,
            "detail": detail,
        })

    # ── Scanning ──────────────────────────────────────────────────────

    def scan_citations(self):
        """Find all in-text citations across the presentation."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text

                # Parenthetical citations: (Author, Year)
                for m in CITE_PATTERN.finditer(text):
                    full = m.group(0)
                    # Extract individual author-year pairs
                    for am in AUTHOR_YEAR_PATTERN.finditer(m.group(1)):
                        author = am.group(1).strip()
                        year = am.group(2)
                        self.citations.append((slide_idx, author, year, full))

                    # Check format
                    issues = check_citation_format(full, slide_idx)
                    for issue in issues:
                        issue.shape_name = shape.name
                        self.log_issue(issue)

                # Narrative citations: Author (Year)
                for m in NARRATIVE_CITE_PATTERN.finditer(text):
                    author = m.group(1).strip()
                    year = m.group(2)
                    full = m.group(0)
                    # Only count if it looks like a real citation (not a random pattern)
                    if len(author) > 2 and author[0].isupper():
                        self.citations.append((slide_idx, author, year, full))

                # "Adapted from Author (Year)" patterns
                for m in ADAPTED_PATTERN.finditer(text):
                    author = m.group(1).strip()
                    year = m.group(2)
                    full = m.group(0)
                    self.citations.append((slide_idx, author, year, full))

    def scan_references(self):
        """Find reference list slides and extract entries."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            full_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    full_text += shape.text_frame.text + "\n"

            # Detect reference list slides
            first_line = full_text.strip().split('\n')[0].strip().lower() if full_text.strip() else ""
            is_ref_slide = first_line in ['references', 'bibliography',
                                           'reference list', 'further reading',
                                           'recommended reading', 'works cited']

            if not is_ref_slide:
                # Also check if slide has title "References"
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip().lower() in ['references', 'bibliography',
                                                               'reference list', 'works cited']:
                                # Check if it's likely a title (larger font or first shape)
                                is_ref_slide = True
                                break
                    if is_ref_slide:
                        break

            if is_ref_slide:
                self.ref_slides.append(slide_idx)

                # Extract individual reference entries from the body text
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    text = shape.text_frame.text

                    # Skip the title shape
                    if text.strip().lower() in ['references', 'bibliography',
                                                  'reference list', 'works cited']:
                        continue

                    # Find reference entries
                    for m in REF_ENTRY_PATTERN.finditer(text):
                        authors = m.group(1)
                        year = m.group(2)
                        # Get the full entry (up to next entry or end)
                        start = m.start()
                        # Find next entry
                        next_m = REF_ENTRY_PATTERN.search(text, m.end())
                        end = next_m.start() if next_m else min(start + 500, len(text))
                        full_entry = text[start:end].strip()

                        surname = extract_surname(authors)
                        self.references.append((slide_idx, surname, year, full_entry))

                        # Check format of this entry
                        issues = check_reference_entry_format(full_entry, slide_idx)
                        for issue in issues:
                            issue.shape_name = shape.name
                            self.log_issue(issue)

    def scan_attributions(self):
        """Find and check image attributions on all slides."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            images = find_images_on_slide(slide)
            attrs = find_attribution_text_on_slide(slide)

            # Flag slides with images but no attribution
            if images and not attrs:
                # Filter out likely decorative/institutional images
                # by checking alt text and shape names for common patterns
                real_images = []
                for img in images:
                    name = (img.name or "").lower()
                    alt = ""
                    try:
                        for el in img._element.iter():
                            if el.tag.endswith('cNvPr'):
                                alt = (el.get('descr', '') or '').lower()
                                break
                    except Exception:
                        pass

                    # Skip likely decorative/institutional images
                    skip = False
                    skip_terms = ['logo', 'icon', 'graphic', 'decorat',
                                  'bullseye', 'meeting', 'newspaper',
                                  'uq', 'background', 'banner']
                    for term in skip_terms:
                        if term in name or term in alt:
                            skip = True
                            break

                    # Skip very small images (likely icons)
                    try:
                        if img.width and img.height:
                            w_inches = img.width / 914400
                            h_inches = img.height / 914400
                            if w_inches < 1.0 and h_inches < 1.0:
                                skip = True
                    except Exception:
                        pass

                    if not skip:
                        real_images.append(img)

                if real_images:
                    self.log_issue(APA7Issue(
                        slide_idx, "missing_attr",
                        f"Slide has {len(real_images)} image(s) but no attribution text",
                        severity="warning",
                    ))

            # Check attribution format
            for attr in attrs:
                original = attr["text"]

                # Skip text that looks like it has other content concatenated
                # (e.g., a reference + attribution jammed together)
                if len(original) > 120 and attr["type"] not in ("adobe_stock", "shutterstock"):
                    self.log_issue(APA7Issue(
                        slide_idx, "attribution",
                        f"Attribution may be concatenated with other text — review manually",
                        original=original[:100] + "...",
                        severity="warning",
                        shape_name=attr["shape_name"],
                    ))
                    continue

                standardised = standardise_attribution(
                    original, attr["type"], attr["stock_id"]
                )

                if original != standardised:
                    self.log_issue(APA7Issue(
                        slide_idx, "attribution",
                        f"Non-standard attribution format",
                        original=original,
                        suggested=standardised,
                        severity="info",
                        shape_name=attr["shape_name"],
                    ))

    def cross_reference(self):
        """Cross-reference citations against the reference list."""
        if not self.ref_slides:
            if self.citations:
                self.log_issue(APA7Issue(
                    0, "cross_ref",
                    f"Found {len(self.citations)} citations but no reference list slides",
                    severity="error",
                ))
            return

        # Build lookup sets
        ref_lookup = {}  # surname_lower → [(year, full_entry), ...]
        for _, surname, year, full in self.references:
            key = surname.lower()
            if key not in ref_lookup:
                ref_lookup[key] = []
            ref_lookup[key].append((year, full))

        cite_lookup = {}  # surname_lower → [(year, slide, full), ...]
        seen_citations = set()  # deduplicate
        for slide_num, author, year, full in self.citations:
            surname = normalise_author_for_matching(author)
            dedup_key = (surname, year)
            if dedup_key in seen_citations:
                continue
            seen_citations.add(dedup_key)
            if surname not in cite_lookup:
                cite_lookup[surname] = []
            cite_lookup[surname].append((year, slide_num, full))

        # Find orphaned citations (cited but not in references)
        for surname, entries in cite_lookup.items():
            for year, slide_num, full_cite in entries:
                found = False
                # Try exact match first
                if surname in ref_lookup:
                    for ref_year, ref_full in ref_lookup[surname]:
                        if ref_year == year or ref_year == "n.d.":
                            found = True
                            break
                # Try partial match (surname appears anywhere in ref entry text)
                if not found:
                    for ref_key, ref_entries in ref_lookup.items():
                        for ref_year, ref_full in ref_entries:
                            if (surname in ref_full.lower()
                                    and (ref_year == year or ref_year == "n.d.")):
                                found = True
                                break
                        if found:
                            break
                if not found:
                    # Skip common false positives (e.g., "Team" from "Team (2002)")
                    if surname in ['team', 'group', 'the', 'this', 'that']:
                        continue
                    self.log_issue(APA7Issue(
                        slide_num, "cross_ref",
                        f'Citation for {surname.title()} ({year}) — no matching reference list entry found',
                        original=full_cite[:80],
                        severity="warning",
                    ))

        # Find orphaned references (in list but never cited)
        for surname, entries in ref_lookup.items():
            for year, full_ref in entries:
                found = False
                # Try exact match
                if surname.lower() in cite_lookup:
                    for cite_year, _, _ in cite_lookup[surname.lower()]:
                        if cite_year == year:
                            found = True
                            break
                # Try partial match (surname appears in any citation author string)
                if not found:
                    for _, author, cite_year, _ in self.citations:
                        if (surname.lower() in author.lower()
                                and cite_year == year):
                            found = True
                            break
                if not found:
                    # This is only informational — many refs may be "further reading"
                    self.log_issue(APA7Issue(
                        self.ref_slides[0], "cross_ref",
                        f'Reference entry for {surname.title()} ({year}) — never cited in slide text',
                        original=full_ref[:100],
                        severity="info",
                    ))

    # ── Auto-fixing ───────────────────────────────────────────────────

    def fix_attributions(self):
        """Standardise image attribution text in the PPTX."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            attrs = find_attribution_text_on_slide(slide)

            for attr in attrs:
                original = attr["text"]
                standardised = standardise_attribution(
                    original, attr["type"], attr["stock_id"]
                )

                if original != standardised:
                    # Find and replace in the shape's text frame
                    shape = attr["shape"]
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            full_para_text = para.text
                            if original in full_para_text:
                                # Replace in runs
                                self._replace_text_in_paragraph(
                                    para, original, standardised
                                )
                                self.log_change(
                                    slide_idx, "attribution",
                                    f'Standardised: "{original}" → "{standardised}"'
                                )
                                self.attribution_fixes += 1
                                break

    def fix_citations(self):
        """Auto-fix citation formatting issues in the PPTX."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for para in shape.text_frame.paragraphs:
                    text = para.text

                    # Fix "and" → "&" in parenthetical citations
                    for m in CITE_PATTERN.finditer(text):
                        citation = m.group(0)
                        if " and " in citation:
                            fixed = citation.replace(" and ", " & ")
                            self._replace_text_in_paragraph(para, citation, fixed)
                            self.log_change(
                                slide_idx, "citation",
                                f'Fixed "and" → "&" in citation: {fixed}'
                            )

                    # Fix "et al" → "et al."
                    if 'et al' in text and 'et al.' not in text:
                        self._replace_text_in_paragraph(para, 'et al', 'et al.')
                        self.log_change(
                            slide_idx, "citation",
                            'Added period: "et al" → "et al."'
                        )

    def _replace_text_in_paragraph(self, paragraph, old_text, new_text):
        """Replace text across runs in a paragraph, preserving formatting."""
        # Build full paragraph text and track run boundaries
        runs = list(paragraph.runs)
        if not runs:
            return

        full_text = "".join(r.text for r in runs)
        if old_text not in full_text:
            return

        # Simple approach: find which runs contain the old text
        # and replace it, preserving the first run's formatting
        new_full = full_text.replace(old_text, new_text, 1)

        # Redistribute text across existing runs
        # Strategy: put all new text in runs proportionally
        pos = 0
        for i, run in enumerate(runs):
            old_len = len(run.text)
            if i == len(runs) - 1:
                # Last run gets the remainder
                run.text = new_full[pos:]
            else:
                # Proportional distribution
                new_len = old_len + (len(new_full) - len(full_text))
                if i == 0:
                    run.text = new_full[pos:pos + old_len + (len(new_full) - len(full_text))]
                    pos += len(run.text)
                else:
                    run.text = new_full[pos:pos + old_len]
                    pos += old_len

    # ── Run All ───────────────────────────────────────────────────────

    def check_all(self):
        """Run all checks (scan only, no modifications)."""
        self.scan_citations()
        self.scan_references()
        self.scan_attributions()
        self.cross_reference()

    def fix_all(self):
        """Run all checks and apply auto-fixes."""
        # Scan first to populate data
        self.scan_citations()
        self.scan_references()
        self.scan_attributions()
        self.cross_reference()

        # Apply fixes
        self.fix_attributions()
        self.fix_citations()

    def generate_report(self):
        """Generate a structured report of all findings."""
        summary = {
            "citations_found": len(self.citations),
            "references_found": len(self.references),
            "ref_slides": self.ref_slides,
            "issues_by_category": defaultdict(int),
            "issues_by_severity": defaultdict(int),
            "attribution_fixes": self.attribution_fixes,
        }

        for issue in self.issues:
            summary["issues_by_category"][issue.category] += 1
            summary["issues_by_severity"][issue.severity] += 1

        return {
            "summary": dict(summary),
            "issues": [i.to_dict() for i in self.issues],
            "changes": self.changes,
            "total_issues": len(self.issues),
            "total_changes": len(self.changes),
            "citations": [
                {"slide": s, "author": a, "year": y, "text": t}
                for s, a, y, t in self.citations
            ],
            "references": [
                {"slide": s, "author": a, "year": y, "text": t[:150]}
                for s, a, y, t in self.references
            ],
        }


# ─── CLI ──────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Check and fix APA 7 references and image attributions in a .pptx"
    )
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("--output", "-o", help="Output .pptx file (with fixes applied)")
    parser.add_argument(
        "--report", "-r", action="store_true",
        help="Print a JSON report of all issues found"
    )
    parser.add_argument(
        "--report-only", action="store_true",
        help="Only report — don't modify the file"
    )
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: {input_path} not found", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(input_path))
    checker = RefChecker(prs, report=True)

    if args.report_only:
        checker.check_all()
    else:
        checker.fix_all()

    report = checker.generate_report()

    # Print summary
    print(f"\n  Reference & Attribution Check")
    print(f"  {'=' * 40}")
    print(f"  Citations found:     {report['summary']['citations_found']}")
    print(f"  References found:    {report['summary']['references_found']}")
    print(f"  Reference slides:    {report['summary']['ref_slides']}")
    print(f"  Total issues:        {report['total_issues']}")
    print(f"  Total changes:       {report['total_changes']}")
    print()

    if report["issues"]:
        by_cat = report["summary"]["issues_by_category"]
        print(f"  Issues by category:")
        for cat, count in sorted(by_cat.items()):
            print(f"    {cat}: {count}")
        print()

        by_sev = report["summary"]["issues_by_severity"]
        print(f"  Issues by severity:")
        for sev, count in sorted(by_sev.items()):
            print(f"    {sev}: {count}")

    # Save output
    if not args.report_only and args.output:
        output_path = Path(args.output)
        prs.save(str(output_path))
        print(f"\n  Fixed file saved to: {output_path}")

    # Save report
    if args.report or args.report_only:
        report_path = input_path.with_suffix('.ref_report.json')
        if args.output:
            report_path = Path(args.output).with_suffix('.ref_report.json')
        with open(report_path, 'w') as f:
            json.dump(report, f, indent=2, default=str)
        print(f"  Report saved to:     {report_path}")


if __name__ == "__main__":
    main()
