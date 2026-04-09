"""
Layout Manager — UQ Slide Compliance Tool
==========================================
Analyses slide content, determines the best UQ template layout,
rebuilds each slide from scratch using the template, then verifies
the result via Claude Vision.

Components:
    LayoutRegistry    — Catalogue of all 46 UQ template layouts
    ContentExtractor  — Extracts structured content from any slide
    LayoutAnalyser    — Uses Claude Vision to pick the best layout
    SlideRebuilder    — Creates fresh slide from template + content
    LayoutVerifier    — Compares original vs rebuilt via Vision
    LayoutManager     — Orchestrates the full pipeline
"""

import os
import io
import re
import json
import copy
import base64
import hashlib
import tempfile
import subprocess
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional
from enum import Enum

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree

try:
    import anthropic
    HAS_ANTHROPIC = True
except ImportError:
    HAS_ANTHROPIC = False


# ============================================================
# Constants
# ============================================================

UQ_TEMPLATE_PATH = os.path.join(
    os.path.dirname(__file__),
    "Powerpoints",
    "UQ PPT Template - February 2026.pptx",
)

# Classification model for Vision calls
VISION_MODEL = "claude-sonnet-4-20250514"
ANALYSIS_MAX_TOKENS = 2000
VERIFY_MAX_TOKENS = 1500

# Cost constants (Sonnet input/output per token)
COST_INPUT_PER_TOKEN = 3.0 / 1_000_000
COST_OUTPUT_PER_TOKEN = 15.0 / 1_000_000


# ============================================================
# Layout Category Enum
# ============================================================

class LayoutCategory(str, Enum):
    COVER = "cover"
    SECTION = "section"
    CONTENTS = "contents"
    CONTENT = "content"              # Title + body text
    TWO_CONTENT = "two_content"      # Title + 2 body areas
    THREE_CONTENT = "three_content"  # Title + 3 body areas
    IMAGE_TEXT = "image_text"        # Text + image combo
    MULTI_IMAGE = "multi_image"      # Multiple images
    TABLE = "table"
    GRAPH = "graph"
    QUOTE = "quote"
    ICONS = "icons"
    SPECIAL = "special"              # Multi-layout, pullouts, etc.
    MINIMAL = "minimal"             # Title only, blank branded
    ENDING = "ending"               # Thank you


# ============================================================
# Layout Registry
# ============================================================

@dataclass
class PlaceholderSpec:
    """Specification for a single placeholder in a layout."""
    idx: int
    ph_type: str          # TITLE, BODY, OBJECT, PICTURE, TABLE, FOOTER, SLIDE_NUMBER
    name: str
    left_in: float        # Position in inches
    top_in: float
    width_in: float
    height_in: float
    role: str = ""        # Semantic role: title, subtitle, body, image, caption, icon, etc.


@dataclass
class LayoutSpec:
    """Full specification for a template layout."""
    index: int                           # Index in template.slide_layouts
    name: str                            # Layout name (e.g. "Title and Content")
    category: LayoutCategory
    placeholders: list                   # List of PlaceholderSpec
    has_solid_bg: bool = False
    has_gradient_bg: bool = False
    non_ph_shape_count: int = 0
    description: str = ""                # Human-readable description
    content_slots: dict = field(default_factory=dict)  # Semantic mapping: role → PH idx


class LayoutRegistry:
    """
    Catalogue of all UQ template layouts, built by introspecting the
    actual template file.  Also provides name-mapping from legacy
    layout names to canonical template names.
    """

    # ── Legacy name → canonical template name mapping ──────────
    # Covers prefixed duplicates (1_, 2_, 3_, 8_) and renamed layouts
    NAME_MAP = {
        # Prefixed duplicates → canonical
        "1_Title and Content":                      "Title and Content",
        "2_Title and Content":                      "Title and Content",
        "3_Title and Content":                      "Title and Content",
        "8_Title and Content":                      "Title and Content",
        "1_Two Content":                            "Two Content",
        "1_Two Content Layout Horizontal":          "Two Content Layout Horizontal",
        "1_One Third Two Third Title and Content":  "One Third Two Third Title and Content",
        "1_Text with Image One Third":              "Text with Image One Third",
        "1_Text with Image One Third Alt":          "Text with Image One Third Alt",
        "1_Text with Image Half":                   "Text with Image Half",
        "1_Text with Image Half Alt":               "Text with Image Half Alt",
        "1_Title Only":                             "Title Only",
        "1_Blank Branded":                          "Blank Branded",
        "1_Multi-layout 2":                         "Multi-layout 2",
        "1_Title and table_purple":                 "Title and Table",
        "2_Title and table":                        "Title and Table",

        # Renamed layouts
        "Section Divider 2":                        "Section Divider",
        "Cover 4":                                  "Cover 3",  # Closest match

        # Near-matches
        "Text with Neutral 1 Block":                "Text with Neutral Block",
        "Text with Purple Block":                   "Text with Dark Purple Block",
        "Heading + Subtitle + Body":                "Title and Content",
    }

    def __init__(self, template_path: str = None):
        self.template_path = template_path or UQ_TEMPLATE_PATH
        self.layouts: dict[str, LayoutSpec] = {}   # name → LayoutSpec
        self._template_prs = None
        self._build_registry()

    def _build_registry(self):
        """Introspect the template and build the full layout catalogue."""
        self._template_prs = Presentation(self.template_path)

        # Category assignment by name patterns
        category_rules = [
            (r"^Cover",              LayoutCategory.COVER),
            (r"^Section Divider",    LayoutCategory.SECTION),
            (r"^Contents",           LayoutCategory.CONTENTS),
            (r"^Thank You",          LayoutCategory.ENDING),
            (r"^Quote",              LayoutCategory.QUOTE),
            (r"^Title and Table",    LayoutCategory.TABLE),
            (r"^Title Only",         LayoutCategory.MINIMAL),
            (r"^Blank Branded",      LayoutCategory.MINIMAL),
            (r"^Icons",              LayoutCategory.ICONS),
            (r"^Order",              LayoutCategory.ICONS),
            (r"Graph",               LayoutCategory.GRAPH),
            (r"^Three content",      LayoutCategory.THREE_CONTENT),
            (r"^Two Content",        LayoutCategory.TWO_CONTENT),
            (r"^One Third Two Third",LayoutCategory.TWO_CONTENT),
            (r"^Two Third One Third",LayoutCategory.TWO_CONTENT),
            (r"^Intro \+ two",       LayoutCategory.TWO_CONTENT),
            (r"^Icons \+ two",       LayoutCategory.TWO_CONTENT),
            (r"^Title, Subtitle, 2", LayoutCategory.TWO_CONTENT),
            (r"Image",               LayoutCategory.IMAGE_TEXT),
            (r"Picture",             LayoutCategory.IMAGE_TEXT),
            (r"collage",             LayoutCategory.MULTI_IMAGE),
            (r"^Text with 4 Images", LayoutCategory.MULTI_IMAGE),
            (r"Three Column Text",   LayoutCategory.MULTI_IMAGE),
            (r"Pullout",             LayoutCategory.SPECIAL),
            (r"Multi-layout",        LayoutCategory.SPECIAL),
            (r"^Title and Content",  LayoutCategory.CONTENT),
            (r"^Text with.*Block",   LayoutCategory.SPECIAL),
            (r"^Three Pullouts",     LayoutCategory.SPECIAL),
        ]

        ph_type_names = {
            1: "TITLE", 2: "BODY", 7: "OBJECT", 12: "TABLE",
            13: "SLIDE_NUMBER", 15: "FOOTER", 18: "PICTURE",
        }

        for i, layout in enumerate(self._template_prs.slide_layouts):
            name = layout.name

            # Determine category
            category = LayoutCategory.CONTENT  # default
            for pattern, cat in category_rules:
                if re.search(pattern, name, re.IGNORECASE):
                    category = cat
                    break

            # Build placeholder specs
            ph_specs = []
            for ph in layout.placeholders:
                pt = ph.placeholder_format.type
                type_val = pt.value if hasattr(pt, 'value') else (pt if isinstance(pt, int) else 0)
                type_name = ph_type_names.get(type_val, f"UNKNOWN({type_val})")

                ph_specs.append(PlaceholderSpec(
                    idx=ph.placeholder_format.idx,
                    ph_type=type_name,
                    name=ph.name,
                    left_in=round(ph.left / 914400, 2) if ph.left else 0,
                    top_in=round(ph.top / 914400, 2) if ph.top else 0,
                    width_in=round(ph.width / 914400, 2) if ph.width else 0,
                    height_in=round(ph.height / 914400, 2) if ph.height else 0,
                ))

            # Background type
            bg_fill = layout.background.fill
            has_solid = bg_fill.type is not None and str(bg_fill.type) == "SOLID (1)"
            has_gradient = bg_fill.type is not None and str(bg_fill.type) == "GRADIENT (3)"

            # Non-placeholder shapes
            non_ph = [s for s in layout.shapes if not s.is_placeholder]

            # Build content_slots: semantic role → placeholder idx
            content_slots = {}
            for ps in ph_specs:
                if ps.ph_type == "TITLE":
                    content_slots["title"] = ps.idx
                elif ps.ph_type == "FOOTER":
                    content_slots.setdefault("footer", ps.idx)
                elif ps.ph_type == "SLIDE_NUMBER":
                    content_slots.setdefault("slide_number", ps.idx)
                elif ps.ph_type == "PICTURE":
                    pics = content_slots.get("pictures", [])
                    pics.append(ps.idx)
                    content_slots["pictures"] = pics
                elif ps.ph_type == "TABLE":
                    content_slots["table"] = ps.idx
                elif ps.ph_type in ("BODY", "OBJECT"):
                    bodies = content_slots.get("bodies", [])
                    bodies.append(ps.idx)
                    content_slots["bodies"] = bodies

            self.layouts[name] = LayoutSpec(
                index=i,
                name=name,
                category=category,
                placeholders=ph_specs,
                has_solid_bg=has_solid,
                has_gradient_bg=has_gradient,
                non_ph_shape_count=len(non_ph),
                content_slots=content_slots,
            )

    def get_layout(self, name: str) -> Optional[LayoutSpec]:
        """Get layout spec by exact name."""
        return self.layouts.get(name)

    def resolve_name(self, legacy_name: str) -> str:
        """Map a legacy/prefixed layout name to the canonical template name."""
        if legacy_name in self.layouts:
            return legacy_name
        return self.NAME_MAP.get(legacy_name, legacy_name)

    def get_template_layout_obj(self, name: str):
        """Get the actual python-pptx SlideLayout object by name."""
        spec = self.layouts.get(name)
        if spec is None:
            return None
        return self._template_prs.slide_layouts[spec.index]

    def get_all_names(self) -> list[str]:
        """Return all template layout names."""
        return list(self.layouts.keys())

    def get_layouts_by_category(self, category: LayoutCategory) -> list[LayoutSpec]:
        """Return all layouts in a given category."""
        return [l for l in self.layouts.values() if l.category == category]

    def guess_layout_from_content(self, slide_content) -> tuple[str, float]:
        """
        Content-based heuristic fallback: given a SlideContent object,
        guess the best template layout based on what's actually on the slide.

        Returns (layout_name, confidence).
        """
        has_title = slide_content.has_title
        has_subtitle = slide_content.has_subtitle
        n_body = slide_content.body_text_count
        n_img = slide_content.image_count
        n_tbl = slide_content.table_count
        title_text = slide_content.title.plain_text.lower() if has_title else ""

        # ── Cover / ending detection ──
        cover_keywords = ["welcome", "thank you", "thanks", "questions"]
        ending_keywords = ["thank you", "thanks", "questions", "q&a", "contact"]
        if any(kw in title_text for kw in ending_keywords):
            return ("Thank You", 0.7)
        if slide_content.slide_number == 1 or any(kw in title_text for kw in cover_keywords):
            if n_img == 0 and n_body <= 2:
                return ("Cover 1", 0.6)

        # ── Section divider detection ──
        # Short title, minimal body, no images
        if has_title and n_body <= 1 and n_img == 0 and n_tbl == 0:
            total_chars = sum(
                len(r.text) for b in slide_content.body_texts
                for p in b.paragraphs for r in p.runs
            ) if slide_content.body_texts else 0
            if total_chars < 100:
                return ("Section Divider", 0.6)

        # ── Table slides ──
        if n_tbl > 0:
            return ("Title and Table", 0.8)

        # ── Image + text combos ──
        if n_img > 0 and n_body > 0:
            if n_img >= 3:
                return ("Three Column Text & Images", 0.5)
            if n_img == 2:
                return ("Two Content", 0.6)
            # 1 image + text
            return ("Text with Image Half", 0.7)

        # ── Image-only slides ──
        if n_img > 0 and n_body == 0:
            if n_img >= 4:
                return ("Text with 4 Images", 0.5)
            return ("Picture with Caption", 0.5)

        # ── Multi-body text ──
        if n_body >= 3:
            return ("Three content layout", 0.6)
        if n_body == 2:
            return ("Two Content", 0.7)

        # ── Single body text (most common) ──
        if n_body == 1 or has_title:
            return ("Title and Content", 0.8)

        # ── Fallback: blank branded ──
        if not has_title and n_body == 0:
            return ("Blank Branded", 0.5)

        return ("Title and Content", 0.5)

    def get_layout_summary_for_prompt(self) -> str:
        """
        Generate a concise summary of all layouts suitable for
        inclusion in a Claude Vision prompt.
        """
        lines = []
        for name, spec in self.layouts.items():
            # Count content placeholders (not footer/slide_number)
            content_phs = [p for p in spec.placeholders
                          if p.ph_type not in ("FOOTER", "SLIDE_NUMBER")]
            title_count = sum(1 for p in content_phs if p.ph_type == "TITLE")
            body_count = sum(1 for p in content_phs if p.ph_type in ("BODY", "OBJECT"))
            pic_count = sum(1 for p in content_phs if p.ph_type == "PICTURE")
            table_count = sum(1 for p in content_phs if p.ph_type == "TABLE")

            parts = []
            if title_count:
                parts.append(f"{title_count} title")
            if body_count:
                parts.append(f"{body_count} text")
            if pic_count:
                parts.append(f"{pic_count} image")
            if table_count:
                parts.append(f"{table_count} table")

            bg = "solid" if spec.has_solid_bg else ("gradient" if spec.has_gradient_bg else "standard")

            lines.append(
                f'- "{name}" [{spec.category.value}]: '
                f'{", ".join(parts) if parts else "no content placeholders"} '
                f'(bg: {bg})'
            )
        return "\n".join(lines)


# ============================================================
# Content Extractor
# ============================================================

@dataclass
class TextRun:
    """A single run of text with consistent formatting."""
    text: str
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    font_name: Optional[str] = None
    font_size_pt: Optional[float] = None
    color_hex: Optional[str] = None
    hyperlink_url: Optional[str] = None


@dataclass
class TextParagraph:
    """A paragraph containing one or more runs."""
    runs: list         # List of TextRun
    alignment: Optional[str] = None    # LEFT, CENTER, RIGHT, JUSTIFY
    level: int = 0                     # Indentation level (for bullets)
    bullet_char: Optional[str] = None  # Custom bullet character if any
    space_before_pt: Optional[float] = None
    space_after_pt: Optional[float] = None
    line_spacing_pt: Optional[float] = None


@dataclass
class TextContent:
    """Full text content from a shape/placeholder."""
    paragraphs: list   # List of TextParagraph
    role: str = ""     # title, subtitle, body, caption, attribution, etc.

    @property
    def plain_text(self) -> str:
        lines = []
        for para in self.paragraphs:
            lines.append("".join(r.text for r in para.runs))
        return "\n".join(lines)

    @property
    def is_empty(self) -> bool:
        return not self.plain_text.strip()


@dataclass
class ImageContent:
    """An image extracted from a slide."""
    image_bytes: bytes
    format: str = "png"           # png, jpg, etc.
    width_in: float = 0
    height_in: float = 0
    left_in: float = 0
    top_in: float = 0
    alt_text: str = ""
    crop_info: Optional[dict] = None
    sha256: str = ""


@dataclass
class TableContent:
    """A table extracted from a slide."""
    rows: list          # List of lists of TextContent (one per cell)
    col_widths: list    # List of widths in inches
    row_heights: list   # List of heights in inches


@dataclass
class ShapeContent:
    """A non-text, non-image, non-table shape."""
    shape_type: str     # AUTO_SHAPE, FREEFORM, GROUP, etc.
    name: str
    left_in: float = 0
    top_in: float = 0
    width_in: float = 0
    height_in: float = 0
    text_content: Optional[TextContent] = None
    has_image: bool = False
    image_content: Optional[ImageContent] = None


@dataclass
class SlideContent:
    """Complete extracted content from a single slide."""
    slide_number: int
    original_layout_name: str
    title: Optional[TextContent] = None
    subtitle: Optional[TextContent] = None
    body_texts: list = field(default_factory=list)    # List of TextContent
    images: list = field(default_factory=list)         # List of ImageContent
    tables: list = field(default_factory=list)         # List of TableContent
    other_shapes: list = field(default_factory=list)   # List of ShapeContent
    speaker_notes: str = ""

    @property
    def has_title(self) -> bool:
        return self.title is not None and not self.title.is_empty

    @property
    def has_subtitle(self) -> bool:
        return self.subtitle is not None and not self.subtitle.is_empty

    @property
    def image_count(self) -> int:
        return len(self.images)

    @property
    def table_count(self) -> int:
        return len(self.tables)

    @property
    def body_text_count(self) -> int:
        return len([b for b in self.body_texts if not b.is_empty])

    def content_summary(self) -> str:
        """One-line summary for debugging."""
        parts = []
        if self.has_title:
            parts.append(f'title="{self.title.plain_text[:40]}"')
        if self.has_subtitle:
            parts.append(f'subtitle="{self.subtitle.plain_text[:30]}"')
        parts.append(f"{self.body_text_count} body")
        parts.append(f"{self.image_count} img")
        parts.append(f"{self.table_count} tbl")
        if self.other_shapes:
            parts.append(f"{len(self.other_shapes)} other")
        return f"Slide {self.slide_number} [{self.original_layout_name}]: {', '.join(parts)}"


class ContentExtractor:
    """
    Extracts all content from a slide into a structured SlideContent object.
    Handles both placeholder-based and freeform content.
    """

    # Patterns that indicate attribution text (not body content)
    ATTRIBUTION_PATTERNS = [
        r"source:\s", r"adobe\s*stock", r"shutterstock", r"microsoft\s*stock",
        r"cc\s*by", r"public\s*domain", r"wikimedia", r"flickr",
        r"image\s*licensed", r"getty\s*images", r"unsplash",
    ]

    def extract_slide(self, slide, slide_number: int) -> SlideContent:
        """Extract all content from a single slide."""
        content = SlideContent(
            slide_number=slide_number,
            original_layout_name=slide.slide_layout.name,
        )

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                content.speaker_notes = notes_tf.text.strip()

        # Process all shapes
        for shape in slide.shapes:
            self._process_shape(shape, content)

        return content

    def _process_shape(self, shape, content: SlideContent):
        """Route a shape to the appropriate extraction method."""
        if shape.is_placeholder:
            self._process_placeholder(shape, content)
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            self._process_table(shape, content)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            self._process_image(shape, content)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            self._process_group(shape, content)
        elif hasattr(shape, "text_frame"):
            self._process_freeform_text(shape, content)
        elif hasattr(shape, "image"):
            self._process_image(shape, content)
        else:
            self._process_other_shape(shape, content)

    def _process_placeholder(self, shape, content: SlideContent):
        """Extract content from a placeholder shape."""
        ph = shape.placeholder_format
        ph_type = ph.type
        # Get the integer type value
        type_val = ph_type.value if hasattr(ph_type, 'value') else (ph_type if isinstance(ph_type, int) else 0)

        # Skip footer and slide number
        if type_val in (13, 15):  # SLIDE_NUMBER, FOOTER
            return

        # TITLE
        if type_val == 1:
            if hasattr(shape, "text_frame"):
                content.title = self._extract_text_content(shape.text_frame, role="title")
            return

        # PICTURE placeholder
        if type_val == 18:
            try:
                if hasattr(shape, "image") and shape.image:
                    content.images.append(self._extract_image(shape))
            except (ValueError, AttributeError):
                pass  # Empty picture placeholder
            return

        # TABLE placeholder
        if type_val == 12:
            if shape.has_table:
                content.tables.append(self._extract_table(shape.table, shape))
            return

        # BODY or OBJECT — could be subtitle or body content
        if hasattr(shape, "text_frame"):
            tc = self._extract_text_content(shape.text_frame, role="body")
            if not tc.is_empty:
                # Check if this is a subtitle (PH 31 in many layouts, or PH 11 in covers)
                if ph.idx in (31, 11) and content.subtitle is None:
                    tc.role = "subtitle"
                    content.subtitle = tc
                else:
                    content.body_texts.append(tc)

    def _process_table(self, shape, content: SlideContent):
        """Extract a table from a shape."""
        if shape.has_table:
            content.tables.append(self._extract_table(shape.table, shape))

    def _process_image(self, shape, content: SlideContent):
        """Extract an image from a shape."""
        try:
            if hasattr(shape, "image") and shape.image:
                content.images.append(self._extract_image(shape))
        except (ValueError, AttributeError):
            pass

    def _process_group(self, shape, content: SlideContent):
        """Process shapes within a group."""
        for child in shape.shapes:
            self._process_shape(child, content)

    def _process_freeform_text(self, shape, content: SlideContent):
        """Extract text from a non-placeholder text shape (text boxes, freeforms, auto shapes)."""
        if not hasattr(shape, "text_frame"):
            # Check if it's an image-bearing shape
            if hasattr(shape, "image"):
                try:
                    content.images.append(self._extract_image(shape))
                except (ValueError, AttributeError):
                    pass
            return

        tc = self._extract_text_content(shape.text_frame, role="body")
        if tc.is_empty:
            # Even empty text shapes might contain images
            if hasattr(shape, "image"):
                try:
                    content.images.append(self._extract_image(shape))
                except (ValueError, AttributeError):
                    pass
            return

        # Check if this looks like attribution text
        plain = tc.plain_text.lower()
        is_attribution = any(re.search(p, plain) for p in self.ATTRIBUTION_PATTERNS)
        if is_attribution:
            tc.role = "attribution"

        content.body_texts.append(tc)

        # Also check if the shape contains an image
        if hasattr(shape, "image"):
            try:
                content.images.append(self._extract_image(shape))
            except (ValueError, AttributeError):
                pass

    def _process_other_shape(self, shape, content: SlideContent):
        """Extract info from other shape types (auto shapes, freeforms, etc.)."""
        sc = ShapeContent(
            shape_type=str(shape.shape_type),
            name=shape.name,
            left_in=round(shape.left / 914400, 2) if shape.left else 0,
            top_in=round(shape.top / 914400, 2) if shape.top else 0,
            width_in=round(shape.width / 914400, 2) if shape.width else 0,
            height_in=round(shape.height / 914400, 2) if shape.height else 0,
        )
        if hasattr(shape, "text_frame"):
            sc.text_content = self._extract_text_content(shape.text_frame, role="other")
        if hasattr(shape, "image"):
            try:
                sc.image_content = self._extract_image(shape)
                sc.has_image = True
            except (ValueError, AttributeError):
                pass
        content.other_shapes.append(sc)

    def _extract_text_content(self, text_frame, role: str = "") -> TextContent:
        """Extract structured text from a text frame."""
        paragraphs = []
        for para in text_frame.paragraphs:
            runs = []
            for run in para.runs:
                # Extract colour
                color_hex = None
                try:
                    if run.font.color and run.font.color.rgb:
                        color_hex = str(run.font.color.rgb)
                except (AttributeError, TypeError):
                    pass

                # Extract font size
                font_size = None
                if run.font.size:
                    font_size = run.font.size.pt

                # Extract hyperlink
                hyperlink = None
                if run.hyperlink and run.hyperlink.address:
                    hyperlink = run.hyperlink.address

                runs.append(TextRun(
                    text=run.text,
                    bold=run.font.bold,
                    italic=run.font.italic,
                    underline=run.font.underline,
                    font_name=run.font.name,
                    font_size_pt=font_size,
                    color_hex=color_hex,
                    hyperlink_url=hyperlink,
                ))

            # Paragraph-level properties
            alignment = None
            if para.alignment is not None:
                alignment = str(para.alignment)

            # Bullet info
            bullet_char = None
            try:
                pPr = para._p.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                if pPr is not None:
                    bullet_char = pPr.get('char')
            except Exception:
                pass

            # Spacing
            space_before = None
            space_after = None
            line_spacing = None
            try:
                if para.space_before:
                    space_before = para.space_before.pt
                if para.space_after:
                    space_after = para.space_after.pt
                if para.line_spacing:
                    line_spacing = para.line_spacing.pt
            except (AttributeError, TypeError):
                pass

            paragraphs.append(TextParagraph(
                runs=runs,
                alignment=alignment,
                level=para.level if para.level else 0,
                bullet_char=bullet_char,
                space_before_pt=space_before,
                space_after_pt=space_after,
                line_spacing_pt=line_spacing,
            ))

        return TextContent(paragraphs=paragraphs, role=role)

    def _extract_image(self, shape) -> ImageContent:
        """Extract image bytes and metadata from a shape."""
        img = shape.image
        img_bytes = img.blob

        # Determine format
        content_type = img.content_type
        fmt = "png"
        if "jpeg" in content_type or "jpg" in content_type:
            fmt = "jpg"
        elif "gif" in content_type:
            fmt = "gif"
        elif "tiff" in content_type:
            fmt = "tiff"
        elif "emf" in content_type:
            fmt = "emf"
        elif "wmf" in content_type:
            fmt = "wmf"

        return ImageContent(
            image_bytes=img_bytes,
            format=fmt,
            width_in=round(shape.width / 914400, 2) if shape.width else 0,
            height_in=round(shape.height / 914400, 2) if shape.height else 0,
            left_in=round(shape.left / 914400, 2) if shape.left else 0,
            top_in=round(shape.top / 914400, 2) if shape.top else 0,
            alt_text=shape.name or "",
            sha256=hashlib.sha256(img_bytes).hexdigest()[:16],
        )

    def _extract_table(self, table, shape) -> TableContent:
        """Extract table data."""
        rows_data = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                tc = self._extract_text_content(cell.text_frame, role="table_cell")
                cells.append(tc)
            rows_data.append(cells)

        col_widths = [round(col.width / 914400, 2) for col in table.columns]
        row_heights = [round(row.height / 914400, 2) for row in table.rows]

        return TableContent(
            rows=rows_data,
            col_widths=col_widths,
            row_heights=row_heights,
        )


# ============================================================
# Slide Renderer (for Vision analysis)
# ============================================================

class SlideRenderer:
    """
    Renders individual slides to images for Claude Vision analysis.
    Uses LibreOffice → PDF → pdftoppm pipeline.
    """

    SOFFICE_SCRIPT = os.path.join(
        os.path.dirname(__file__),
        "..", ".claude", "skills", "pptx", "scripts", "office", "soffice.py"
    )

    @staticmethod
    def render_deck_to_images(pptx_path: str, output_dir: str, dpi: int = 150) -> list[str]:
        """
        Render all slides in a PPTX to individual JPEG images.
        Returns list of image file paths.
        """
        os.makedirs(output_dir, exist_ok=True)

        # Step 1: Convert PPTX → PDF
        pdf_path = os.path.join(output_dir, "deck.pdf")

        # Try soffice directly
        try:
            subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf",
                 "--outdir", output_dir, pptx_path],
                capture_output=True, timeout=120, check=True,
            )
            # soffice names the output after the input file
            src_pdf = os.path.join(output_dir,
                                   Path(pptx_path).stem + ".pdf")
            if os.path.exists(src_pdf) and src_pdf != pdf_path:
                os.rename(src_pdf, pdf_path)
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError) as e:
            raise RuntimeError(f"LibreOffice PDF conversion failed: {e}")

        if not os.path.exists(pdf_path):
            raise RuntimeError(f"PDF not created at {pdf_path}")

        # Step 2: PDF → individual JPEG images
        subprocess.run(
            ["pdftoppm", "-jpeg", "-r", str(dpi), pdf_path,
             os.path.join(output_dir, "slide")],
            capture_output=True, timeout=120, check=True,
        )

        # Collect generated images
        images = sorted(
            [os.path.join(output_dir, f) for f in os.listdir(output_dir)
             if f.startswith("slide") and f.endswith(".jpg")],
        )
        return images

    @staticmethod
    def render_single_slide(pptx_bytes: bytes, slide_index: int,
                            output_path: str, dpi: int = 150) -> str:
        """Render a single slide from PPTX bytes to a JPEG image."""
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = os.path.join(tmpdir, "deck.pptx")
            with open(pptx_path, "wb") as f:
                f.write(pptx_bytes)

            images = SlideRenderer.render_deck_to_images(pptx_path, tmpdir, dpi)
            if slide_index < len(images):
                # Copy the target image to output
                import shutil
                shutil.copy2(images[slide_index], output_path)
                return output_path

        raise ValueError(f"Slide {slide_index} not rendered (only {len(images)} slides)")


# ============================================================
# Layout Analyser (Claude Vision)
# ============================================================

class LayoutAnalyser:
    """
    Uses Claude Vision to analyse each source slide image and
    recommend the best UQ template layout.
    """

    def __init__(self, api_key: str = None, registry: LayoutRegistry = None):
        if not HAS_ANTHROPIC:
            raise ImportError("anthropic package required for Vision analysis")
        self.client = anthropic.Anthropic(api_key=api_key or os.environ.get("ANTHROPIC_API_KEY"))
        self.registry = registry or LayoutRegistry()
        self._layout_catalogue = self.registry.get_layout_summary_for_prompt()

    def analyse_slide(self, slide_image_path: str, slide_content: SlideContent) -> dict:
        """
        Send a slide image to Claude Vision for layout recommendation.

        Returns dict with:
            - recommended_layout: str (template layout name)
            - confidence: float (0-1)
            - reasoning: str
            - content_inventory: dict (what Vision sees on the slide)
            - input_tokens: int
            - output_tokens: int
        """
        # Read image
        with open(slide_image_path, "rb") as f:
            image_data = base64.standard_b64encode(f.read()).decode("utf-8")

        # Build the content inventory from extraction (helps Vision confirm)
        extraction_summary = slide_content.content_summary()

        prompt = f"""You are analysing a PowerPoint slide to determine which UQ template layout would best suit its content.

## Extracted Content Summary
{extraction_summary}

## Available UQ Template Layouts
{self._layout_catalogue}

## Your Task
Look at this slide image and determine:

1. **What content is on this slide?** List: title text, subtitle, body text blocks (how many), images (how many and approximate position — left/right/top/bottom/full), tables, charts/diagrams, icons.

2. **What type of slide is this?** Cover/title slide, section divider, content slide, image+text, table, quote, thank you/ending, etc.

3. **Which template layout is the BEST fit?** Consider:
   - Does the layout have the right number and type of content placeholders?
   - Does the image position match (left vs right, half vs one-third vs two-thirds)?
   - Is the overall structure compatible?
   - For text-heavy slides with no images, prefer "Title and Content" or "Two Content"
   - For slides with an image on one side, choose the appropriate "Text with Image" variant
   - For slides with coloured side panels, consider "Graph/Text with Dark Purple/Neutral/Grey Block"

4. **How confident are you?** (0.0 to 1.0)

Respond in this exact JSON format:
```json
{{
    "recommended_layout": "exact layout name from the list above",
    "confidence": 0.85,
    "slide_type": "content|cover|section|image_text|table|quote|ending|other",
    "reasoning": "Brief explanation of why this layout fits",
    "content_inventory": {{
        "has_title": true,
        "title_text": "first 50 chars of title",
        "has_subtitle": false,
        "body_text_blocks": 1,
        "image_count": 0,
        "image_positions": [],
        "has_table": false,
        "has_chart": false,
        "has_icons": false,
        "special_elements": []
    }}
}}
```

Respond with ONLY the JSON, no other text."""

        response = self.client.messages.create(
            model=VISION_MODEL,
            max_tokens=ANALYSIS_MAX_TOKENS,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": "image/jpeg",
                            "data": image_data,
                        },
                    },
                    {
                        "type": "text",
                        "text": prompt,
                    },
                ],
            }],
        )

        # Parse response
        result_text = response.content[0].text.strip()

        # Extract JSON from response (handle markdown code blocks)
        json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', result_text, re.DOTALL)
        if json_match:
            result_text = json_match.group(1)
        elif result_text.startswith("{"):
            pass  # Already JSON
        else:
            # Try to find JSON object
            start = result_text.find("{")
            end = result_text.rfind("}") + 1
            if start >= 0 and end > start:
                result_text = result_text[start:end]

        try:
            result = json.loads(result_text)
        except json.JSONDecodeError:
            result = {
                "recommended_layout": "Title and Content",
                "confidence": 0.3,
                "reasoning": f"Failed to parse Vision response: {result_text[:200]}",
                "content_inventory": {},
                "slide_type": "content",
            }

        # Validate the recommended layout exists
        rec = result.get("recommended_layout", "")
        if rec not in self.registry.layouts:
            # Try to fuzzy match
            best_match = None
            for name in self.registry.layouts:
                if name.lower() == rec.lower():
                    best_match = name
                    break
            if best_match:
                result["recommended_layout"] = best_match
            else:
                result["recommended_layout_original"] = rec
                result["recommended_layout"] = "Title and Content"
                result["confidence"] = min(result.get("confidence", 0.5), 0.5)

        # Add token usage
        result["input_tokens"] = response.usage.input_tokens
        result["output_tokens"] = response.usage.output_tokens

        return result


# ============================================================
# Slide Rebuilder
# ============================================================

class SlideRebuilder:
    """
    Creates a fresh slide from a UQ template layout and populates
    it with extracted content.
    """

    def __init__(self, registry: LayoutRegistry = None):
        self.registry = registry or LayoutRegistry()

    def rebuild_slide(self, target_prs: Presentation, slide_content: SlideContent,
                      layout_name: str) -> tuple:
        """
        Create a new slide using the specified layout and populate
        it with the extracted content.

        Args:
            target_prs: The target Presentation to add the slide to
            slide_content: Extracted content from the source slide
            layout_name: Name of the template layout to use

        Returns:
            (slide, report) where report is a dict of what was done
        """
        layout_spec = self.registry.get_layout(layout_name)
        if layout_spec is None:
            raise ValueError(f"Layout '{layout_name}' not found in registry")

        # Get the actual layout object from the TARGET presentation
        # We need to find the matching layout in the target presentation
        target_layout = None
        for layout in target_prs.slide_layouts:
            if layout.name == layout_name:
                target_layout = layout
                break

        if target_layout is None:
            raise ValueError(
                f"Layout '{layout_name}' not found in target presentation. "
                f"Make sure the target was created from the UQ template."
            )

        # Create new slide
        slide = target_prs.slides.add_slide(target_layout)

        report = {
            "layout_applied": layout_name,
            "content_placed": [],
            "content_skipped": [],
            "warnings": [],
        }

        slots = layout_spec.content_slots

        # ── Place title ──
        if slide_content.has_title and "title" in slots:
            self._place_text_in_placeholder(
                slide, slots["title"], slide_content.title
            )
            report["content_placed"].append("title")
        elif slide_content.has_title:
            report["content_skipped"].append("title (no placeholder)")

        # ── Place subtitle ──
        if slide_content.has_subtitle:
            bodies = slots.get("bodies", [])
            # Subtitle typically goes in PH 31 or 11
            subtitle_ph = None
            for idx in [31, 11]:
                if idx in bodies:
                    subtitle_ph = idx
                    break
            if subtitle_ph is not None:
                self._place_text_in_placeholder(
                    slide, subtitle_ph, slide_content.subtitle
                )
                report["content_placed"].append("subtitle")
                # Remove from bodies list so we don't double-use
                bodies = [b for b in bodies if b != subtitle_ph]
            else:
                report["content_skipped"].append("subtitle (no placeholder)")

        # ── Place images ──
        pic_slots = slots.get("pictures", [])
        for i, img in enumerate(slide_content.images):
            if i < len(pic_slots):
                self._place_image_in_placeholder(slide, pic_slots[i], img)
                report["content_placed"].append(f"image_{i}")
            else:
                # Place as free-form shape
                self._place_image_freeform(slide, img)
                report["content_placed"].append(f"image_{i} (freeform)")

        # ── Place tables ──
        if slide_content.tables:
            table_slot = slots.get("table")
            if table_slot is not None:
                self._place_table_in_placeholder(
                    slide, table_slot, slide_content.tables[0]
                )
                report["content_placed"].append("table")
            else:
                # Place as freeform table
                self._place_table_freeform(slide, slide_content.tables[0])
                report["content_placed"].append("table (freeform)")
            if len(slide_content.tables) > 1:
                report["warnings"].append(
                    f"{len(slide_content.tables) - 1} additional tables skipped"
                )

        # ── Place body texts ──
        body_slots = slots.get("bodies", [])
        # Remove subtitle slot if we already used it
        if slide_content.has_subtitle:
            body_slots = [b for b in body_slots if b not in (31, 11)]

        # Filter out non-empty body texts, separate attribution
        body_texts = [b for b in slide_content.body_texts
                      if not b.is_empty and b.role != "attribution"]
        attributions = [b for b in slide_content.body_texts
                       if b.role == "attribution"]

        for i, body in enumerate(body_texts):
            if i < len(body_slots):
                self._place_text_in_placeholder(slide, body_slots[i], body)
                report["content_placed"].append(f"body_{i}")
            else:
                # Place as freeform text box
                self._place_text_freeform(slide, body, slide_content)
                report["content_placed"].append(f"body_{i} (freeform)")

        # Place attribution text (small, near bottom)
        for attr in attributions:
            self._place_attribution(slide, attr)
            report["content_placed"].append("attribution")

        # ── Place other shapes with text ──
        for shape_content in slide_content.other_shapes:
            if shape_content.text_content and not shape_content.text_content.is_empty:
                report["warnings"].append(
                    f"Shape '{shape_content.name}' has text content that may need manual placement"
                )

        # ── Speaker notes ──
        if slide_content.speaker_notes:
            if slide.has_notes_slide or True:  # Always try
                try:
                    notes_slide = slide.notes_slide
                    notes_tf = notes_slide.notes_text_frame
                    notes_tf.text = slide_content.speaker_notes
                    report["content_placed"].append("speaker_notes")
                except Exception:
                    report["warnings"].append("Could not place speaker notes")

        return slide, report

    def _place_text_in_placeholder(self, slide, ph_idx: int, text_content: TextContent):
        """Place text content into a placeholder by index."""
        # Find the placeholder on the slide
        ph = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == ph_idx:
                ph = shape
                break

        if ph is None or not hasattr(ph, "text_frame"):
            return

        tf = ph.text_frame

        # Clear existing text
        for i in range(len(tf.paragraphs) - 1, 0, -1):
            p = tf.paragraphs[i]._p
            p.getparent().remove(p)

        # Populate with extracted content
        for i, para_content in enumerate(text_content.paragraphs):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()

            # Set paragraph properties
            if para_content.alignment:
                try:
                    align_map = {
                        "LEFT (0)": PP_ALIGN.LEFT,
                        "CENTER (1)": PP_ALIGN.CENTER,
                        "RIGHT (2)": PP_ALIGN.RIGHT,
                        "JUSTIFY (3)": PP_ALIGN.JUSTIFY,
                    }
                    for key, val in align_map.items():
                        if key in str(para_content.alignment):
                            para.alignment = val
                            break
                except Exception:
                    pass

            para.level = para_content.level

            # Set spacing
            if para_content.space_before_pt is not None:
                try:
                    para.space_before = Pt(para_content.space_before_pt)
                except Exception:
                    pass
            if para_content.space_after_pt is not None:
                try:
                    para.space_after = Pt(para_content.space_after_pt)
                except Exception:
                    pass

            # Add runs
            for j, run_content in enumerate(para_content.runs):
                if j == 0 and i == 0:
                    # Use existing first run if possible
                    if len(para.runs) > 0:
                        run = para.runs[0]
                    else:
                        run = para.add_run()
                else:
                    run = para.add_run()

                run.text = run_content.text

                # Apply formatting
                if run_content.bold is not None:
                    run.font.bold = run_content.bold
                if run_content.italic is not None:
                    run.font.italic = run_content.italic
                if run_content.underline is not None:
                    run.font.underline = run_content.underline
                if run_content.font_name:
                    run.font.name = "Arial"  # Always use Arial for UQ compliance
                if run_content.font_size_pt:
                    run.font.size = Pt(run_content.font_size_pt)
                if run_content.color_hex:
                    try:
                        run.font.color.rgb = RGBColor.from_string(run_content.color_hex)
                    except (ValueError, AttributeError):
                        pass

            # Set bullet if specified
            if para_content.bullet_char:
                try:
                    pPr = para._p.get_or_add_pPr()
                    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    buChar = etree.SubElement(pPr, f'{{{nsmap["a"]}}}buChar')
                    buChar.set('char', para_content.bullet_char)
                except Exception:
                    pass

    def _place_image_in_placeholder(self, slide, ph_idx: int, image_content: ImageContent):
        """Place an image into a picture placeholder."""
        ph = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == ph_idx:
                ph = shape
                break

        if ph is None:
            self._place_image_freeform(slide, image_content)
            return

        # Insert image into the placeholder
        try:
            # Save image to temp file
            ext = image_content.format
            if ext == "jpg":
                ext = "jpeg"
            with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                tmp.write(image_content.image_bytes)
                tmp_path = tmp.name
            try:
                ph.insert_picture(tmp_path)
            finally:
                os.unlink(tmp_path)
        except Exception:
            # Fall back to freeform placement
            self._place_image_freeform(slide, image_content)

    def _place_image_freeform(self, slide, image_content: ImageContent):
        """Place an image as a free-form shape on the slide."""
        ext = image_content.format
        if ext == "jpg":
            ext = "jpeg"
        with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
            tmp.write(image_content.image_bytes)
            tmp_path = tmp.name
        try:
            left = Inches(image_content.left_in) if image_content.left_in else Inches(0.5)
            top = Inches(image_content.top_in) if image_content.top_in else Inches(1.5)
            width = Inches(image_content.width_in) if image_content.width_in else Inches(4)
            height = Inches(image_content.height_in) if image_content.height_in else Inches(3)
            slide.shapes.add_picture(tmp_path, left, top, width, height)
        finally:
            os.unlink(tmp_path)

    def _place_table_in_placeholder(self, slide, ph_idx: int, table_content: TableContent):
        """Place a table — python-pptx doesn't support table placeholders directly,
        so we create a freeform table at the placeholder's position."""
        # Find placeholder position
        ph = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == ph_idx:
                ph = shape
                break

        if ph is not None:
            left, top, width, height = ph.left, ph.top, ph.width, ph.height
        else:
            left, top = Inches(0.5), Inches(2.5)
            width, height = Inches(12.3), Inches(4.5)

        self._create_table(slide, table_content, left, top, width, height)

    def _place_table_freeform(self, slide, table_content: TableContent):
        """Place a table as a freeform shape."""
        self._create_table(
            slide, table_content,
            Inches(0.5), Inches(2.5), Inches(12.3), Inches(4.5)
        )

    def _create_table(self, slide, table_content: TableContent,
                      left, top, width, height):
        """Create an actual table shape on the slide."""
        rows = len(table_content.rows)
        cols = len(table_content.rows[0]) if rows > 0 else 1

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Set column widths if available
        if table_content.col_widths:
            for i, w in enumerate(table_content.col_widths):
                if i < len(table.columns):
                    table.columns[i].width = Inches(w)

        # Populate cells
        for r_idx, row_data in enumerate(table_content.rows):
            for c_idx, cell_content in enumerate(row_data):
                if r_idx < rows and c_idx < cols:
                    cell = table.cell(r_idx, c_idx)
                    # Set text from the TextContent
                    if cell_content and not cell_content.is_empty:
                        tf = cell.text_frame
                        # Clear and populate
                        for p_idx, para in enumerate(cell_content.paragraphs):
                            if p_idx == 0:
                                p = tf.paragraphs[0]
                            else:
                                p = tf.add_paragraph()
                            text = "".join(r.text for r in para.runs)
                            p.text = text

    def _place_text_freeform(self, slide, text_content: TextContent,
                             slide_content: SlideContent):
        """Place a text block as a free-form text box, using original position if available."""
        from pptx.util import Inches, Pt

        # Try to use a reasonable position based on the content area
        # Count existing freeform shapes to avoid overlap
        existing_freeform = len([s for s in slide.shapes
                                if not s.is_placeholder and hasattr(s, 'text_frame')])

        # Estimate text length to set appropriate height
        total_chars = sum(len(r.text) for p in text_content.paragraphs for r in p.runs)
        num_paragraphs = len(text_content.paragraphs)
        est_height = max(0.5, min(3.0, num_paragraphs * 0.3 + total_chars / 200))

        # Stack overflow text below the main content area
        left = Inches(0.5)
        top = Inches(2.5 + existing_freeform * (est_height + 0.15))
        width = Inches(12.3)
        height = Inches(est_height)

        # Clamp to stay on slide
        if top + height > Inches(6.8):
            top = Inches(6.8) - height

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, para_content in enumerate(text_content.paragraphs):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()

            # Preserve alignment
            if para_content.alignment:
                try:
                    align_map = {
                        "LEFT (0)": PP_ALIGN.LEFT,
                        "CENTER (1)": PP_ALIGN.CENTER,
                        "RIGHT (2)": PP_ALIGN.RIGHT,
                        "JUSTIFY (3)": PP_ALIGN.JUSTIFY,
                    }
                    for key, val in align_map.items():
                        if key in str(para_content.alignment):
                            para.alignment = val
                            break
                except Exception:
                    pass

            para.level = para_content.level

            for j, run_content in enumerate(para_content.runs):
                run = para.add_run()
                run.text = run_content.text
                run.font.name = "Arial"
                if run_content.font_size_pt:
                    run.font.size = Pt(run_content.font_size_pt)
                if run_content.bold:
                    run.font.bold = True
                if run_content.italic:
                    run.font.italic = True
                if run_content.color_hex:
                    try:
                        run.font.color.rgb = RGBColor.from_string(run_content.color_hex)
                    except (ValueError, AttributeError):
                        pass

    def _place_attribution(self, slide, text_content: TextContent):
        """Place attribution text near the bottom of the slide."""
        from pptx.util import Inches, Pt

        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(5)
        height = Inches(0.4)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        para = tf.paragraphs[0]

        text = text_content.plain_text
        run = para.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)


# ============================================================
# Layout Verifier (Claude Vision)
# ============================================================

class LayoutVerifier:
    """
    Compares original and rebuilt slide images via Claude Vision
    to verify the rebuild quality.
    """

    def __init__(self, api_key: str = None):
        if not HAS_ANTHROPIC:
            raise ImportError("anthropic package required for Vision verification")
        self.client = anthropic.Anthropic(api_key=api_key or os.environ.get("ANTHROPIC_API_KEY"))

    def verify_rebuild(self, original_image_path: str,
                       rebuilt_image_path: str,
                       slide_number: int,
                       layout_name: str) -> dict:
        """
        Compare original and rebuilt slide images.

        Returns dict with:
            - pass: bool
            - score: float (0-1)
            - issues: list of str
            - input_tokens: int
            - output_tokens: int
        """
        # Read both images
        with open(original_image_path, "rb") as f:
            original_b64 = base64.standard_b64encode(f.read()).decode("utf-8")
        with open(rebuilt_image_path, "rb") as f:
            rebuilt_b64 = base64.standard_b64encode(f.read()).decode("utf-8")

        prompt = f"""You are verifying a slide rebuild. The FIRST image is the ORIGINAL slide. The SECOND image is the REBUILT slide using the "{layout_name}" UQ template layout.

Compare them and check:
1. **Content preservation**: Is ALL text content from the original present in the rebuild? (titles, body text, captions)
2. **Image preservation**: Are all images from the original present in the rebuild?
3. **Table preservation**: Are all tables present with correct data?
4. **Layout quality**: Does the rebuilt slide look well-organised? Is text readable? Are elements properly positioned?
5. **Missing content**: Is anything from the original MISSING in the rebuild?

Note: The rebuild will look DIFFERENT from the original (that's the point — it's applying a new template). We expect:
- Different backgrounds, colours, and decorative elements
- Different text positions (aligned to template placeholders)
- Professional UQ branding applied

What we DON'T want:
- Missing text content
- Missing images
- Overlapping elements
- Text cut off or overflowing
- Unreadable text

Respond in this exact JSON format:
```json
{{
    "pass": true,
    "score": 0.9,
    "issues": ["issue 1 if any", "issue 2 if any"],
    "content_preserved": true,
    "images_preserved": true,
    "layout_quality": "good"
}}
```

Respond with ONLY the JSON."""

        response = self.client.messages.create(
            model=VISION_MODEL,
            max_tokens=VERIFY_MAX_TOKENS,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": "image/jpeg",
                            "data": original_b64,
                        },
                    },
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": "image/jpeg",
                            "data": rebuilt_b64,
                        },
                    },
                    {
                        "type": "text",
                        "text": prompt,
                    },
                ],
            }],
        )

        result_text = response.content[0].text.strip()

        # Parse JSON
        json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', result_text, re.DOTALL)
        if json_match:
            result_text = json_match.group(1)
        elif not result_text.startswith("{"):
            start = result_text.find("{")
            end = result_text.rfind("}") + 1
            if start >= 0 and end > start:
                result_text = result_text[start:end]

        try:
            result = json.loads(result_text)
        except json.JSONDecodeError:
            result = {
                "pass": False,
                "score": 0.0,
                "issues": [f"Failed to parse verification response: {result_text[:200]}"],
            }

        result["input_tokens"] = response.usage.input_tokens
        result["output_tokens"] = response.usage.output_tokens
        result["slide_number"] = slide_number

        return result


# ============================================================
# Layout Manager — Full Pipeline Orchestrator
# ============================================================

@dataclass
class SlideResult:
    """Result for a single slide in the pipeline."""
    slide_number: int
    original_layout: str
    recommended_layout: str
    confidence: float
    analysis: dict = field(default_factory=dict)
    rebuild_report: dict = field(default_factory=dict)
    verification: dict = field(default_factory=dict)
    status: str = "pending"  # pending, analysed, rebuilt, verified, failed


class LayoutManager:
    """
    Orchestrates the full layout auto-apply pipeline:
      1. Render source slides to images
      2. Extract content from each slide
      3. Analyse each slide via Claude Vision
      4. Rebuild each slide using the recommended template layout
      5. Verify each rebuild via Claude Vision
    """

    def __init__(self, template_path: str = None, api_key: str = None):
        self.api_key = api_key or os.environ.get("ANTHROPIC_API_KEY")
        self.registry = LayoutRegistry(template_path)
        self.extractor = ContentExtractor()
        self.rebuilder = SlideRebuilder(registry=self.registry)

        # Vision components — created lazily, only when needed
        self._analyser = None
        self._verifier = None

        # Tracking
        self.results: list[SlideResult] = []
        self.total_input_tokens = 0
        self.total_output_tokens = 0

    @property
    def analyser(self):
        if self._analyser is None:
            self._analyser = LayoutAnalyser(api_key=self.api_key, registry=self.registry)
        return self._analyser

    @property
    def verifier(self):
        if self._verifier is None:
            self._verifier = LayoutVerifier(api_key=self.api_key)
        return self._verifier

    def run_pipeline(self, pptx_bytes: bytes,
                     progress_callback=None,
                     skip_verification: bool = False,
                     skip_vision: bool = False,
                     slide_limit: int = None) -> dict:
        """
        Run the full layout pipeline on a PPTX file.

        Args:
            pptx_bytes: Raw bytes of the source PPTX
            progress_callback: Optional callable(step, detail, progress_pct)
            skip_verification: Skip the Vision verification step
            skip_vision: Skip ALL Vision calls (use name-mapping fallback only)
            slide_limit: Only process first N slides (for testing/cost control)

        Returns:
            dict with:
                - output_pptx_bytes: bytes of the rebuilt PPTX
                - results: list of SlideResult
                - summary: dict with totals
                - total_cost_usd: float
        """
        self.results = []
        self.total_input_tokens = 0
        self.total_output_tokens = 0

        with tempfile.TemporaryDirectory() as tmpdir:
            # ── Step 1: Save source and render to images ──
            if progress_callback:
                progress_callback("render", "Rendering slides to images...", 0.05)

            src_path = os.path.join(tmpdir, "source.pptx")
            with open(src_path, "wb") as f:
                f.write(pptx_bytes)

            source_images = []
            if not skip_vision:
                try:
                    render_dir = os.path.join(tmpdir, "source_images")
                    source_images = SlideRenderer.render_deck_to_images(src_path, render_dir)
                except Exception as e:
                    if progress_callback:
                        progress_callback("render", f"Image rendering failed ({e}), using name mapping", 0.08)
                    skip_vision = True

            # ── Step 2: Load source and extract content ──
            if progress_callback:
                progress_callback("extract", "Extracting slide content...", 0.10)

            source_prs = Presentation(io.BytesIO(pptx_bytes))
            slides = list(source_prs.slides)
            num_slides = min(len(slides), slide_limit) if slide_limit else len(slides)

            contents = []
            for i in range(num_slides):
                sc = self.extractor.extract_slide(slides[i], i + 1)
                contents.append(sc)

            # ── Step 3: Create target presentation from template ──
            target_prs = Presentation(self.registry.template_path)
            # Remove all example slides from the template
            # (Keep only the slide layouts, remove actual slides)
            while len(target_prs.slides) > 0:
                rId = target_prs.slides._sldIdLst[0].get(
                    '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
                )
                if rId:
                    target_prs.part.drop_rel(rId)
                target_prs.slides._sldIdLst.remove(target_prs.slides._sldIdLst[0])

            # ── Step 4: Analyse + Rebuild each slide ──
            for i in range(num_slides):
                pct = 0.15 + (0.65 * i / num_slides)
                if progress_callback:
                    progress_callback(
                        "process",
                        f"Processing slide {i+1}/{num_slides}...",
                        pct,
                    )

                sc = contents[i]
                result = SlideResult(
                    slide_number=i + 1,
                    original_layout=sc.original_layout_name,
                    recommended_layout="",
                    confidence=0.0,
                )

                try:
                    # ── 4a: Layout analysis ──
                    if not skip_vision and i < len(source_images):
                        try:
                            analysis = self.analyser.analyse_slide(
                                source_images[i], sc
                            )
                        except Exception as vision_err:
                            # Vision failed — fall back to name mapping + heuristic
                            resolved = self.registry.resolve_name(sc.original_layout_name)
                            if resolved in self.registry.layouts:
                                analysis = {
                                    "recommended_layout": resolved,
                                    "confidence": 0.7 if resolved != sc.original_layout_name else 0.9,
                                    "reasoning": f"Name mapping fallback (Vision error: {vision_err})",
                                    "input_tokens": 0,
                                    "output_tokens": 0,
                                }
                            else:
                                guessed, conf = self.registry.guess_layout_from_content(sc)
                                analysis = {
                                    "recommended_layout": guessed,
                                    "confidence": conf,
                                    "reasoning": f"Content heuristic fallback (Vision error: {vision_err})",
                                    "input_tokens": 0,
                                    "output_tokens": 0,
                                }
                    else:
                        # No Vision — use name mapping, then heuristic fallback
                        resolved = self.registry.resolve_name(sc.original_layout_name)
                        if resolved in self.registry.layouts:
                            analysis = {
                                "recommended_layout": resolved,
                                "confidence": 0.7 if resolved != sc.original_layout_name else 0.9,
                                "reasoning": "Name mapping",
                                "input_tokens": 0,
                                "output_tokens": 0,
                            }
                        else:
                            # Name not in registry — use content heuristic
                            guessed, conf = self.registry.guess_layout_from_content(sc)
                            analysis = {
                                "recommended_layout": guessed,
                                "confidence": conf,
                                "reasoning": f"Content heuristic (original: '{sc.original_layout_name}' not in template)",
                                "input_tokens": 0,
                                "output_tokens": 0,
                            }

                    result.recommended_layout = analysis["recommended_layout"]
                    result.confidence = analysis.get("confidence", 0.5)
                    result.analysis = analysis
                    result.status = "analysed"

                    self.total_input_tokens += analysis.get("input_tokens", 0)
                    self.total_output_tokens += analysis.get("output_tokens", 0)

                    # ── 4b: Rebuild ──
                    slide, rebuild_report = self.rebuilder.rebuild_slide(
                        target_prs, sc, result.recommended_layout
                    )
                    result.rebuild_report = rebuild_report
                    result.status = "rebuilt"

                except Exception as e:
                    result.status = "failed"
                    result.rebuild_report = {"error": str(e)}

                self.results.append(result)

            # ── Step 5: Save rebuilt PPTX ──
            if progress_callback:
                progress_callback("save", "Saving rebuilt presentation...", 0.82)

            output_buf = io.BytesIO()
            target_prs.save(output_buf)
            output_bytes = output_buf.getvalue()

            # ── Step 6: Verification (optional, requires Vision) ──
            if not skip_verification and not skip_vision:
                if progress_callback:
                    progress_callback("verify", "Rendering rebuilt slides for verification...", 0.85)

                rebuild_path = os.path.join(tmpdir, "rebuilt.pptx")
                with open(rebuild_path, "wb") as f:
                    f.write(output_bytes)

                rebuild_render_dir = os.path.join(tmpdir, "rebuilt_images")
                rebuilt_images = SlideRenderer.render_deck_to_images(
                    rebuild_path, rebuild_render_dir
                )

                for i, result in enumerate(self.results):
                    if result.status == "rebuilt" and i < len(rebuilt_images) and i < len(source_images):
                        pct = 0.85 + (0.13 * i / len(self.results))
                        if progress_callback:
                            progress_callback(
                                "verify",
                                f"Verifying slide {i+1}/{len(self.results)}...",
                                pct,
                            )

                        try:
                            verification = self.verifier.verify_rebuild(
                                source_images[i],
                                rebuilt_images[i],
                                result.slide_number,
                                result.recommended_layout,
                            )
                            result.verification = verification
                            result.status = "verified"

                            self.total_input_tokens += verification.get("input_tokens", 0)
                            self.total_output_tokens += verification.get("output_tokens", 0)

                        except Exception as e:
                            result.verification = {"error": str(e)}

            # ── Step 7: Summary ──
            if progress_callback:
                progress_callback("done", "Complete!", 1.0)

            total_cost = (
                self.total_input_tokens * COST_INPUT_PER_TOKEN +
                self.total_output_tokens * COST_OUTPUT_PER_TOKEN
            )

            summary = {
                "total_slides": num_slides,
                "rebuilt": sum(1 for r in self.results if r.status in ("rebuilt", "verified")),
                "verified_pass": sum(1 for r in self.results
                                    if r.verification.get("pass", False)),
                "verified_fail": sum(1 for r in self.results
                                    if r.status == "verified" and not r.verification.get("pass", True)),
                "failed": sum(1 for r in self.results if r.status == "failed"),
                "low_confidence": sum(1 for r in self.results if r.confidence < 0.6),
                "total_input_tokens": self.total_input_tokens,
                "total_output_tokens": self.total_output_tokens,
                "total_cost_usd": round(total_cost, 4),
            }

            return {
                "output_pptx_bytes": output_bytes,
                "results": self.results,
                "summary": summary,
                "total_cost_usd": total_cost,
            }

    def get_results_report(self) -> str:
        """Generate a human-readable results report."""
        lines = ["=" * 60, "LAYOUT AUTO-APPLY RESULTS", "=" * 60, ""]

        for r in self.results:
            status_emoji = {
                "verified": "PASS" if r.verification.get("pass") else "WARN",
                "rebuilt": "BUILT",
                "analysed": "ANALYSED",
                "failed": "FAIL",
                "pending": "PENDING",
            }.get(r.status, "?")

            conf = f"{r.confidence:.0%}" if r.confidence else "?"

            lines.append(
                f"[{status_emoji}] Slide {r.slide_number}: "
                f"{r.original_layout} → {r.recommended_layout} "
                f"(confidence: {conf})"
            )

            if r.rebuild_report.get("warnings"):
                for w in r.rebuild_report["warnings"]:
                    lines.append(f"  ⚠ {w}")

            if r.verification.get("issues"):
                for issue in r.verification["issues"]:
                    lines.append(f"  ! {issue}")

            if r.status == "failed":
                lines.append(f"  ✗ Error: {r.rebuild_report.get('error', 'unknown')}")

        lines.append("")
        lines.append(f"Total: {len(self.results)} slides processed")

        return "\n".join(lines)
