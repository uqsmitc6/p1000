#!/usr/bin/env python3
"""
UQ Brand Compliance Fixer — Step 1
===================================
Takes a messy .pptx submitted by an academic and automatically corrects
brand-level formatting to the February 2026 UQ template standard.

Fixes:
  1. Font normalisation (all → Arial, preserving bold/italic/size)
  2. Text colour correction (smart rules — respects white-on-dark)
  3. Table restyling (UQ header + alternating rows)
  4. Footer standardisation
  5. Heading size normalisation (titles 28–44pt, body 14–20pt)
  6. Bullet style consistency

Usage:
  python brand_fixer.py input.pptx --output fixed.pptx [--report]
"""

import argparse
import copy
import json
import re
import sys
from pathlib import Path
from collections import defaultdict

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ─── UQ Brand Constants ───────────────────────────────────────────────

UQ_FONT = "Arial"

# Colour palette
UQ_PURPLE = RGBColor(0x51, 0x24, 0x7A)
UQ_MAGENTA = RGBColor(0x96, 0x2A, 0x8B)
UQ_DARK = RGBColor(0x2B, 0x1D, 0x37)
UQ_LIGHT = RGBColor(0xD7, 0xD1, 0xCC)
UQ_RED = RGBColor(0xE6, 0x26, 0x45)
UQ_BLUE = RGBColor(0x40, 0x85, 0xC6)
UQ_GOLD = RGBColor(0xFB, 0xB8, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# All approved colours (for checking if a colour is "close enough")
UQ_APPROVED_COLOURS = {
    UQ_PURPLE, UQ_MAGENTA, UQ_DARK, UQ_LIGHT,
    UQ_RED, UQ_BLUE, UQ_GOLD, WHITE, BLACK,
}

# Colours that are "dark" enough to have white text on them
DARK_BACKGROUNDS = {UQ_PURPLE, UQ_MAGENTA, UQ_DARK, BLACK}

# Common near-black colours that should be auto-corrected to UQ_DARK
# These are frequently used dark greys/blacks that are clearly meant to be
# "body text black" — not intentional accent colours
NEAR_BLACK_COLOURS = {
    "#000000", "#111111", "#1A1A1A", "#202020", "#202124",
    "#222222", "#262626", "#2C2C2C", "#303030", "#333333",
    "#363636", "#3C3C3C", "#404040", "#444444", "#464646",
    "#4A4A4A", "#4D4D4D", "#505050", "#515151", "#545454",
}

# Standard heading size range (pt)
TITLE_SIZE_MIN = Pt(28)
TITLE_SIZE_MAX = Pt(44)
BODY_SIZE_MIN = Pt(14)
BODY_SIZE_MAX = Pt(20)

# Theme font reference codes that resolve to Arial in UQ template
THEME_FONT_CODES = {"+mj-lt", "+mn-lt", "+mj-ea", "+mn-ea", "+mj-cs", "+mn-cs"}

# Namespaces for XML operations
NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


# ─── Helper Functions ──────────────────────────────────────────────────

def rgb_to_hex(colour):
    """Convert RGBColor to hex string like #51247A."""
    if colour is None:
        return None
    return f"#{colour[0]:02X}{colour[1]:02X}{colour[2]:02X}"


def hex_to_rgb(hex_str):
    """Convert hex string to RGBColor."""
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def colour_distance(c1, c2):
    """Simple Euclidean distance between two RGBColors."""
    return ((c1[0] - c2[0]) ** 2 + (c1[1] - c2[1]) ** 2 + (c1[2] - c2[2]) ** 2) ** 0.5


def is_dark_colour(colour):
    """Check if a colour is dark (luminance-based)."""
    if colour is None:
        return False
    # Relative luminance formula
    r, g, b = colour[0] / 255, colour[1] / 255, colour[2] / 255
    luminance = 0.2126 * r + 0.7152 * g + 0.0722 * b
    return luminance < 0.4


def is_light_colour(colour):
    """Check if a colour is light."""
    if colour is None:
        return True  # Default assumption
    r, g, b = colour[0] / 255, colour[1] / 255, colour[2] / 255
    luminance = 0.2126 * r + 0.7152 * g + 0.0722 * b
    return luminance > 0.6


def is_approved_colour(colour, tolerance=15):
    """Check if a colour is close to any approved UQ colour."""
    if colour is None:
        return True
    for approved in UQ_APPROVED_COLOURS:
        if colour_distance(colour, approved) < tolerance:
            return True
    return False


def get_shape_background_colour(shape):
    """Try to determine the background colour of a shape or its parent slide."""
    # Check shape fill
    try:
        fill = shape.fill
        if fill.type is not None:
            fg = fill.fore_color
            if fg and fg.type is not None:
                return fg.rgb
    except Exception:
        pass
    return None


def get_slide_background_colour(slide):
    """Try to determine the slide's background colour."""
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            fg = fill.fore_color
            if fg and fg.type is not None:
                return fg.rgb
    except Exception:
        pass
    return None


def is_placeholder_title(shape):
    """Check if a shape is a title placeholder."""
    try:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            # Title types: TITLE (15), CENTER_TITLE (3), SUBTITLE (4)
            from pptx.enum.shapes import PP_PLACEHOLDER
            return ph_type in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
                PP_PLACEHOLDER.SUBTITLE,
            )
    except Exception:
        pass
    return False


def is_placeholder_footer(shape):
    """Check if a shape is a footer, slide number, or date placeholder."""
    try:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            from pptx.enum.shapes import PP_PLACEHOLDER
            return ph_type in (
                PP_PLACEHOLDER.FOOTER,
                PP_PLACEHOLDER.SLIDE_NUMBER,
                PP_PLACEHOLDER.DATE,
            )
    except Exception:
        pass
    return False


def get_run_font_name(run):
    """Get the effective font name for a run, resolving theme references."""
    # Check the run's XML directly for theme font references
    rPr = run._r.find("{http://schemas.openxmlformats.org/drawingml/2006/main}rPr")
    if rPr is not None:
        # Check for latin font
        latin = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}latin")
        if latin is not None:
            typeface = latin.get("typeface", "")
            if typeface in THEME_FONT_CODES:
                return None  # Theme font — resolves to Arial, leave it alone
            if typeface:
                return typeface

    # Fall back to python-pptx's font.name
    font_name = run.font.name
    if font_name in THEME_FONT_CODES:
        return None
    return font_name


def get_run_colour(run):
    """Get the effective text colour of a run, or None if inherited."""
    try:
        if run.font.color and run.font.color.type is not None:
            return run.font.color.rgb
    except Exception:
        pass
    return None


# ─── Fix Functions ─────────────────────────────────────────────────────

class BrandFixer:
    """Applies UQ brand compliance fixes to a presentation."""

    def __init__(self, prs, report=False, footer_text=None):
        self.prs = prs
        self.report = report
        self.footer_text = footer_text  # Optional: set all footers to this text
        self.stats = defaultdict(int)
        self.changes = []  # Detailed change log for report

    def log_change(self, slide_num, category, detail):
        """Record a change for the report."""
        self.stats[category] += 1
        if self.report:
            self.changes.append({
                "slide": slide_num,
                "category": category,
                "detail": detail,
            })

    # ── 1. Font Normalisation ──────────────────────────────────────

    def fix_fonts(self):
        """Replace all non-Arial fonts with Arial, preserving attributes."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self._fix_text_frame_fonts(shape.text_frame, slide_idx)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            self._fix_text_frame_fonts(cell.text_frame, slide_idx)
            # Also check grouped shapes
            self._fix_group_fonts(slide, slide_idx)

    def _fix_group_fonts(self, parent, slide_idx):
        """Recursively fix fonts in grouped shapes."""
        try:
            if hasattr(parent, "shapes"):
                for shape in parent.shapes:
                    if shape.has_text_frame:
                        self._fix_text_frame_fonts(shape.text_frame, slide_idx)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                self._fix_text_frame_fonts(cell.text_frame, slide_idx)
                    self._fix_group_fonts(shape, slide_idx)
        except Exception:
            pass

    def _fix_text_frame_fonts(self, text_frame, slide_idx):
        """Fix all fonts in a text frame."""
        for para in text_frame.paragraphs:
            for run in para.runs:
                font_name = get_run_font_name(run)
                if font_name is not None and font_name != UQ_FONT:
                    old_font = font_name
                    run.font.name = UQ_FONT
                    # Also fix the XML directly to ensure latin typeface is set
                    rPr = run._r.find(
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
                    )
                    if rPr is not None:
                        latin = rPr.find(
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}latin"
                        )
                        if latin is not None:
                            latin.set("typeface", UQ_FONT)
                        # Also fix cs (complex script) and ea (East Asian) if they
                        # have non-theme typefaces
                        for tag in ("cs", "ea"):
                            el = rPr.find(
                                f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}"
                            )
                            if el is not None:
                                tf = el.get("typeface", "")
                                if tf and tf not in THEME_FONT_CODES and tf != UQ_FONT:
                                    el.set("typeface", UQ_FONT)
                    self.log_change(
                        slide_idx,
                        "font",
                        f"Changed '{old_font}' → 'Arial'",
                    )

    # ── 2. Text Colour Correction ──────────────────────────────────

    def fix_colours(self):
        """Smart colour correction — respects context."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            slide_bg = get_slide_background_colour(slide)
            slide_is_dark = is_dark_colour(slide_bg) if slide_bg else False

            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_bg = get_shape_background_colour(shape)
                    is_title = is_placeholder_title(shape)
                    self._fix_text_frame_colours(
                        shape.text_frame, slide_idx, is_title,
                        shape_bg, slide_bg, slide_is_dark,
                    )
                if shape.has_table:
                    self._fix_table_text_colours(shape.table, slide_idx)
                # Recurse into groups
                self._fix_group_colours(shape, slide_idx, slide_bg, slide_is_dark)

    def _fix_group_colours(self, parent, slide_idx, slide_bg, slide_is_dark):
        """Recursively fix colours in grouped shapes."""
        try:
            if hasattr(parent, "shapes"):
                for shape in parent.shapes:
                    if shape.has_text_frame:
                        shape_bg = get_shape_background_colour(shape)
                        is_title = is_placeholder_title(shape)
                        self._fix_text_frame_colours(
                            shape.text_frame, slide_idx, is_title,
                            shape_bg, slide_bg, slide_is_dark,
                        )
                    self._fix_group_colours(shape, slide_idx, slide_bg, slide_is_dark)
        except Exception:
            pass

    def _fix_text_frame_colours(self, text_frame, slide_idx, is_title,
                                 shape_bg, slide_bg, slide_is_dark):
        """Fix text colours in a text frame with smart context awareness."""
        for para in text_frame.paragraphs:
            for run in para.runs:
                colour = get_run_colour(run)
                if colour is None:
                    # Inherited from theme — normally leave it.
                    # BUT for titles, explicitly set to UQ Purple so they
                    # survive layout changes (otherwise they inherit whatever
                    # the new theme defaults to, often black).
                    if is_title and run.text.strip():
                        run.font.color.rgb = UQ_PURPLE
                        self.log_change(
                            slide_idx,
                            "colour",
                            f"Set title colour → #51247A (was theme-inherited)",
                        )
                    continue

                # Already an approved UQ colour? Skip.
                if is_approved_colour(colour):
                    continue

                # Determine the effective background
                effective_bg = shape_bg or slide_bg

                # Is the text white/very light on a dark background? Leave it.
                if is_light_colour(colour) and (
                    is_dark_colour(effective_bg) or slide_is_dark
                ):
                    continue

                # Is this text nearly white (e.g., #FAFAFA) on a light bg?
                # That's probably broken — fix to appropriate colour.

                old_hex = rgb_to_hex(colour)

                # Decide what to fix to
                if is_title:
                    new_colour = UQ_PURPLE
                elif is_dark_colour(effective_bg) or slide_is_dark:
                    new_colour = WHITE
                else:
                    # Standard body text on light background → UQ dark
                    # Check if it's a near-black colour (common dark greys
                    # used as body text) — auto-correct these to UQ_DARK.
                    # Uses both an explicit list of common hex values AND
                    # a luminance check (dark colours with low saturation
                    # are almost certainly "meant to be black" body text).
                    r, g, b = colour[0], colour[1], colour[2]
                    is_near_black = (
                        old_hex.upper() in NEAR_BLACK_COLOURS
                        or colour_distance(colour, BLACK) < 60
                        or (
                            # Low luminance + low saturation = dark grey
                            max(r, g, b) < 100
                            and (max(r, g, b) - min(r, g, b)) < 30
                        )
                    )
                    if is_near_black:
                        new_colour = UQ_DARK
                    else:
                        # It's a non-approved accent colour — leave it and flag
                        # (don't auto-correct potentially intentional accent colours)
                        self.log_change(
                            slide_idx,
                            "colour_flagged",
                            f"Non-UQ colour {old_hex} — review manually",
                        )
                        continue

                run.font.color.rgb = new_colour
                new_hex = rgb_to_hex(new_colour)
                self.log_change(
                    slide_idx,
                    "colour",
                    f"Changed text colour {old_hex} → {new_hex}"
                    + (" (title)" if is_title else " (body)"),
                )

    def _fix_table_text_colours(self, table, slide_idx):
        """Fix text colours inside tables — handled separately since table cells
        have their own background context."""
        for row_idx, row in enumerate(table.rows):
            for cell in row.cells:
                # Determine if this is a header row (first row)
                is_header = row_idx == 0
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        colour = get_run_colour(run)
                        if is_header:
                            # Header text MUST be white (purple background).
                            # Set explicitly even if colour is None (inherited),
                            # because inherited colours won't survive layout changes.
                            if colour is None or (colour != WHITE and not is_light_colour(colour)):
                                old_hex = rgb_to_hex(colour) if colour else "theme"
                                run.font.color.rgb = WHITE
                                self.log_change(
                                    slide_idx,
                                    "colour",
                                    f"Table header text {old_hex} → #FFFFFF",
                                )
                        else:
                            if colour is None:
                                continue  # Body rows — inherited colour is fine
                            # Body rows — dark text
                            if not is_approved_colour(colour):
                                old_hex = rgb_to_hex(colour)
                                run.font.color.rgb = UQ_DARK
                                self.log_change(
                                    slide_idx,
                                    "colour",
                                    f"Table body text {old_hex} → #2B1D37",
                                )

    # ── 3. Table Restyling ─────────────────────────────────────────

    def fix_tables(self):
        """Apply UQ brand colours to tables."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_table:
                    self._fix_table(shape.table, slide_idx)

    def _fix_table(self, table, slide_idx):
        """Restyle a single table to UQ brand standards."""
        changes_made = False

        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                tc_elem = cell._tc

                # Get or create tcPr
                tcPr = tc_elem.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr"
                )
                if tcPr is None:
                    tcPr = etree.SubElement(
                        tc_elem,
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr",
                    )
                    # Move tcPr to be the last child (after txBody)
                    tc_elem.append(tcPr)

                if row_idx == 0:
                    # Header row — UQ Purple background, white text
                    self._set_cell_fill(tcPr, "51247A")
                    changes_made = True
                else:
                    # Alternating rows
                    if row_idx % 2 == 1:
                        self._set_cell_fill(tcPr, "FFFFFF")
                    else:
                        self._set_cell_fill(tcPr, "D7D1CC")
                    changes_made = True

        if changes_made:
            self.log_change(slide_idx, "table", "Restyled table to UQ brand colours")

    def _set_cell_fill(self, tcPr, hex_colour):
        """Set the solid fill of a table cell via XML."""
        # Remove existing fill elements
        for child in list(tcPr):
            tag = etree.QName(child.tag).localname
            if tag in ("solidFill", "noFill", "gradFill", "pattFill"):
                tcPr.remove(child)

        # Add solid fill
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        solidFill = etree.SubElement(tcPr, f"{{{ns}}}solidFill")
        srgbClr = etree.SubElement(solidFill, f"{{{ns}}}srgbClr")
        srgbClr.set("val", hex_colour)

    # ── 4. Footer Standardisation ──────────────────────────────────

    def fix_footers(self):
        """Ensure footer placeholders use consistent formatting.

        If self.footer_text is set, also replaces the text content of footer
        placeholders (excluding slide number and date placeholders) with that text.
        """
        from pptx.enum.shapes import PP_PLACEHOLDER

        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if is_placeholder_footer(shape) and shape.has_text_frame:
                    # Set footer text if configured (only for FOOTER type, not slide numbers/dates)
                    if self.footer_text is not None:
                        try:
                            ph_type = shape.placeholder_format.type
                            if ph_type == PP_PLACEHOLDER.FOOTER:
                                current_text = shape.text_frame.text.strip()
                                if current_text != self.footer_text:
                                    # Clear existing text and set new
                                    for para in shape.text_frame.paragraphs:
                                        for run in para.runs:
                                            run.text = ""
                                    # Set first run of first paragraph
                                    if shape.text_frame.paragraphs:
                                        para = shape.text_frame.paragraphs[0]
                                        if para.runs:
                                            para.runs[0].text = self.footer_text
                                        else:
                                            para.text = self.footer_text
                                    self.log_change(
                                        slide_idx,
                                        "footer",
                                        f"Footer text → '{self.footer_text}'",
                                    )
                        except (ValueError, AttributeError):
                            pass  # Shape isn't a placeholder or doesn't have type

                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            font_name = get_run_font_name(run)
                            if font_name is not None and font_name != UQ_FONT:
                                run.font.name = UQ_FONT
                                self.log_change(
                                    slide_idx,
                                    "footer",
                                    f"Footer font → Arial",
                                )
                            # Footers should be small and in UQ_DARK
                            if run.font.size and run.font.size > Pt(12):
                                run.font.size = Pt(10)
                                self.log_change(
                                    slide_idx,
                                    "footer",
                                    "Footer font size → 10pt",
                                )

    # ── 5. Heading Size Normalisation ──────────────────────────────

    def fix_heading_sizes(self):
        """Ensure title sizes are within 28-44pt range."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if is_placeholder_title(shape) and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            size = run.font.size
                            if size is None:
                                continue
                            if size < TITLE_SIZE_MIN:
                                old_pt = size / 12700
                                run.font.size = TITLE_SIZE_MIN
                                self.log_change(
                                    slide_idx,
                                    "heading_size",
                                    f"Title size {old_pt:.0f}pt → 28pt (was below minimum)",
                                )
                            elif size > TITLE_SIZE_MAX:
                                old_pt = size / 12700
                                run.font.size = TITLE_SIZE_MAX
                                self.log_change(
                                    slide_idx,
                                    "heading_size",
                                    f"Title size {old_pt:.0f}pt → 44pt (was above maximum)",
                                )

    # ── 6. Body Text Size Check ────────────────────────────────────

    BODY_SIZE_MIN = Pt(12)
    BODY_SIZE_MAX = Pt(24)

    # Patterns that indicate text is intentionally small (attributions, captions)
    _SMALL_TEXT_EXCEPTIONS = re.compile(
        r"(source:|image\s+(licensed|source)|adapted\s+from|"
        r"photo\s+(?:by|credit)|cc\s+by|creative\s+commons|"
        r"public\s+domain|wikimedia|adobe\s+stock|shutterstock|"
        r"©|\d{4}\s+\w+\s+(pty|ltd|inc|corp)|all\s+rights\s+reserved)",
        re.IGNORECASE,
    )

    def flag_body_text_sizes(self):
        """Flag body text that falls outside the 14–24pt range.

        This is a flag-only check — no auto-correction, since small or large
        body text may be intentional (captions, callout boxes, etc.).
        Excludes titles, footers, and attribution/caption text.
        """
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                # Skip titles and footers (handled separately)
                if is_placeholder_title(shape) or is_placeholder_footer(shape):
                    continue
                if not shape.has_text_frame:
                    continue

                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        size = run.font.size
                        if size is None:
                            continue

                        text = run.text.strip()
                        if not text:
                            continue

                        # Skip intentionally small text (attributions, captions)
                        if self._SMALL_TEXT_EXCEPTIONS.search(text):
                            continue

                        pt_val = round(size / 12700, 1)

                        if size < self.BODY_SIZE_MIN:
                            self.log_change(
                                slide_idx,
                                "body_size_flagged",
                                f"Body text at {pt_val}pt (below 14pt minimum) — "
                                f"'{text[:50]}{'...' if len(text) > 50 else ''}'",
                            )
                        elif size > self.BODY_SIZE_MAX:
                            self.log_change(
                                slide_idx,
                                "body_size_flagged",
                                f"Body text at {pt_val}pt (above 24pt maximum) — "
                                f"'{text[:50]}{'...' if len(text) > 50 else ''}'",
                            )

    # ── 7. Bullet Style Consistency ────────────────────────────────

    def fix_bullets(self):
        """Normalise bullet characters for consistency.

        Strategy: Don't force bullets where none exist. Where bullets are
        already used, normalise to a consistent character (en dash –).
        """
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self._fix_text_frame_bullets(shape.text_frame, slide_idx)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            self._fix_text_frame_bullets(cell.text_frame, slide_idx)

    def _fix_text_frame_bullets(self, text_frame, slide_idx):
        """Normalise bullet characters in a text frame."""
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        for para in text_frame.paragraphs:
            pPr = para._p.find(f"{{{ns}}}pPr")
            if pPr is not None:
                buChar = pPr.find(f"{{{ns}}}buChar")
                if buChar is not None:
                    current_char = buChar.get("char", "")
                    # Normalise various bullet chars to en dash
                    # but leave common standard ones alone
                    standard_chars = {"–", "—", "•", "-", "▪", "►", "‣", "→"}
                    if current_char and current_char not in standard_chars:
                        buChar.set("char", "–")
                        self.log_change(
                            slide_idx,
                            "bullet",
                            f"Normalised bullet '{current_char}' → '–'",
                        )

    # ── Run All Fixes ──────────────────────────────────────────────

    def fix_all(self):
        """Run all brand compliance fixes in order."""
        print("  [1/7] Fixing fonts...")
        self.fix_fonts()
        print("  [2/7] Fixing text colours...")
        self.fix_colours()
        print("  [3/7] Restyling tables...")
        self.fix_tables()
        print("  [4/7] Standardising footers...")
        self.fix_footers()
        print("  [5/7] Normalising heading sizes...")
        self.fix_heading_sizes()
        print("  [6/7] Checking body text sizes...")
        self.flag_body_text_sizes()
        print("  [7/7] Fixing bullet styles...")
        self.fix_bullets()

    def print_summary(self):
        """Print a summary of all changes made."""
        total = sum(self.stats.values())
        print(f"\n{'='*60}")
        print(f"  UQ Brand Compliance Fixer — Summary")
        print(f"{'='*60}")
        print(f"  Total changes: {total}")
        print()
        categories = {
            "font": "Font fixes",
            "colour": "Colour corrections",
            "colour_flagged": "Colours flagged for review",
            "table": "Tables restyled",
            "footer": "Footer fixes",
            "heading_size": "Heading size adjustments",
            "body_size_flagged": "Body text sizes flagged",
            "bullet": "Bullet normalisations",
        }
        for key, label in categories.items():
            count = self.stats.get(key, 0)
            if count > 0:
                print(f"  {label:.<40} {count}")
        print(f"{'='*60}")
        return total

    def generate_report(self):
        """Generate a JSON-serialisable report of all changes."""
        return {
            "summary": dict(self.stats),
            "total_changes": sum(self.stats.values()),
            "changes": self.changes,
        }


# ─── Main ──────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="UQ Brand Compliance Fixer — automatically correct brand formatting"
    )
    parser.add_argument("input", help="Path to input .pptx file")
    parser.add_argument(
        "--output", "-o",
        help="Path to output .pptx file (default: input_FIXED.pptx)",
    )
    parser.add_argument(
        "--report",
        action="store_true",
        help="Generate a detailed JSON change report",
    )
    parser.add_argument(
        "--footer-text",
        help="Set all footer placeholders to this text (e.g. 'UQ Business School')",
    )
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file not found: {input_path}")
        sys.exit(1)

    if args.output:
        output_path = Path(args.output)
    else:
        output_path = input_path.with_name(f"{input_path.stem}_FIXED{input_path.suffix}")

    print(f"\nUQ Brand Compliance Fixer")
    print(f"{'─'*40}")
    print(f"  Input:  {input_path}")
    print(f"  Output: {output_path}")
    print()

    # Load presentation
    print("Loading presentation...")
    prs = Presentation(str(input_path))
    print(f"  {len(prs.slides)} slides loaded")
    print()

    # Run fixes
    print("Applying brand fixes...")
    fixer = BrandFixer(prs, report=args.report, footer_text=args.footer_text)
    fixer.fix_all()
    total = fixer.print_summary()

    # Save
    print(f"\nSaving to {output_path}...")
    prs.save(str(output_path))
    print("Done!")

    # Generate report if requested
    if args.report:
        report_path = output_path.with_suffix(".json")
        report_data = fixer.generate_report()
        with open(report_path, "w") as f:
            json.dump(report_data, f, indent=2)
        print(f"Report saved to {report_path}")

    return total


if __name__ == "__main__":
    main()
