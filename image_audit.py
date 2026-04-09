#!/usr/bin/env python3
"""
UQ Image Audit Report — Step 2
================================
Extracts every image from a .pptx deck, classifies each using Claude's
vision API, and produces a structured HTML compliance report.

Requires: ANTHROPIC_API_KEY environment variable

Usage:
  python image_audit.py input.pptx --output audit_report.html
  python image_audit.py input.pptx --output audit_report.html --limit 10  # first 10 images only
"""

import argparse
import base64
import hashlib
import io
import json
import os
import sys
import time
from collections import defaultdict
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from PIL import Image

try:
    import anthropic
except ImportError:
    print("Error: anthropic package required. Install with: pip install anthropic")
    sys.exit(1)


# ─── Configuration ─────────────────────────────────────────────────────

CLAUDE_MODEL = "claude-sonnet-4-20250514"
MAX_IMAGE_DIMENSION = 1024  # Resize large images before sending to API

# Pricing per million tokens (USD) — Sonnet 4, as of 2025-05
# Update these if the model or pricing changes
COST_PER_M_INPUT = 3.00    # $3 per 1M input tokens
COST_PER_M_OUTPUT = 15.00  # $15 per 1M output tokens
SUPPORTED_IMAGE_TYPES = {"image/png", "image/jpeg", "image/gif", "image/webp", "image/bmp", "image/tiff"}

# Minimum image size to bother classifying (skip tiny icons/spacers)
MIN_IMAGE_WIDTH = 50
MIN_IMAGE_HEIGHT = 50

CLASSIFICATION_PROMPT = """You are a copyright compliance analyst for the University of Queensland Business School.
You are reviewing images extracted from executive education PowerPoint slide decks.

IMPORTANT LEGAL CONTEXT: These decks are used in COMMERCIAL executive education programs
that fall OUTSIDE Australia's statutory education licence. Every uncleared copyrighted image
is genuine legal exposure. Err on the side of flagging — false positives are acceptable,
false negatives are dangerous.

LICENSING CONTEXT — READ CAREFULLY:
UQ holds licences for the following stock image services:
- Adobe Stock (institutional licence)
- Shutterstock (institutional licence)
- Microsoft 365 stock images (included with institutional M365 subscription)

Images from these providers are LEGALLY COVERED for commercial use at UQ, provided they are
properly attributed. If an image looks like a stock photo AND the nearby slide text contains
an attribution (e.g. "Source: Adobe Stock 123456789", "Image licensed through Shutterstock: 456",
"Images: Microsoft Stock"), treat it as LICENSED — not as a copyright risk.

Analyse this image and provide a JSON response with EXACTLY these fields:

{
  "image_type": "<one of: STOCK_PHOTO, SCREENSHOT_PUBLISHED, SCREENSHOT_TOOL, DIAGRAM_ORIGINAL, DIAGRAM_PUBLISHED, ICON, AI_GENERATED, UQ_BRANDED, DECORATIVE_PHOTO, DATA_CHART, OTHER>",
  "risk_level": "<one of: CRITICAL, HIGH, MEDIUM, LOW, CLEAR>",
  "confidence": "<one of: HIGH, MEDIUM, LOW>",
  "reasoning": "<2-3 sentences explaining your classification>",
  "watermark_text": "<any visible watermark text, or null>",
  "copyright_notice": "<any visible copyright notice, or null>",
  "brand_visible": "<any publication/brand name visible, or null>",
  "attribution_found": "<the attribution text found on the slide, if any, or null>",
  "is_decorative": <true if purely decorative, false if conveys specific content>,
  "content_description": "<brief description of what the image shows — useful for finding a replacement if needed>",
  "alt_text": "<concise, accessible alt text for screen readers (max 125 characters). Describe what the image shows, not what it is. For decorative images, use empty string. For diagrams/charts, summarise the key data or relationship shown.>",
  "recommended_action": "<one of: REPLACE_IMMEDIATELY, VERIFY_LICENCE, ADD_ATTRIBUTION, CHECK_SOURCE, REVIEW_MANUALLY, NO_ACTION>"
}

CLASSIFICATION GUIDE:
- STOCK_PHOTO: Professional photograph with clean lighting, generic/posed subjects
- SCREENSHOT_PUBLISHED: Screenshot of a published infographic, diagram, book page, or commercial content
- SCREENSHOT_TOOL: Screenshot from a software tool, website, or online assessment
- DIAGRAM_ORIGINAL: Simple diagram likely created by the author (basic shapes, text boxes, PowerPoint SmartArt)
- DIAGRAM_PUBLISHED: Complex/professionally designed diagram from a published source
- ICON: Simple icon or symbol (usually fine)
- AI_GENERATED: Image with AI generation characteristics
- UQ_BRANDED: Contains UQ logo, UQ colours, or institutional branding
- DECORATIVE_PHOTO: Atmospheric/mood photo for visual appeal
- DATA_CHART: Chart, graph, or data visualisation
- OTHER: Anything else

PRIORITY GUIDE — what matters most for copyright compliance:
Published diagrams, figures, and infographics from academic or commercial sources (DIAGRAM_PUBLISHED,
SCREENSHOT_PUBLISHED) are HIGHER priority than decorative/atmospheric photos. A copyrighted diagram
reproduced without permission is a clearer infringement than a generic stock photo used decoratively.
Prioritise flagging published figures/diagrams over decorative images.

GOVERNMENT & CORPORATE WEBSITE SCREENSHOTS:
Screenshots from government (.gov, .gov.au) or corporate websites are generally acceptable for
educational/commercial use PROVIDED they are acknowledged (attributed). They are typically published
for public information. Classify as LOW risk if acknowledged, MEDIUM if not acknowledged.
Do NOT classify government/corporate website screenshots as HIGH or CRITICAL unless they contain
content that is clearly behind a paywall or subscription.

RISK GUIDE:
- CRITICAL: Visible watermark, copyright notice, or clearly identifiable commercial publication branding.
  Exception: Adobe Stock / Shutterstock / Microsoft Stock watermarks on images that also have an
  on-slide attribution are NOT critical — they are licensed. Classify as LOW.
- HIGH: Published diagrams, figures, or infographics from academic/commercial sources with no
  attribution or licence evidence. Also: content clearly from behind a paywall or subscription.
  If the nearby text DOES contain an attribution mentioning Adobe Stock, Shutterstock, or
  Microsoft Stock, downgrade to LOW.
- MEDIUM: Professional stock photo of uncertain origin (no attribution found), or unacknowledged
  screenshots from corporate/government websites. If attribution IS present for a licensed
  provider, downgrade to LOW.
- LOW: Licensed stock photo with proper attribution, acknowledged government/corporate screenshots,
  appears original, AI-generated, UQ-branded, or a simple icon.
- CLEAR: Clearly UQ institutional asset, simple shapes, or obviously author-created.

RECOMMENDED ACTION GUIDE:
- REPLACE_IMMEDIATELY: Watermarked, clearly copyrighted, and NOT from a UQ-licensed provider
- VERIFY_LICENCE: Looks like stock but no attribution found — may be licensed, needs checking
- ADD_ATTRIBUTION: Appears to be a licensed stock photo but the slide has no attribution text
- CHECK_SOURCE: Uncertain origin — could be original or from a published source
- REVIEW_MANUALLY: Complex case that needs human judgement
- NO_ACTION: Licensed and attributed, UQ-branded, original, or clearly fine

Respond with ONLY the JSON object, no other text."""

CONTEXT_PROMPT_TEMPLATE = """
ADDITIONAL CONTEXT from the slide:
- Slide title: {slide_title}
- Nearby text: {nearby_text}
- Image alt text: {alt_text}
- Shape name: {shape_name}
- Slide notes: {slide_notes}
- Detected attribution text: {detected_attribution}

IMPORTANT: Check the "Nearby text" and "Detected attribution text" fields carefully for
attribution patterns like "Source: Adobe Stock 123456789", "Image licensed through Shutterstock",
"Images: Microsoft Stock", "Source: Wikimedia Commons", "Source: Public domain", etc.
If you find such an attribution, report it in the "attribution_found" field of your response
and factor it into your risk assessment (licensed + attributed = LOW risk)."""


# ─── Image Extraction ──────────────────────────────────────────────────

def extract_images(pptx_path, output_dir=None, limit=None):
    """Extract all images from a PPTX file with metadata."""
    prs = Presentation(str(pptx_path))
    images = []
    image_hashes = set()  # Deduplicate identical images

    for slide_idx, slide in enumerate(prs.slides, 1):
        # Get slide-level context
        slide_title = ""
        slide_notes = ""
        slide_texts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_texts.append(text)
                try:
                    if shape.is_placeholder:
                        from pptx.enum.shapes import PP_PLACEHOLDER
                        if shape.placeholder_format.type in (
                            PP_PLACEHOLDER.TITLE,
                            PP_PLACEHOLDER.CENTER_TITLE,
                        ):
                            slide_title = text
                except Exception:
                    pass

        try:
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                slide_notes = notes_tf.text.strip()
        except Exception:
            pass

        # Extract images
        img_count = 0
        for shape in slide.shapes:
            image_info = _extract_image_from_shape(
                shape, slide_idx, img_count,
                slide_title, slide_texts, slide_notes,
                image_hashes, output_dir,
            )
            if image_info:
                images.append(image_info)
                img_count += 1

            # Check grouped shapes
            try:
                if hasattr(shape, "shapes"):
                    for child in shape.shapes:
                        image_info = _extract_image_from_shape(
                            child, slide_idx, img_count,
                            slide_title, slide_texts, slide_notes,
                            image_hashes, output_dir,
                        )
                        if image_info:
                            images.append(image_info)
                            img_count += 1
            except Exception:
                pass

        if limit and len(images) >= limit:
            images = images[:limit]
            break

    return images


def _extract_image_from_shape(shape, slide_idx, img_count,
                               slide_title, slide_texts, slide_notes,
                               image_hashes, output_dir):
    """Extract image data from a single shape, if it contains an image."""
    try:
        if not hasattr(shape, "image"):
            return None

        image = shape.image
        content_type = image.content_type

        # Skip non-image content types
        if content_type not in SUPPORTED_IMAGE_TYPES:
            # Also handle WMF/EMF which python-pptx might report
            if "wmf" in str(content_type).lower() or "emf" in str(content_type).lower():
                return None
            # Try anyway for edge cases
            pass

        image_bytes = image.blob
        image_hash = hashlib.sha256(image_bytes).hexdigest()[:16]

        # Deduplicate
        if image_hash in image_hashes:
            return None
        image_hashes.add(image_hash)

        # Check dimensions — skip tiny images
        try:
            img = Image.open(io.BytesIO(image_bytes))
            width, height = img.size
            if width < MIN_IMAGE_WIDTH or height < MIN_IMAGE_HEIGHT:
                return None
        except Exception:
            width, height = 0, 0

        # Get shape metadata
        shape_name = getattr(shape, "name", "")
        alt_text = ""
        try:
            # Alt text is in the shape's XML
            desc_elem = shape._element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetml}cNvPr"
            )
            if desc_elem is None:
                # Try the presentation namespace
                from lxml import etree
                for elem in shape._element.iter():
                    if "cNvPr" in elem.tag:
                        alt_text = elem.get("descr", "")
                        break
        except Exception:
            pass

        # Get hyperlink if present
        hyperlink = None
        try:
            if hasattr(shape, "click_action") and shape.click_action:
                hyperlink = getattr(shape.click_action, "hyperlink", None)
                if hyperlink:
                    hyperlink = str(hyperlink.address) if hyperlink.address else None
        except Exception:
            pass

        # Determine file extension
        ext_map = {
            "image/png": "png",
            "image/jpeg": "jpg",
            "image/gif": "gif",
            "image/webp": "webp",
            "image/bmp": "bmp",
            "image/tiff": "tiff",
        }
        ext = ext_map.get(content_type, "png")
        filename = f"slide{slide_idx:03d}_image{img_count:02d}.{ext}"

        # Save to disk if output_dir specified
        if output_dir:
            out_path = Path(output_dir) / filename
            out_path.write_bytes(image_bytes)

        # Build nearby text context (first 500 chars)
        nearby_text = " | ".join(slide_texts)[:500]

        # Check EXIF metadata for copyright info
        exif_copyright = None
        try:
            img = Image.open(io.BytesIO(image_bytes))
            exif_data = img._getexif()
            if exif_data:
                # Tag 33432 = Copyright
                exif_copyright = exif_data.get(33432, None)
        except Exception:
            pass

        return {
            "slide_number": slide_idx,
            "image_index": img_count,
            "filename": filename,
            "content_type": content_type,
            "width": width,
            "height": height,
            "hash": image_hash,
            "shape_name": shape_name,
            "alt_text": alt_text,
            "hyperlink": hyperlink,
            "slide_title": slide_title,
            "nearby_text": nearby_text,
            "slide_notes": slide_notes,
            "exif_copyright": exif_copyright,
            "image_bytes": image_bytes,  # Keep in memory for API call
        }

    except (AttributeError, KeyError):
        # Shape doesn't have an image
        return None
    except Exception as e:
        print(f"  Warning: Error extracting image from slide {slide_idx}: {e}")
        return None


# ─── Attribution Detection ────────────────────────────────────────────

import re

# Patterns that indicate an image attribution is present on the slide
_ATTRIBUTION_PATTERNS = [
    re.compile(r"Source:\s*Adobe\s*Stock\s*\d+", re.IGNORECASE),
    re.compile(r"Image\s+licensed\s+through\s+Adobe\s*Stock[:\s]*\d*", re.IGNORECASE),
    re.compile(r"Source:\s*Shutterstock\s*\d+", re.IGNORECASE),
    re.compile(r"Image\s+licensed\s+through\s+Shutterstock[:\s]*\d*", re.IGNORECASE),
    re.compile(r"Images?:\s*Microsoft\s+Stock", re.IGNORECASE),
    re.compile(r"Source:\s*Microsoft\s+Stock", re.IGNORECASE),
    re.compile(r"Source:.*Wikimedia\s+Commons", re.IGNORECASE),
    re.compile(r"Source:\s*Public\s+domain", re.IGNORECASE),
    re.compile(r"Source:.*CC\s+BY", re.IGNORECASE),
    re.compile(r"Source:.*Creative\s+Commons", re.IGNORECASE),
    re.compile(r"Image\s+source:", re.IGNORECASE),
    re.compile(r"Photo\s+(?:by|credit|courtesy)", re.IGNORECASE),
    re.compile(r"Source:.*Flickr", re.IGNORECASE),
    re.compile(r"Source:.*Unsplash", re.IGNORECASE),
    re.compile(r"Source:.*Pexels", re.IGNORECASE),
    re.compile(r"Source:.*Getty", re.IGNORECASE),
    re.compile(r"Source:.*iStock", re.IGNORECASE),
]


def _detect_attribution(image_info):
    """Scan slide text and notes for image attribution patterns.

    Returns the matched attribution text if found, or None.
    """
    text_sources = [
        image_info.get("nearby_text", ""),
        image_info.get("slide_notes", ""),
        image_info.get("alt_text", ""),
    ]
    combined = " ".join(t for t in text_sources if t)

    for pattern in _ATTRIBUTION_PATTERNS:
        m = pattern.search(combined)
        if m:
            # Return the matched text plus some surrounding context
            start = max(0, m.start() - 10)
            end = min(len(combined), m.end() + 30)
            return combined[start:end].strip()

    return None


# ─── AI Classification ─────────────────────────────────────────────────

def classify_image(client, image_info):
    """Send an image to Claude's vision API for classification."""
    image_bytes = image_info["image_bytes"]
    content_type = image_info["content_type"]

    # Ensure content type is supported by the API
    supported_api_types = {"image/png", "image/jpeg", "image/gif", "image/webp"}
    if content_type not in supported_api_types:
        # Convert to PNG
        try:
            img = Image.open(io.BytesIO(image_bytes))
            buf = io.BytesIO()
            img.convert("RGB").save(buf, format="PNG")
            image_bytes = buf.getvalue()
            content_type = "image/png"
        except Exception as e:
            return {"error": f"Could not convert image: {e}"}

    # Resize if too large (saves API costs and avoids timeouts)
    try:
        img = Image.open(io.BytesIO(image_bytes))
        w, h = img.size
        if max(w, h) > MAX_IMAGE_DIMENSION:
            ratio = MAX_IMAGE_DIMENSION / max(w, h)
            new_size = (int(w * ratio), int(h * ratio))
            img = img.resize(new_size, Image.LANCZOS)
            buf = io.BytesIO()
            fmt = "PNG" if content_type == "image/png" else "JPEG"
            if img.mode == "RGBA" and fmt == "JPEG":
                img = img.convert("RGB")
            img.save(buf, format=fmt)
            image_bytes = buf.getvalue()
            if fmt == "JPEG":
                content_type = "image/jpeg"
    except Exception:
        pass

    # Build the base64 data
    b64_data = base64.b64encode(image_bytes).decode("utf-8")

    # Detect attribution text from slide context
    detected_attribution = _detect_attribution(image_info)

    # Build context string
    context = CONTEXT_PROMPT_TEMPLATE.format(
        slide_title=image_info.get("slide_title", "N/A"),
        nearby_text=image_info.get("nearby_text", "N/A")[:300],
        alt_text=image_info.get("alt_text", "N/A"),
        shape_name=image_info.get("shape_name", "N/A"),
        slide_notes=image_info.get("slide_notes", "N/A")[:300],
        detected_attribution=detected_attribution or "None found",
    )

    full_prompt = CLASSIFICATION_PROMPT + context

    try:
        response = client.messages.create(
            model=CLAUDE_MODEL,
            max_tokens=1024,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": content_type,
                                "data": b64_data,
                            },
                        },
                        {
                            "type": "text",
                            "text": full_prompt,
                        },
                    ],
                }
            ],
        )

        # Parse the JSON response
        response_text = response.content[0].text.strip()

        # Handle markdown code blocks
        if response_text.startswith("```"):
            lines = response_text.split("\n")
            response_text = "\n".join(lines[1:-1])

        result = json.loads(response_text)
        result["api_tokens_used"] = {
            "input": response.usage.input_tokens,
            "output": response.usage.output_tokens,
        }
        return result

    except json.JSONDecodeError as e:
        return {"error": f"Failed to parse API response as JSON: {e}", "raw_response": response_text}
    except Exception as e:
        return {"error": f"API call failed: {e}"}


# ─── HTML Report Generation ───────────────────────────────────────────

def generate_html_report(images, classifications, pptx_name, output_path, extracted_dir, image_bytes_map=None):
    """Generate a styled HTML compliance report.

    Args:
        images: List of image metadata dicts
        classifications: List of classification result dicts
        pptx_name: Name of the source PPTX file
        output_path: Path to write the HTML report
        extracted_dir: Directory containing extracted images (for file-based src)
        image_bytes_map: Optional dict of {filename: bytes} for self-contained
                        base64-embedded images. When provided, images are embedded
                        directly in the HTML as data URIs instead of file references.
    """
    now = datetime.now().strftime("%Y-%m-%d %H:%M")

    # Calculate summary stats
    risk_counts = defaultdict(int)
    type_counts = defaultdict(int)
    total_images = len(images)
    errors = 0

    for cls in classifications:
        if "error" in cls:
            errors += 1
        else:
            risk_counts[cls.get("risk_level", "UNKNOWN")] += 1
            type_counts[cls.get("image_type", "UNKNOWN")] += 1

    # Calculate total API tokens and cost
    total_input_tokens = sum(
        c.get("api_tokens_used", {}).get("input", 0) for c in classifications
    )
    total_output_tokens = sum(
        c.get("api_tokens_used", {}).get("output", 0) for c in classifications
    )
    cost_input = (total_input_tokens / 1_000_000) * COST_PER_M_INPUT
    cost_output = (total_output_tokens / 1_000_000) * COST_PER_M_OUTPUT
    total_cost_usd = cost_input + cost_output

    # Risk colour mapping
    risk_colours = {
        "CRITICAL": "#DC2626",
        "HIGH": "#EA580C",
        "MEDIUM": "#D97706",
        "LOW": "#16A34A",
        "CLEAR": "#059669",
    }

    risk_bg_colours = {
        "CRITICAL": "#FEF2F2",
        "HIGH": "#FFF7ED",
        "MEDIUM": "#FFFBEB",
        "LOW": "#F0FDF4",
        "CLEAR": "#ECFDF5",
    }

    # Build image cards HTML
    cards_html = []

    # Sort by risk level (CRITICAL first)
    risk_order = {"CRITICAL": 0, "HIGH": 1, "MEDIUM": 2, "LOW": 3, "CLEAR": 4, "UNKNOWN": 5}
    paired = list(zip(images, classifications))
    paired.sort(key=lambda x: risk_order.get(x[1].get("risk_level", "UNKNOWN"), 5))

    for img_info, cls in paired:
        if "error" in cls:
            risk_level = "ERROR"
            risk_colour = "#6B7280"
            risk_bg = "#F3F4F6"
        else:
            risk_level = cls.get("risk_level", "UNKNOWN")
            risk_colour = risk_colours.get(risk_level, "#6B7280")
            risk_bg = risk_bg_colours.get(risk_level, "#F3F4F6")

        # Build image src — use base64 data URI if image bytes available,
        # otherwise fall back to relative file path
        img_src = ""
        filename = img_info.get('filename', '')
        if image_bytes_map and filename in image_bytes_map:
            img_bytes = image_bytes_map[filename]
            content_type = img_info.get('content_type', 'image/png')
            b64 = base64.b64encode(img_bytes).decode('ascii')
            img_src = f"data:{content_type};base64,{b64}"
        elif extracted_dir:
            img_src = f"{extracted_dir}/{filename}"

        flags = []
        if cls.get("watermark_text"):
            flags.append(f"<span class='flag flag-critical'>Watermark: {cls['watermark_text']}</span>")
        if cls.get("copyright_notice"):
            flags.append(f"<span class='flag flag-critical'>Copyright: {cls['copyright_notice']}</span>")
        if cls.get("brand_visible"):
            flags.append(f"<span class='flag flag-high'>Brand: {cls['brand_visible']}</span>")

        card = f"""
        <div class="card" style="border-left: 4px solid {risk_colour}; background: {risk_bg};">
            <div class="card-header">
                <div class="card-meta">
                    <span class="slide-num">Slide {img_info['slide_number']}</span>
                    <span class="risk-badge" style="background: {risk_colour};">{risk_level}</span>
                    <span class="type-badge">{cls.get('image_type', 'ERROR')}</span>
                </div>
                <div class="card-action">{cls.get('recommended_action', 'N/A')}</div>
            </div>
            <div class="card-body">
                <div class="card-image">
                    <img src="{img_src}" alt="Slide {img_info['slide_number']} image"
                         onerror="this.style.display='none'"/>
                </div>
                <div class="card-details">
                    <p class="reasoning">{cls.get('reasoning', cls.get('error', 'Classification failed'))}</p>
                    <p class="description"><strong>Content:</strong> {cls.get('content_description', 'N/A')}</p>
                    {f'<p class="description"><strong>Suggested alt text:</strong> <em>{cls.get("alt_text", "")}</em></p>' if cls.get('alt_text') else ''}
                    {''.join(flags)}
                    <div class="card-context">
                        <small>
                            <strong>Slide title:</strong> {img_info.get('slide_title', 'N/A')}<br/>
                            <strong>Shape:</strong> {img_info.get('shape_name', 'N/A')} &nbsp;|&nbsp;
                            <strong>Size:</strong> {img_info.get('width', '?')}×{img_info.get('height', '?')}px &nbsp;|&nbsp;
                            <strong>Confidence:</strong> {cls.get('confidence', 'N/A')}
                        </small>
                    </div>
                </div>
            </div>
        </div>
        """
        cards_html.append(card)

    # Assemble full HTML
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Image Audit Report — {pptx_name}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: Arial, sans-serif;
            background: #F8F9FA;
            color: #2B1D37;
            line-height: 1.5;
            padding: 2rem;
        }}
        .container {{ max-width: 1100px; margin: 0 auto; }}

        /* Header */
        .header {{
            background: #51247A;
            color: white;
            padding: 2rem;
            border-radius: 8px;
            margin-bottom: 2rem;
        }}
        .header h1 {{ font-size: 1.8rem; margin-bottom: 0.5rem; }}
        .header .subtitle {{ opacity: 0.9; font-size: 0.95rem; }}

        /* Summary */
        .summary {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
            gap: 1rem;
            margin-bottom: 2rem;
        }}
        .summary-card {{
            background: white;
            border-radius: 8px;
            padding: 1.2rem;
            text-align: center;
            box-shadow: 0 1px 3px rgba(0,0,0,0.1);
        }}
        .summary-card .number {{ font-size: 2rem; font-weight: bold; }}
        .summary-card .label {{ font-size: 0.85rem; color: #666; margin-top: 0.3rem; }}

        /* Risk bar */
        .risk-bar {{
            display: flex;
            height: 24px;
            border-radius: 12px;
            overflow: hidden;
            margin-bottom: 2rem;
            box-shadow: 0 1px 3px rgba(0,0,0,0.1);
        }}
        .risk-bar div {{ display: flex; align-items: center; justify-content: center;
                         font-size: 0.75rem; color: white; font-weight: bold; }}

        /* Section headers */
        .section-header {{
            font-size: 1.3rem;
            font-weight: bold;
            color: #51247A;
            margin: 2rem 0 1rem;
            padding-bottom: 0.5rem;
            border-bottom: 2px solid #D7D1CC;
        }}

        /* Cards */
        .card {{
            background: white;
            border-radius: 8px;
            margin-bottom: 1rem;
            padding: 1.2rem;
            box-shadow: 0 1px 3px rgba(0,0,0,0.1);
        }}
        .card-header {{
            display: flex;
            justify-content: space-between;
            align-items: center;
            margin-bottom: 0.8rem;
        }}
        .card-meta {{ display: flex; align-items: center; gap: 0.5rem; }}
        .slide-num {{ font-weight: bold; font-size: 0.9rem; }}
        .risk-badge {{
            color: white;
            padding: 0.15rem 0.6rem;
            border-radius: 4px;
            font-size: 0.75rem;
            font-weight: bold;
        }}
        .type-badge {{
            background: #E5E7EB;
            padding: 0.15rem 0.6rem;
            border-radius: 4px;
            font-size: 0.75rem;
        }}
        .card-action {{
            font-size: 0.8rem;
            font-weight: bold;
            color: #51247A;
        }}
        .card-body {{
            display: flex;
            gap: 1.2rem;
        }}
        .card-image {{
            flex-shrink: 0;
            width: 180px;
            height: 120px;
            overflow: hidden;
            border-radius: 4px;
            background: #F3F4F6;
            display: flex;
            align-items: center;
            justify-content: center;
        }}
        .card-image img {{
            max-width: 100%;
            max-height: 100%;
            object-fit: contain;
        }}
        .card-details {{ flex: 1; }}
        .reasoning {{ margin-bottom: 0.5rem; }}
        .description {{ margin-bottom: 0.5rem; font-size: 0.9rem; }}
        .card-context {{
            margin-top: 0.8rem;
            padding-top: 0.5rem;
            border-top: 1px solid #E5E7EB;
            color: #666;
        }}
        .flag {{
            display: inline-block;
            padding: 0.1rem 0.5rem;
            border-radius: 3px;
            font-size: 0.8rem;
            margin-right: 0.3rem;
            margin-bottom: 0.3rem;
        }}
        .flag-critical {{ background: #FEE2E2; color: #DC2626; }}
        .flag-high {{ background: #FFEDD5; color: #EA580C; }}

        /* Footer */
        .footer {{
            text-align: center;
            padding: 2rem;
            color: #999;
            font-size: 0.85rem;
        }}

        /* Print styles */
        @media print {{
            body {{ padding: 0; }}
            .card {{ break-inside: avoid; }}
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>Image Copyright Audit Report</h1>
            <div class="subtitle">
                <strong>Deck:</strong> {pptx_name} &nbsp;|&nbsp;
                <strong>Generated:</strong> {now} &nbsp;|&nbsp;
                <strong>Images scanned:</strong> {total_images}
            </div>
        </div>

        <div class="summary">
            <div class="summary-card">
                <div class="number" style="color: #DC2626;">{risk_counts.get('CRITICAL', 0)}</div>
                <div class="label">Critical</div>
            </div>
            <div class="summary-card">
                <div class="number" style="color: #EA580C;">{risk_counts.get('HIGH', 0)}</div>
                <div class="label">High Risk</div>
            </div>
            <div class="summary-card">
                <div class="number" style="color: #D97706;">{risk_counts.get('MEDIUM', 0)}</div>
                <div class="label">Medium Risk</div>
            </div>
            <div class="summary-card">
                <div class="number" style="color: #16A34A;">{risk_counts.get('LOW', 0)}</div>
                <div class="label">Low Risk</div>
            </div>
            <div class="summary-card">
                <div class="number" style="color: #059669;">{risk_counts.get('CLEAR', 0)}</div>
                <div class="label">Clear</div>
            </div>
            <div class="summary-card">
                <div class="number">{total_images}</div>
                <div class="label">Total Images</div>
            </div>
        </div>

        <div class="risk-bar">
            {"".join(f'<div style="flex: {risk_counts.get(r, 0)}; background: {risk_colours.get(r, "#ccc")};">{risk_counts.get(r, 0) if risk_counts.get(r, 0) > 0 else ""}</div>' for r in ["CRITICAL", "HIGH", "MEDIUM", "LOW", "CLEAR"])}
        </div>

        <div class="section-header">Flagged Images — Requires Action</div>
        {"".join(c for c, (img, cls) in zip(cards_html, paired) if cls.get('risk_level') in ('CRITICAL', 'HIGH', 'MEDIUM'))}

        <div class="section-header">Low Risk / Clear — For Reference</div>
        {"".join(c for c, (img, cls) in zip(cards_html, paired) if cls.get('risk_level') in ('LOW', 'CLEAR', None) or 'error' in cls)}

        <div class="footer">
            <p>Generated by UQ Slide Compliance Tool &nbsp;|&nbsp;
               API tokens used: {total_input_tokens:,} input + {total_output_tokens:,} output &nbsp;|&nbsp;
               <strong>Estimated cost: USD ${total_cost_usd:.4f}</strong> &nbsp;|&nbsp;
               {errors} classification errors</p>
            <p style="margin-top: 0.5rem; color: #bbb;">
                This report is an automated triage — it does not constitute legal advice.
                All flagged images should be reviewed by a human before final clearance.
            </p>
        </div>
    </div>
</body>
</html>"""

    Path(output_path).write_text(html)
    return {
        "total_images": total_images,
        "risk_counts": dict(risk_counts),
        "type_counts": dict(type_counts),
        "errors": errors,
        "tokens": {"input": total_input_tokens, "output": total_output_tokens},
        "cost_usd": round(total_cost_usd, 4),
    }


# ─── Main ──────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="UQ Image Audit Report — classify images for copyright risk"
    )
    parser.add_argument("input", help="Path to input .pptx file")
    parser.add_argument(
        "--output", "-o",
        default="audit_report.html",
        help="Path to output HTML report (default: audit_report.html)",
    )
    parser.add_argument(
        "--limit", "-l",
        type=int,
        default=None,
        help="Limit to first N images (useful for testing)",
    )
    parser.add_argument(
        "--extract-dir",
        default=None,
        help="Directory to save extracted images (default: auto-created next to report)",
    )
    parser.add_argument(
        "--no-classify",
        action="store_true",
        help="Extract images only — skip AI classification",
    )
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file not found: {input_path}")
        sys.exit(1)

    output_path = Path(args.output)

    # Set up extraction directory
    if args.extract_dir:
        extract_dir = Path(args.extract_dir)
    else:
        extract_dir = output_path.parent / f"{output_path.stem}_images"
    extract_dir.mkdir(parents=True, exist_ok=True)

    print(f"\nUQ Image Audit Tool")
    print(f"{'─'*40}")
    print(f"  Input:   {input_path}")
    print(f"  Output:  {output_path}")
    print(f"  Images:  {extract_dir}/")
    if args.limit:
        print(f"  Limit:   {args.limit} images")
    print()

    # Step 1: Extract images
    print("Extracting images from presentation...")
    images = extract_images(input_path, output_dir=extract_dir, limit=args.limit)
    print(f"  Extracted {len(images)} unique images")

    if not images:
        print("  No images found — nothing to audit.")
        sys.exit(0)

    # Show extraction summary
    print(f"\n  Images by slide:")
    slide_counts = defaultdict(int)
    for img in images:
        slide_counts[img["slide_number"]] += 1
    for slide_num in sorted(slide_counts.keys()):
        print(f"    Slide {slide_num}: {slide_counts[slide_num]} image(s)")

    if args.no_classify:
        print(f"\n  --no-classify flag set. Images extracted to {extract_dir}/")
        # Strip image bytes before returning
        for img in images:
            img.pop("image_bytes", None)
        # Save metadata
        meta_path = extract_dir / "metadata.json"
        with open(meta_path, "w") as f:
            json.dump(images, f, indent=2, default=str)
        print(f"  Metadata saved to {meta_path}")
        sys.exit(0)

    # Step 2: Classify with Claude API
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        print("\nError: ANTHROPIC_API_KEY environment variable not set.")
        print("Set it with: export ANTHROPIC_API_KEY=sk-ant-...")
        print("\nTo extract images only, use --no-classify flag.")
        sys.exit(1)

    client = anthropic.Anthropic(api_key=api_key)

    print(f"\nClassifying {len(images)} images with Claude API...")
    print(f"  Model: {CLAUDE_MODEL}")
    print()

    classifications = []
    for i, img in enumerate(images, 1):
        print(f"  [{i}/{len(images)}] Slide {img['slide_number']}: {img['filename']}...", end=" ", flush=True)
        result = classify_image(client, img)
        classifications.append(result)

        if "error" in result:
            print(f"ERROR: {result['error'][:60]}")
        else:
            risk = result.get("risk_level", "?")
            img_type = result.get("image_type", "?")
            print(f"{risk} ({img_type})")

        # Brief pause to avoid rate limiting
        time.sleep(0.5)

    # Step 3: Generate report
    print(f"\nGenerating report...")

    # Strip image bytes before report generation
    for img in images:
        img.pop("image_bytes", None)

    # Use relative path for images in report
    rel_extract_dir = extract_dir.name

    summary = generate_html_report(
        images, classifications, input_path.name,
        output_path, rel_extract_dir,
    )

    # Also save raw classification data as JSON
    json_path = output_path.with_suffix(".json")
    raw_data = {
        "metadata": {
            "input_file": str(input_path),
            "generated": datetime.now().isoformat(),
            "model": CLAUDE_MODEL,
            "total_images": summary["total_images"],
        },
        "summary": summary,
        "images": [
            {**img, **cls}
            for img, cls in zip(images, classifications)
        ],
    }
    with open(json_path, "w") as f:
        json.dump(raw_data, f, indent=2, default=str)

    print(f"\n{'='*60}")
    print(f"  Image Audit Report — Complete")
    print(f"{'='*60}")
    print(f"  Total images:  {summary['total_images']}")
    print(f"  Critical:      {summary['risk_counts'].get('CRITICAL', 0)}")
    print(f"  High risk:     {summary['risk_counts'].get('HIGH', 0)}")
    print(f"  Medium risk:   {summary['risk_counts'].get('MEDIUM', 0)}")
    print(f"  Low risk:      {summary['risk_counts'].get('LOW', 0)}")
    print(f"  Clear:         {summary['risk_counts'].get('CLEAR', 0)}")
    print(f"  Errors:        {summary['errors']}")
    print(f"  API tokens:    {summary['tokens']['input']:,} in / {summary['tokens']['output']:,} out")
    print(f"\n  HTML report:   {output_path}")
    print(f"  JSON data:     {json_path}")
    print(f"  Images:        {extract_dir}/")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
